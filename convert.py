#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
English Worksheet Converter  –  GUI Edition
문제 폴더에서 .docx 파일을 선택하여 4종 시험지로 자동 변환합니다.
"""

import sys, re, os, threading, subprocess, datetime, shutil, json, hashlib, hmac, random, traceback
from collections import defaultdict
from pathlib import Path

# 패키지 경로 탐색 (개발환경 / 배포환경 모두 지원)
for _p in [
    Path(__file__).parent / 'packages',   # 배포 시 로컬 packages 폴더
    Path('D:/Python/packages'),            # 개발 환경
]:
    if _p.exists() and str(_p) not in sys.path:
        sys.path.insert(0, str(_p))

import tkinter as tk
from tkinter import ttk, filedialog

# 드래그 앤 드롭 (설치 시 지원, 없으면 파일 대화상자로 대체)
try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    _ROOT = TkinterDnD.Tk
except ImportError:
    _ROOT = tk.Tk
    DND_FILES = None

# ══════════════════════════════════════════════════════════════════
# 경로 설정
# ══════════════════════════════════════════════════════════════════

BASE_DIR   = Path(__file__).parent
INPUT_DIR  = BASE_DIR / '문제'
OUTPUT_DIR = BASE_DIR / '결과'
DATA_DIR   = BASE_DIR / '데이터'
INPUT_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)
DATA_DIR.mkdir(exist_ok=True)


def _init_custom_data_dir():
    """config.json에 저장된 커스텀 데이터 경로가 있으면 DATA_DIR에 반영."""
    global DATA_DIR, STUDENTS_FILE, SCORES_FILE
    cfg_file = BASE_DIR / 'config.json'
    if not cfg_file.exists():
        return
    try:
        cfg  = json.loads(cfg_file.read_text(encoding='utf-8'))
        cust = cfg.get('data_dir', '')
        if cust:
            p = Path(cust)
            p.mkdir(parents=True, exist_ok=True)
            DATA_DIR      = p
            STUDENTS_FILE = p / 'students.json'
            SCORES_FILE   = p / 'scores.json'
    except Exception:
        pass

_init_custom_data_dir()   # 앱 로드 시 즉시 적용

# 계정별 데이터 디렉토리 (로그인 후 _set_account()로 설정)
_ACCT_DIR     = None
STUDENTS_FILE = DATA_DIR / 'students.json'   # 기본값 (로그인 전)
SCORES_FILE   = DATA_DIR / 'scores.json'

# 계정별 vs 글로벌 설정 키 분류
_ACCT_KEYS   = frozenset({'template_path', 'fmt', 'custom_problem_types',
                           'sheets_id', 'oauth_key'})


def _set_account(uid: str):
    """로그인 계정 설정 — 계정별 전용 디렉토리 & 파일 경로 지정."""
    global _ACCT_DIR, STUDENTS_FILE, SCORES_FILE

    uid_hash = hashlib.sha256(uid.encode('utf-8')).hexdigest()[:16]
    acct_dir = DATA_DIR / 'accounts' / uid_hash
    is_new   = not acct_dir.exists()
    acct_dir.mkdir(parents=True, exist_ok=True)

    global TOKEN_FILE
    _ACCT_DIR     = acct_dir
    STUDENTS_FILE = acct_dir / 'students.json'
    SCORES_FILE   = acct_dir / 'scores.json'
    TOKEN_FILE    = acct_dir / 'google_token.json'

    # 최초 접속 시: 기존 루트 데이터 마이그레이션
    if is_new:
        for src_name, dst in (('students.json', STUDENTS_FILE),
                               ('scores.json',   SCORES_FILE)):
            src = DATA_DIR / src_name
            if src.exists() and not dst.exists():
                shutil.copy2(src, dst)


def _acct_settings_file() -> Path:
    return _ACCT_DIR / 'settings.json' if _ACCT_DIR else None


def _load_acct_settings() -> dict:
    """현재 계정의 개인 설정 로드."""
    f = _acct_settings_file()
    if f and f.exists():
        try:
            return json.loads(f.read_text(encoding='utf-8'))
        except Exception:
            pass
    return {}


def _save_acct_settings(data: dict):
    """현재 계정의 개인 설정 저장."""
    f = _acct_settings_file()
    if f:
        f.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')

# ══════════════════════════════════════════════════════════════════
# 라이센스 키 시스템
# ══════════════════════════════════════════════════════════════════
# ※ 배포 전에 아래 _LIC_SECRET 값을 반드시 변경하세요!
_LIC_SECRET = b'xogus0226!-biobank1717!-whwnsxxod'


def _is_valid_font(name) -> bool:
    """유효한 글꼴 이름 여부 — 테마 참조('+…') 및 nil 제외."""
    return bool(name) and not str(name).startswith('+') and name not in ('nil', '')


def _validate_license(key: str) -> bool:
    """입력된 라이센스 키 유효성 검증."""
    key = key.strip().upper()
    if not key.startswith('EWS-'):
        return False
    parts = key.split('-')
    if len(parts) != 5:
        return False
    inner = ''.join(parts[1:])  # 16자 hex
    if len(inner) != 16 or not all(c in '0123456789ABCDEF' for c in inner):
        return False
    # inner 값이 어떤 teacher_id로부터 나왔는지 역추적 불필요
    # → keygen이 만든 키는 HMAC(secret, id)의 첫 16자이므로,
    #   키 자체가 HMAC 구조임을 확인
    # 저장된 해시와 비교 (설치 시 등록된 키)
    cfg = load_config()
    stored = cfg.get('license_hash', '')
    if stored:
        return hmac.compare_digest(
            stored,
            hashlib.sha256(key.encode()).hexdigest()
        )
    return True   # 첫 등록 시 (아직 저장 안 됨)


def _reset_user_settings():
    """새 라이센스 키 또는 새 아이디로 변경 시 계정/계정설정 초기화.

    글로벌 config에서: app_user_id, app_password_hash, password_last_changed 제거
    계정별 settings:   완전 삭제
    (다른 계정의 데이터는 건드리지 않음)
    """
    # 글로벌에서 계정 인증 정보 제거
    _pop_config('app_user_id', 'app_password_hash', 'password_last_changed')
    # 현재 계정 설정 파일 삭제 (새 계정이 들어올 것이므로)
    f = _acct_settings_file()
    if f and f.exists():
        try: f.unlink()
        except Exception: pass
    # Google 토큰 삭제
    try:
        if TOKEN_FILE.exists(): TOKEN_FILE.unlink()
    except Exception:
        pass


# ══════════════════════════════════════════════════════════════════
# 테마 색상 (라이트)
# ══════════════════════════════════════════════════════════════════

BG    = '#f0f2f7'   # 전체 배경
BG2   = '#ffffff'   # 헤더/상태바
CARD  = '#ffffff'   # 카드 배경
ITEM  = '#f7f8fc'   # 리스트 아이템
HOV   = '#eaecf8'   # 호버
SEL   = '#dce6ff'   # 선택
BDR   = '#d0d5e8'   # 테두리
GOLD  = '#b8932a'   # 골드 (라이트 배경 대비용)
GOLDB = '#d4a730'   # 활성 골드
TEXT  = '#1a1d2e'   # 기본 텍스트
SUB   = '#5c6585'   # 보조 텍스트
DIM   = '#a8b0cc'   # 비활성
OK    = '#16803c'   # 성공
WARN  = '#b45309'   # 경고
ERR   = '#dc2626'   # 에러
PRFG  = '#b8932a'   # 프로그레스바 채움
PRBG  = '#e4e8f5'   # 프로그레스바 배경

KO   = '맑은 고딕'
MONO = 'Consolas'

# ══════════════════════════════════════════════════════════════════
# 파싱 로직
# ══════════════════════════════════════════════════════════════════

ITEM_CODE_RE = re.compile(r'\b([A-Z]+-\d+-(?:TF|MC|EN|AD|WR)\d+)\.')
DASH_SEP_RE  = re.compile(r'^[-:]{2,}$')


def _split_by_codes(text):
    parts = ITEM_CODE_RE.split(text)
    items = []
    i = 1
    while i + 1 < len(parts):
        items.append((parts[i].strip(), parts[i+1].strip()))
        i += 2
    return items


def _parse_mc_item(body):
    syms = ['①','②','③','④','⑤']
    pos  = [(body.find(s), s) for s in syms if body.find(s) != -1]
    if not pos:
        return body.strip(), []
    pos.sort()
    q = body[:pos[0][0]].strip()
    choices = []
    for idx, (p, _) in enumerate(pos):
        end = pos[idx+1][0] if idx+1 < len(pos) else len(body)
        choices.append(body[p+1:end].strip())
    return q, choices


def _parse_pipe_table(text):
    raw = [t.strip() for t in text.split('|')]
    while raw and raw[0] == '': raw.pop(0)
    while raw and raw[-1] == '': raw.pop()
    if not raw: return []
    SEP    = '\x00'
    tokens = [SEP if DASH_SEP_RE.match(t) else t for t in raw]
    header = []
    hi = 0
    while hi < len(tokens) and tokens[hi] not in ('', SEP):
        header.append(tokens[hi]); hi += 1
    n_cols = len(header)
    if n_cols == 0: return []
    rows = [header]
    i = hi
    while i < len(tokens) and tokens[i] in ('', SEP): i += 1
    cur = []
    for tok in tokens[i:]:
        if tok == SEP:
            if cur: rows.append(cur); cur = []
        elif tok == '':
            if   len(cur) == 0:       pass
            elif len(cur) >= n_cols:  pass
            else:                     cur.append('')
        else:
            cur.append(tok)
            if len(cur) >= n_cols:
                rows.append(cur); cur = []
    if cur: rows.append(cur)
    return rows


def _parse_keyword_block(text):
    pipe = text.find('|')
    if pipe == -1:
        return {'intro': re.sub(r'^[━\s📝]+', '', text).strip(), 'table': None}
    intro = re.sub(r'^[━\s📝]+', '', text[:pipe]).strip()
    traw  = re.sub(r'━+\s*$', '', text[pipe:]).strip()
    return {'intro': intro, 'table': _parse_pipe_table(traw)}


def _parse_chunk_items(text):
    end = text.find('①')
    if end == -1: return {'intro': text, 'items': []}
    intro = text[:end].strip()
    rest  = text[end:]
    syms  = ['①','②','③','④']
    pos   = [(rest.find(s), s) for s in syms if rest.find(s) != -1]
    pos.sort()
    items = []
    for idx, (p, _) in enumerate(pos):
        chunk = re.sub(r'━+\s*$', '', rest[p+1:(pos[idx+1][0] if idx+1 < len(pos) else len(rest))]).strip()
        hint  = (re.search(r'\(([^)]+)\)', chunk) or type('',(),{'group':lambda *_:''})()).group(1)
        items.append((str(idx+1), re.sub(r'\([^)]+\)', '', chunk).strip(), hint))
    return {'intro': intro, 'items': items}


def _parse_answer_table(text):
    p = text.find('|')
    return _parse_pipe_table(text[p:]) if p != -1 else []


def detect_section(text):
    if re.search(r'STANDARD\s*B\s*시험지|🟡\s*STANDARD\s*B', text): return 'standard_b'
    if re.search(r'STANDARD\s*A\s*시험지|🟡\s*STANDARD\s*A', text): return 'standard_a'
    if re.search(r'EASY\s*시험지|🟢\s*EASY', text) or re.match(r'^EASY\s*$', text.strip()): return 'easy'
    if re.search(r'HARD\s*시험지|🔴\s*HARD', text): return 'hard'
    return None


def parse_docx(filepath):
    NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    doc  = Document(filepath)

    # 단락과 Word 표를 문서 순서대로 수집
    body_items: list[tuple] = []
    for child in doc.element.body:
        local = child.tag.split('}')[-1]
        if local == 'p':
            txt = ''.join(t.text or '' for t in child.findall(f'.//{{{NS_W}}}t'))
            body_items.append(('p', txt))
        elif local == 'tbl':
            rows = []
            for tr in child.findall(f'.//{{{NS_W}}}tr'):
                row = [''.join(t.text or '' for t in
                               tc.findall(f'.//{{{NS_W}}}t')).strip()
                       for tc in tr.findall(f'{{{NS_W}}}tc')]
                if any(row):
                    rows.append(row)
            if rows:
                body_items.append(('tbl', rows))

    secs = {k: {'subtitle': '', 'passages': []} for k in ('easy','standard_a','standard_b','hard')}
    cur_sec = None; cur_pass = None; last_bt = None

    def new_pass(n): return {'num': n, 'blocks': [], 'answer_table': []}
    def flush():
        nonlocal cur_pass
        if cur_pass and cur_sec: secs[cur_sec]['passages'].append(cur_pass)
        cur_pass = None

    def _peek_text(j):
        """body_items[j] 가 단락이면 텍스트, 아니면 None."""
        if j < len(body_items) and body_items[j][0] == 'p':
            return body_items[j][1].strip()
        return None

    i = 0
    while i < len(body_items):
        kind, payload = body_items[i]; i += 1

        # Word 표 아이템: 키워드 블록 직후에 오면 표 데이터로 채움
        if kind == 'tbl':
            if (last_bt == 'keyword' and cur_pass and cur_pass['blocks']
                    and cur_pass['blocks'][-1]['type'] == 'keyword'
                    and not cur_pass['blocks'][-1].get('table')):
                cur_pass['blocks'][-1]['table'] = payload
            continue

        text = payload.strip()
        if not text: continue
        sec = detect_section(text)
        if sec:
            flush(); cur_sec = sec; last_bt = None; continue
        if cur_sec is None: continue
        if not secs[cur_sec]['passages'] and cur_pass is None:
            if not re.match(r'^(?:##\s*)?지문\s*\d+', text):
                secs[cur_sec]['subtitle'] = (secs[cur_sec]['subtitle'] + ' | ' + text).lstrip(' | '); continue
        pm = re.match(r'^(?:##\s*)?지문\s*(\d+)\s*$', text)
        if pm:
            flush(); cur_pass = new_pass(int(pm.group(1))); last_bt = None; continue
        if cur_pass is None: continue
        if re.match(r'^[-─━═\s]{10,}$', text) and not re.search(r'[가-힣a-zA-Z]', text): continue
        if re.search(r'\[📋|정답표\]', text):
            cur_pass['answer_table'] = _parse_answer_table(text); continue

        if re.search(r'오늘의 변형 예상 키워드', text):
            intro = re.sub(r'^[━\s📝]+', '', text).strip()
            kw_table = None
            if '|' in text:
                blk = _parse_keyword_block(text)
                intro, kw_table = blk['intro'], blk.get('table')
            else:
                # 이후 아이템에서 Word 표 또는 파이프 행 탐색
                j = i
                while j < len(body_items):
                    jkind, jpayload = body_items[j]
                    if jkind == 'tbl':           # ← Word 표 발견
                        kw_table = jpayload; i = j + 1; break
                    nxt = jpayload.strip()
                    if not nxt: j += 1; continue
                    if re.match(r'^[-─━═\s]{4,}$', nxt) and not re.search(r'[가-힣a-zA-Z0-9|]', nxt):
                        j += 1; continue
                    if '|' not in nxt: break     # 관련 없는 단락 → 중단
                    combined_kw = text
                    while j < len(body_items):
                        jkind2, jpayload2 = body_items[j]
                        if jkind2 == 'tbl': break
                        nxt2 = jpayload2.strip()
                        if not nxt2: j += 1; continue
                        if re.match(r'^[-─━═\s]{4,}$', nxt2) and not re.search(r'[가-힣a-zA-Z0-9|]', nxt2):
                            j += 1; continue
                        if '|' in nxt2: combined_kw += ' ' + nxt2; j += 1
                        else: break
                    blk = _parse_keyword_block(combined_kw)
                    intro, kw_table = blk['intro'], blk.get('table')
                    i = j; break
            cur_pass['blocks'].append({'type': 'keyword', 'intro': intro, 'table': kw_table})
            last_bt = 'keyword'; continue

        if re.search(r'✍️\s*서술형 대비.*핵심 Chunk', text):
            combined = re.sub(r'^[━\s]+|━+\s*$', '', text).strip()
            if '①' not in combined:
                nxt = _peek_text(i)
                if nxt and '①' in nxt: combined += ' ' + nxt; i += 1
            cur_pass['blocks'].append({'type': 'chunk', **_parse_chunk_items(combined)})
            last_bt = 'chunk'; continue
        if re.match(r'^[①②]', text) and last_bt == 'chunk':
            if cur_pass['blocks'] and cur_pass['blocks'][-1]['type'] == 'chunk' and not cur_pass['blocks'][-1]['items']:
                cur_pass['blocks'][-1].update(_parse_chunk_items(text))
            continue
        if re.match(r'^\[T/F', text):
            hm = re.match(r'^(\[[^\]]+\])', text); hdr = hm.group(1) if hm else ''
            items = []
            for n, (_, body) in enumerate(_split_by_codes(text[len(hdr):].strip()), 1):
                body = body.strip()
                body = (body[:-5].rstrip() if body.endswith(('(T/F)','( T / F )')) else body.rstrip('.')) + ' ( T / F )'
                items.append((str(n), body))
            cur_pass['blocks'].append({'type': 'tf', 'header': hdr, 'items': items})
            last_bt = 'tf'; continue
        if re.match(r'^\[한국어 객관식', text):
            hm = re.match(r'^(\[[^\]]+\])', text); hdr = hm.group(1) if hm else ''
            items = [(str(n), *_parse_mc_item(b)) for n, (_, b) in enumerate(_split_by_codes(text[len(hdr):].strip()), 1)]
            cur_pass['blocks'].append({'type': 'ko_mc', 'header': hdr, 'items': items}); last_bt = 'ko_mc'; continue
        if re.match(r'^\[영어 객관식', text):
            hm = re.match(r'^(\[[^\]]+\])', text); hdr = hm.group(1) if hm else ''
            items = [(str(n), *_parse_mc_item(b)) for n, (_, b) in enumerate(_split_by_codes(text[len(hdr):].strip()), 1)]
            cur_pass['blocks'].append({'type': 'en_mc', 'header': hdr, 'items': items}); last_bt = 'en_mc'; continue
        if re.match(r'^\[Hard 전용 추가 영어', text):
            hm = re.match(r'^(\[[^\]]+\])', text); hdr = hm.group(1) if hm else ''
            items = []
            for n, (_, body) in enumerate(_split_by_codes(text[len(hdr):].strip()), 1):
                sm = re.search(r'\[(\d+점)\]', body); score = sm.group(1) if sm else ''
                b2 = (body[:sm.start()] + body[sm.end():]).strip() if sm else body
                q, ch = _parse_mc_item(b2); items.append((str(n), q, score, ch))
            cur_pass['blocks'].append({'type': 'hard_add', 'header': hdr, 'items': items}); last_bt = 'hard_add'; continue
        if re.match(r'^\[Hard 전용 융합형 논술', text):
            hm = re.match(r'^(\[[^\]]+\])', text); hdr = hm.group(1) if hm else ''
            raw = _split_by_codes(text[len(hdr):].strip()); instr = raw[0][1] if raw else text[len(hdr):].strip()
            tpl = score_ln = ''
            while i < len(body_items):
                jkind, jpayload = body_items[i]
                if jkind == 'tbl': break
                nxt = jpayload.strip()
                if not nxt: i += 1; continue
                if re.search(r'\[📋|정답표\]', nxt) or detect_section(nxt) or re.match(r'^(?:##\s*)?지문\s*\d+', nxt): break
                i += 1
                if re.search(r'부분 점수', nxt): score_ln = nxt
                elif not tpl and nxt: tpl = nxt
            cur_pass['blocks'].append({'type': 'hard_wr', 'header': hdr, 'instruction': instr, 'template': tpl, 'score_line': score_ln})
            last_bt = 'hard_wr'; continue
        raw = _split_by_codes(text)
        if raw and cur_pass['blocks']:
            lb = cur_pass['blocks'][-1]
            if last_bt in ('ko_mc','en_mc') and lb['type'] == last_bt:
                start = len(lb['items']) + 1
                for off, (_, b) in enumerate(raw): lb['items'].append((str(start+off), *_parse_mc_item(b)))
                continue
            if last_bt == 'hard_add' and lb['type'] == 'hard_add':
                start = len(lb['items']) + 1
                for off, (_, body) in enumerate(raw):
                    sm = re.search(r'\[(\d+점)\]', body); score = sm.group(1) if sm else ''
                    b2 = (body[:sm.start()] + body[sm.end():]).strip() if sm else body
                    q, ch = _parse_mc_item(b2); lb['items'].append((str(start+off), q, score, ch))
                continue
    flush()
    return secs


# ══════════════════════════════════════════════════════════════════
# Word 출력 로직
# ══════════════════════════════════════════════════════════════════

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph

SEC_LABELS = {'easy': 'Easy', 'standard_a': 'Standard A',
              'standard_b': 'Standard B', 'hard': 'Hard'}
KO_FONT = '맑은 고딕'


def _add_page_break(doc):
    """명시적 페이지 나누기 삽입."""
    p = doc.add_paragraph()
    r = p.add_run()
    br = OxmlElement('w:br')
    br.set(qn('w:type'), 'page')
    r._r.append(br)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(0)


def _content_table(doc, rows):
    """정답표용 5열 표 생성."""
    if len(rows) < 2: return
    hdr, data = rows[0], rows[1:]
    nc  = max(len(r) for r in rows)
    tbl = doc.add_table(rows=1 + len(data), cols=nc)
    tbl.style = 'Table Grid'

    def _set_cell(cell, text, bold=False, bg=None):
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(1)
        p.paragraph_format.space_after  = Pt(1)
        r = p.add_run(text)
        r.font.size = Pt(8.5); r.font.bold = bold
        r.font.name = KO_FONT
        r.font.element.rPr.rFonts.set(qn('w:eastAsia'), KO_FONT)
        if bg:
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd  = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), bg)
            tcPr.append(shd)

    for ci, ct in enumerate(hdr[:nc]):
        _set_cell(tbl.rows[0].cells[ci], ct, bold=True, bg='E8EDF7')

    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd[:nc]):
            _set_cell(tbl.rows[ri+1].cells[ci], ct)

    doc.add_paragraph().paragraph_format.space_after = Pt(4)


CONFIG_FILE = None   # 앱 시작 시 BASE_DIR 확정 후 설정됨


def load_config() -> dict:
    """설정 로드: 글로벌(config.json) + 계정별(settings.json) 합산."""
    cfg = {}
    if CONFIG_FILE and CONFIG_FILE.exists():
        try:
            cfg = json.loads(CONFIG_FILE.read_text(encoding='utf-8'))
        except Exception:
            pass
    # 계정별 설정으로 덮어쓰기 (계정별이 우선)
    cfg.update(_load_acct_settings())
    return cfg


def save_config(cfg: dict):
    """설정 저장: 키 종류에 따라 글로벌/계정별로 분리 저장."""
    # ① 글로벌 설정 (config.json)
    if CONFIG_FILE:
        try:
            existing = {}
            if CONFIG_FILE.exists():
                existing = json.loads(CONFIG_FILE.read_text(encoding='utf-8'))
        except Exception:
            existing = {}
        for k, v in cfg.items():
            if k not in _ACCT_KEYS:       # 계정별 키는 global에 저장 안 함
                existing[k] = v
        for k in _ACCT_KEYS:             # 혹시 남아있는 계정 키 제거
            existing.pop(k, None)
        CONFIG_FILE.write_text(
            json.dumps(existing, ensure_ascii=False, indent=2), encoding='utf-8')

    # ② 계정별 설정 (accounts/{uid}/settings.json)
    acct_updates = {k: v for k, v in cfg.items() if k in _ACCT_KEYS}
    if acct_updates:
        acct = _load_acct_settings()
        acct.update(acct_updates)
        _save_acct_settings(acct)


def _pop_config(*keys):
    """config에서 특정 키 제거 (글로벌 + 계정별 모두)."""
    # 글로벌
    if CONFIG_FILE and CONFIG_FILE.exists():
        try:
            c = json.loads(CONFIG_FILE.read_text(encoding='utf-8'))
            changed = False
            for k in keys:
                if k in c:
                    del c[k]; changed = True
            if changed:
                CONFIG_FILE.write_text(json.dumps(c, ensure_ascii=False, indent=2), encoding='utf-8')
        except Exception:
            pass
    # 계정별
    acct = _load_acct_settings()
    changed = False
    for k in keys:
        if k in acct:
            del acct[k]; changed = True
    if changed:
        _save_acct_settings(acct)


def read_template_fmt(path):
    """파일 형식별 서식 설정 추출 (docx / hwp / pdf / pptx·ppt)."""
    path = Path(path)
    ext  = path.suffix.lower()

    if ext == '.docx':
        return _fmt_from_docx(path)
    elif ext == '.hwp':
        return _fmt_from_hwp(path)
    elif ext in ('.pptx', '.ppt'):
        return _fmt_from_pptx(path)
    return {}


def _fmt_from_docx(path):
    """Word(.docx) → 글꼴·여백 추출.

    우선순위: Normal 스타일 → rPrDefault → 실제 run eastAsia → run name → 테마
    한국어 문서는 eastAsia 글꼴을 우선 탐색한다.
    """
    try:
        doc = Document(str(path))
        sec = doc.sections[0]
        is_landscape = sec.page_width > sec.page_height
        fmt = {
            'source_type':      'docx',
            'orientation':      'landscape' if is_landscape else 'portrait',
            'page_width_cm':    round(float(sec.page_width.cm),    1),
            'page_height_cm':   round(float(sec.page_height.cm),   1),
            'top_margin_cm':    round(float(sec.top_margin.cm),    1),
            'bottom_margin_cm': round(float(sec.bottom_margin.cm), 1),
            'left_margin_cm':   round(float(sec.left_margin.cm),   1),
            'right_margin_cm':  round(float(sec.right_margin.cm),  1),
        }

        _ok = _is_valid_font   # 모듈 레벨 헬퍼 참조

        # ① Normal 스타일 직접 글꼴 (eastAsia 우선)
        try:
            ns = doc.styles['Normal']
            # python-docx font.name 은 ASCII 기준이므로 XML에서 직접 추출
            try:
                rf = ns.font.element.find(qn('w:rFonts'))
                if rf is not None:
                    fn = rf.get(qn('w:eastAsia')) or rf.get(qn('w:ascii'))
                    if _ok(fn):
                        fmt['font_name'] = fn
            except Exception:
                pass
            if 'font_name' not in fmt and _ok(ns.font.name):
                fmt['font_name'] = ns.font.name
            if ns.font.size:
                fmt['font_size'] = round(float(ns.font.size.pt), 1)
        except Exception:
            pass

        # ② rPrDefault XML (eastAsia 우선)
        if 'font_name' not in fmt:
            try:
                NS_W = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                rpr = doc.element.find(f'.//{{{NS_W}}}rPrDefault')
                if rpr is not None:
                    rf = rpr.find(f'.//{{{NS_W}}}rFonts')
                    if rf is not None:
                        for attr in (qn('w:eastAsia'), qn('w:ascii'), qn('w:hAnsi')):
                            fn = rf.get(attr)
                            if _ok(fn):
                                fmt['font_name'] = fn
                                break
            except Exception:
                pass

        # ③ 첫 번째 실제 run의 eastAsia 글꼴 (한국어 문서 핵심)
        if 'font_name' not in fmt:
            try:
                for p in doc.paragraphs:
                    for r in p.runs:
                        # eastAsia 먼저
                        fn = None
                        try:
                            rf = r.font.element.rPr.rFonts
                            fn = rf.get(qn('w:eastAsia'))
                        except Exception:
                            pass
                        if not _ok(fn):
                            fn = r.font.name   # Latin fallback
                        if _ok(fn):
                            fmt['font_name'] = fn
                            if r.font.size:
                                fmt['font_size'] = round(float(r.font.size.pt), 1)
                            break
                    if 'font_name' in fmt:
                        break
            except Exception:
                pass

        # ④ 테마 XML → 최후 수단 (Latin 전용이라 한국어에는 부적합하나 없는 것보단 낫다)
        if 'font_name' not in fmt:
            try:
                for rel in doc.part.rels.values():
                    if 'theme' in str(rel.reltype):
                        xml = rel.target_part.blob.decode('utf-8', errors='ignore')
                        m = re.search(r'<a:latin\s+typeface="([^"+ ][^"]*)"', xml)
                        if m:
                            fmt['font_name'] = m.group(1)
                        break
            except Exception:
                pass

        # ── 단락별 서식 추출 ────────────────────────────────────
        fmt['para_specs'] = _extract_para_specs(doc)
        return fmt
    except Exception:
        return {'source_type': 'docx'}


def _extract_para_specs(doc) -> dict:
    """서식 파일(.docx)에서 요소별 단락 서식 추출.
    반환: {key: {size, bold, sb, sa, ind, align}, ...}
    """
    _ALIGN = {
        WD_ALIGN_PARAGRAPH.CENTER:  'CENTER',
        WD_ALIGN_PARAGRAPH.RIGHT:   'RIGHT',
        WD_ALIGN_PARAGRAPH.LEFT:    'LEFT',
        WD_ALIGN_PARAGRAPH.JUSTIFY: 'JUSTIFY',
    }

    def _info(p):
        pf = p.paragraph_format
        fn, fs, bd = '맑은 고딕', None, False
        for r in p.runs:
            if r.text.strip():
                try:
                    ea = r.font.element.rPr.rFonts.get(qn('w:eastAsia'))
                    fn = ea or r.font.name or '맑은 고딕'
                except Exception:
                    fn = r.font.name or '맑은 고딕'
                fs = round(r.font.size.pt, 1) if r.font.size else None
                bd = bool(r.font.bold)
                break
        return {
            'font':  fn,
            'size':  fs,
            'bold':  bd,
            'sb':    round(pf.space_before.pt, 1) if pf.space_before else 0,
            'sa':    round(pf.space_after.pt, 1)  if pf.space_after  else 0,
            'ind':   round(pf.left_indent.cm, 2)  if pf.left_indent  else 0,
            'align': _ALIGN.get(p.alignment),
        }

    specs = {}
    past_ans = False   # 정답표 이후 여부

    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        info = _info(p)

        # ── 패턴 매칭 ──────────────────────────────────────────
        if 'WORKSHEET' in t and '이름' in t and 'ws_header' not in specs:
            specs['ws_header'] = info

        elif '점수' in t and '100' in t and 'score' not in specs:
            specs['score'] = info

        elif t in ('Easy', 'Standard A', 'Standard B', 'Hard') and 'section' not in specs:
            specs['section'] = info

        elif re.match(r'^지문\s*\d+\s*$', t):
            if not past_ans and 'passage' not in specs:
                specs['passage'] = info
            elif past_ans and 'ans_passage' not in specs:
                specs['ans_passage'] = info

        elif t.startswith('[T/F') and 'blk_header' not in specs:
            specs['blk_header'] = info

        elif (t.startswith('[한국어') or t.startswith('[영어')) and 'mc_blk_header' not in specs:
            specs['mc_blk_header'] = info

        elif ('[서술형' in t or 'Chunk' in t) and 'chunk_hdr' not in specs:
            specs['chunk_hdr'] = info

        elif t.startswith('[오늘의 변형') and 'keyword_hdr' not in specs:
            specs['keyword_hdr'] = info

        elif re.search(r'\(\s*T\s*/\s*F\s*\)\s*$', t) and 'tf_item' not in specs:
            specs['tf_item'] = info

        elif re.match(r'^[①②③④⑤]\s', t):
            if '___' in t and 'chunk_item' not in specs:
                specs['chunk_item'] = info
            elif '___' not in t and 'mc_choice' not in specs:
                specs['mc_choice'] = info

        elif re.match(r'^\d+\.\s', t) and '( T / F )' not in t and 'mc_body' not in specs:
            specs['mc_body'] = info

        elif '정' in t and '답' in t and '표' in t and len(t) < 15 and 'ans_title' not in specs:
            specs['ans_title'] = info
            past_ans = True

        elif t.startswith('[Hard 전용') and 'hard_blk_header' not in specs:
            specs['hard_blk_header'] = info

    # 누락 키 보완 (mc_blk_header 없으면 blk_header 복사 등)
    for src, dst in (('blk_header', 'mc_blk_header'),
                     ('blk_header', 'hard_blk_header'),
                     ('keyword_hdr', 'chunk_hdr'),
                     ('passage', 'ans_passage')):
        if src in specs and dst not in specs:
            specs[dst] = dict(specs[src])

    return specs


def _fmt_from_hwp(path):
    """HWP(.hwp) → 페이지 크기·여백·글꼴 직접 바이너리 추출.

    검증된 태그 번호 (HWP5):
      DocInfo  – HWPTAG_FACE_NAME = 19  (0x13)  : 글꼴 이름
      BodyText – HWPTAG_PAGE_DEF  = 73  (0x49)  : 용지 크기/여백
    PageDef 오프셋: [0]paper_w [4]paper_h [8]left [12]right [20]bottom [24]top
    1 HWPUNIT = 1/7200 inch ≈ 0.0003528 cm
    스트림은 zlib raw-deflate 압축되어 있음.
    """
    fmt = {'source_type': 'hwp'}
    HU  = 0.0003528

    try:
        import struct, olefile, zlib

        def _decompress(raw):
            try:    return zlib.decompress(raw, -15)   # raw deflate
            except: pass
            try:    return zlib.decompress(raw)        # zlib with header
            except: return raw

        def _records(data):
            pos = 0
            while pos + 4 <= len(data):
                hdr    = struct.unpack_from('<I', data, pos)[0]
                tag_id = hdr & 0x3FF
                size   = (hdr >> 20) & 0xFFF
                pos   += 4
                if size == 0xFFF:
                    if pos + 4 > len(data): break
                    size = struct.unpack_from('<I', data, pos)[0]; pos += 4
                yield tag_id, data[pos:pos + size] if pos + size <= len(data) else b''
                pos += size

        with olefile.OleFileIO(str(path)) as ole:

            # ① 글꼴: DocInfo – HWPTAG_FACE_NAME (tag 19)
            if ole.exists('DocInfo'):
                doc_data = _decompress(ole.openstream('DocInfo').read())
                for tid, rec in _records(doc_data):
                    if tid == 19 and len(rec) >= 3:   # FaceName
                        nlen = struct.unpack_from('<H', rec, 1)[0]
                        name = rec[3:3 + nlen * 2].decode('utf-16-le', errors='replace').rstrip('\x00')
                        if name and not name.startswith('@') and len(name) > 1:
                            fmt['font_name'] = name
                            break

            # ② 용지: BodyText/Section0 – HWPTAG_PAGE_DEF (tag 73)
            sec_stream = None
            for entry in ole.listdir():
                joined = '/'.join(entry)
                if 'BodyText' in joined and 'Section' in joined:
                    sec_stream = joined; break

            if sec_stream and ole.exists(sec_stream):
                sec_data = _decompress(ole.openstream(sec_stream).read())
                for tid, rec in _records(sec_data):
                    if tid == 73 and len(rec) >= 32:   # PageDef
                        pw = abs(struct.unpack_from('<i', rec,  0)[0])
                        ph = abs(struct.unpack_from('<i', rec,  4)[0])
                        lm = abs(struct.unpack_from('<i', rec,  8)[0])
                        rm = abs(struct.unpack_from('<i', rec, 12)[0])
                        bm = abs(struct.unpack_from('<i', rec, 20)[0])
                        tm = abs(struct.unpack_from('<i', rec, 24)[0])
                        pw_cm = round(pw * HU, 1)
                        ph_cm = round(ph * HU, 1)
                        if 10 < pw_cm < 100 and 10 < ph_cm < 100:
                            fmt.update({
                                'page_width_cm':    pw_cm,
                                'page_height_cm':   ph_cm,
                                'left_margin_cm':   round(lm * HU, 1),
                                'right_margin_cm':  round(rm * HU, 1),
                                'top_margin_cm':    round(tm * HU, 1),
                                'bottom_margin_cm': round(bm * HU, 1),
                                'orientation': 'landscape' if pw_cm > ph_cm else 'portrait',
                            })
                        break

    except Exception:
        pass
    return fmt



def _fmt_from_pptx(path):
    """PPTX/PPT(.pptx/.ppt) → 슬라이드 크기·글꼴 추출."""
    fmt = {'source_type': 'pptx'}
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        # EMU → cm  (1 EMU = 1/914400 m = 1/360000 inch = 0.000002540005 m)
        w_cm = round(prs.slide_width  / 914400 * 2.54, 1)
        h_cm = round(prs.slide_height / 914400 * 2.54, 1)
        fmt['page_width_cm']  = w_cm
        fmt['page_height_cm'] = h_cm
        fmt['orientation']    = 'landscape' if w_cm > h_cm else 'portrait'

        _ok = _is_valid_font   # 모듈 레벨 헬퍼 참조

        # ① 슬라이드 마스터 테마 XML – 동아시아(한국어) 폰트 우선
        try:
            xml = prs.slide_master._element.xml
            # <a:ea typeface="..."> 동아시아 폰트
            m = re.search(r'<a:ea\s+typeface="([^"+ ][^"]*)"', xml)
            if m and _ok(m.group(1)):
                fmt['font_name'] = m.group(1)
            else:
                # <a:latin typeface="..."> 라틴 폰트
                m = re.search(r'<a:latin\s+typeface="([^"+ ][^"]*)"', xml)
                if m and _ok(m.group(1)):
                    fmt['font_name'] = m.group(1)
        except Exception:
            pass

        # ② 슬라이드 텍스트 실제 run에서 폰트 추출 (최후 수단)
        if 'font_name' not in fmt:
            try:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame: continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                fn = run.font.name
                                if _ok(fn):
                                    fmt['font_name'] = fn
                                    if run.font.size:
                                        fmt['font_size'] = round(run.font.size.pt, 1)
                                    break
                            if 'font_name' in fmt: break
                        if 'font_name' in fmt: break
                    if 'font_name' in fmt: break
            except Exception:
                pass

    except Exception:
        pass
    return fmt


# ══════════════════════════════════════════════════════════════════
# 문서 생성 — 서식 파일(standard_a.docx) 정밀 측정 기반
# 측정값 출처: 문제 서식/standard_a.docx 직접 분석 (2026-05-12)
# ══════════════════════════════════════════════════════════════════

# ── 기본값 (서식 파일 미지정 시 fallback) ────────────────────────
# standard_a.docx 실측 기반. 서식 파일 설정 시 para_specs가 우선.
_DEF_SPECS = {
    'ws_header':       {'size': 12, 'bold': True,  'sb': 0,  'sa': 0,  'ind': 0,   'align': 'CENTER'},
    'score':           {'size': 10, 'bold': False, 'sb': 2,  'sa': 8,  'ind': 0,   'align': 'RIGHT'},
    'section':         {'size': 11, 'bold': True,  'sb': 0,  'sa': 6,  'ind': 0},
    'passage':         {'size': 11, 'bold': True,  'sb': 8,  'sa': 3,  'ind': 0},
    'blk_header':      {'size': 10, 'bold': True,  'sb': 4,  'sa': 2,  'ind': 0},
    'mc_blk_header':   {'size': 10, 'bold': True,  'sb': 4,  'sa': 2,  'ind': 0},
    'hard_blk_header': {'size': 10, 'bold': True,  'sb': 4,  'sa': 2,  'ind': 0},
    'tf_item':         {'size':  9, 'bold': False, 'sb': 1,  'sa': 2,  'ind': 0},
    'mc_body':         {'size': 10, 'bold': False, 'sb': 3,  'sa': 1,  'ind': 0},
    'mc_choice':       {'size': 10, 'bold': False, 'sb': 0,  'sa': 1,  'ind': 0.5},
    'keyword_hdr':     {'size': 10, 'bold': True,  'sb': 5,  'sa': 2,  'ind': 0},
    'chunk_hdr':       {'size': 10, 'bold': True,  'sb': 5,  'sa': 2,  'ind': 0},
    'chunk_item':      {'size': 10, 'bold': False, 'sb': 1,  'sa': 2,  'ind': 0.5},
    'hard_body':       {'size': 10, 'bold': False, 'sb': 3,  'sa': 1,  'ind': 0},
    'hard_choice':     {'size': 10, 'bold': False, 'sb': 0,  'sa': 1,  'ind': 0.5},
    'wr_instr':        {'size': 10, 'bold': False, 'sb': 2,  'sa': 2,  'ind': 0},
    'wr_template':     {'size': 10, 'bold': True,  'sb': 3,  'sa': 3,  'ind': 0.3},
    'wr_score':        {'size':  9, 'bold': False, 'sb': 1,  'sa': 4,  'ind': 0},
    'ans_title':       {'size': 11, 'bold': True,  'sb': 0,  'sa': 10, 'ind': 0,   'align': 'CENTER'},
    'ans_passage':     {'size': 10, 'bold': True,  'sb': 6,  'sa': 3,  'ind': 0},
}


def _bd_para(doc, text, size, bold=False,
             align=None, sb=0, sa=2, ind=None, color=None):
    """단락 한 개 생성 — 크기/간격을 인자로 완전 제어."""
    p = doc.add_paragraph()
    f = p.paragraph_format
    f.space_before = Pt(sb); f.space_after = Pt(sa)
    if align: f.alignment = align
    if ind is not None: f.left_indent = Cm(ind)
    r = p.add_run(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = KO_FONT
    r.font.element.rPr.rFonts.set(qn('w:eastAsia'), KO_FONT)
    if color: r.font.color.rgb = RGBColor(*color)
    return p


def _setup_page_header(doc):
    """모든 페이지 상단에 WORKSHEET 헤더 자동 반복 (Word 페이지 헤더 기능).
    기본값 고정: 12pt Bold 가운데 / 10pt 오른쪽
    """
    for section in doc.sections:
        hdr      = section.header
        hdr_body = hdr._element

        # 기존 단락 제거
        for old_p in list(hdr_body.findall(qn('w:p'))):
            hdr_body.remove(old_p)

        def _hp(text_parts, align, sb, sa):
            p_el = OxmlElement('w:p')
            hdr_body.append(p_el)
            p = Paragraph(p_el, hdr)
            p.alignment = align
            p.paragraph_format.space_before = Pt(sb)
            p.paragraph_format.space_after  = Pt(sa)
            for text, sz, bold in text_parts:
                r = p.add_run(text)
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.name = KO_FONT
                try:
                    r.font.element.rPr.rFonts.set(qn('w:eastAsia'), KO_FONT)
                except Exception:
                    pass
            return p

        # 줄 1: WORKSHEET / 이름 (가운데, 12pt Bold)
        _hp([
            ('WORKSHEET  /  이름 :  ', 12, True),
            ('_' * 22,                 12, True),
            ('     (이름 미기입시 0점 처리)', 12, True),
        ], align=WD_ALIGN_PARAGRAPH.CENTER, sb=0, sa=2)

        # 줄 2: 점수 (오른쪽, 10pt)
        _hp([
            ('점수 :      / 100 점', 10, False),
        ], align=WD_ALIGN_PARAGRAPH.RIGHT, sb=2, sa=6)


def _bd_clear(doc):
    """문서 body 콘텐츠 초기화, sectPr(페이지 설정)만 보존."""
    body    = doc.element.body
    sect_pr = body.find(qn('w:sectPr'))
    for child in list(body): body.remove(child)
    if sect_pr is not None: body.append(sect_pr)


def build_doc(section_key, section_data, output_dir, fmt=None):
    """워크시트 Word 문서 생성.

    서식 파일의 para_specs(단락별 서식)를 사용.
    서식 파일 미지정 시 _DEF_SPECS(standard_a.docx 실측)로 폴백.
    고정 동작: 지문별 페이지 나누기 / 블록 유형 변경 시 구분 공백
    """
    fmt = fmt or {}

    # ── 서식 스펙 준비 ────────────────────────────────────────────
    # 서식 파일 추출값이 있으면 _DEF_SPECS 위에 덮어씀
    _file_ps = fmt.get('para_specs', {})
    _PS = {}
    for k, defaults in _DEF_SPECS.items():
        merged = dict(defaults)
        merged.update({kk: vv for kk, vv in _file_ps.get(k, {}).items()
                       if vv is not None})
        _PS[k] = merged

    _ALIGN_TO_WDA = {'CENTER': WD_ALIGN_PARAGRAPH.CENTER, 'RIGHT': WD_ALIGN_PARAGRAPH.RIGHT,
                     'LEFT': WD_ALIGN_PARAGRAPH.LEFT, 'JUSTIFY': WD_ALIGN_PARAGRAPH.JUSTIFY}

    def _g(key):
        """스펙 dict 반환 (WDA enum 변환 포함)."""
        s = _PS.get(key, _DEF_SPECS.get(key, {}))
        a = s.get('align')
        return {
            'size': s.get('size', 10),
            'bold': s.get('bold', False),
            'sb':   s.get('sb', 0),
            'sa':   s.get('sa', 2),
            'ind':  s.get('ind', 0) or None,
            'align': _ALIGN_TO_WDA.get(a) if a else None,
        }

    def _p(doc, text, key):
        """스펙 기반 단락 생성."""
        g = _g(key)
        return _bd_para(doc, text, g['size'], bold=g['bold'],
                        sb=g['sb'], sa=g['sa'],
                        ind=g['ind'], align=g['align'])

    # ── 글꼴: 서식 파일 font_name 사용 (없으면 맑은 고딕) ─────────
    _font = fmt.get('font_name', '맑은 고딕')
    global KO_FONT
    _orig_font = KO_FONT
    KO_FONT    = _font

    # ── 베이스 문서: .docx 서식 파일이면 스타일째 불러옴 ───────────
    tpl_path = load_config().get('template_path', '')
    if (tpl_path and Path(tpl_path).exists() and
            str(tpl_path).lower().endswith('.docx') and
            fmt.get('source_type') == 'docx'):
        doc = Document(str(tpl_path))
        _bd_clear(doc)
    else:
        doc = Document()
        for p in doc.paragraphs:
            p._element.getparent().remove(p._element)

    # ── 페이지 설정: 방향·크기·여백 ────────────────────────────────
    _orient = fmt.get('orientation', 'portrait')
    _pw     = fmt.get('page_width_cm',  21.0)
    _ph     = fmt.get('page_height_cm', 29.7)
    if _orient == 'portrait':
        _pw, _ph = min(_pw, _ph), max(_pw, _ph)
    else:
        _pw, _ph = max(_pw, _ph), min(_pw, _ph)

    for sec in doc.sections:
        sec.orientation   = (WD_ORIENT.LANDSCAPE if _orient == 'landscape'
                              else WD_ORIENT.PORTRAIT)
        sec.page_width    = Cm(_pw)
        sec.page_height   = Cm(_ph)
        sec.top_margin    = Cm(fmt.get('top_margin_cm',    2.0))
        sec.bottom_margin = Cm(fmt.get('bottom_margin_cm', 2.0))
        sec.left_margin   = Cm(fmt.get('left_margin_cm',   1.5))
        sec.right_margin  = Cm(fmt.get('right_margin_cm',  1.5))

    C = ['①','②','③','④','⑤']

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 1. 페이지 헤더 (모든 페이지 자동 반복 — 기본값 고정)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    _setup_page_header(doc)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 2. 섹션 제목 (본문 첫 줄)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    _p(doc, SEC_LABELS[section_key], 'section')

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 3. 지문별 문제  ← 고정: 지문마다 새 페이지
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    for pas_idx, pas in enumerate(section_data['passages']):

        if pas_idx > 0:            # ← 고정: 지문별 페이지 나누기
            _add_page_break(doc)

        _p(doc, f'지문 {pas["num"]}', 'passage')

        prev_bt = None

        for blk in pas['blocks']:
            bt = blk['type']

            # ← 고정: 블록 유형 변경 시 구분 공백
            if prev_bt is not None and prev_bt != bt:
                sep = doc.add_paragraph()
                sep.paragraph_format.space_before = Pt(0)
                sep.paragraph_format.space_after  = Pt(8)
            prev_bt = bt

            # ── T/F 블록 ──────────────────────────────────────────
            if bt == 'tf':
                _p(doc, blk['header'], 'blk_header')
                for num, body_txt in blk['items']:
                    clean = re.sub(r'\s*\(\s*T\s*/\s*F\s*\)\s*$', '', body_txt).strip()
                    _p(doc, f'{num}. {clean}  ( T / F )', 'tf_item')

            # ── 한국어/영어 객관식 ──────────────────────────────────
            elif bt in ('ko_mc', 'en_mc'):
                _p(doc, blk['header'], 'mc_blk_header')
                for num, q, chs in blk['items']:
                    _p(doc, f'{num}. {q}', 'mc_body')
                    for si, ch in enumerate(chs):
                        _p(doc, f'{C[si]} {ch}', 'mc_choice')

            # ── 키워드 코너 (Standard 전용) ───────────────────────
            elif bt == 'keyword':
                intro = re.sub(r'^[📝✏️\s━]+', '', blk['intro']).strip()
                _p(doc, f'[{intro}]', 'keyword_hdr')
                if blk.get('table'):
                    tbl = blk['table']
                    header = tbl[0] if tbl else []
                    # 구분선·빈 행 제외: 첫 열이 비어있거나 대시/밑줄만 있는 행 필터
                    data_rows = [
                        list(r) for r in tbl[1:]
                        if r and r[0].strip()
                        and not re.match(r'^[-─━=_\s]+$', r[0].strip())
                    ]
                    if not data_rows:
                        pass  # 표 데이터 없으면 헤더만 있는 셈 → 출력 생략
                    else:
                        # 오른쪽 열 인덱스: 헤더 기준 마지막 텍스트 열
                        ri = max((j for j, c in enumerate(header) if c.strip()), default=len(header)-1)
                        ri = min(ri, len(min(data_rows, key=len)) - 1) if data_rows else 2
                        right = [r[ri] if len(r) > ri else '' for r in data_rows]
                        # 2가지 이상 값이면 셔플
                        if len(set(right)) > 1:
                            shuffled = right[:]
                            for _ in range(200):
                                random.shuffle(shuffled)
                                if shuffled != right:
                                    break
                            for idx, r in enumerate(data_rows):
                                if len(r) > ri: r[ri] = shuffled[idx]
                        _content_table(doc, [header] + data_rows)

            # ── Chunk 코너 (Standard 전용) ────────────────────────
            elif bt == 'chunk':
                intro = re.sub(r'^[✍️\s]+', '',
                               blk.get('intro') or '핵심 Chunk 적어보기').strip()
                _p(doc, f'[{intro}]', 'chunk_hdr')
                for it in blk.get('items', []):
                    num, _, hint = it
                    si  = int(num) - 1
                    sym = C[si] if si < 5 else f'{num}.'
                    _p(doc, f'{sym} _________________________  ({hint})',
                       'chunk_item')

            # ── Hard 추가 영어 ────────────────────────────────────
            elif bt == 'hard_add':
                _p(doc, blk['header'], 'hard_blk_header')
                for num, q, score, chs in blk['items']:
                    tag = f'  [{score}]' if score else ''
                    _p(doc, f'{num}. {q}{tag}', 'hard_body')
                    for si, ch in enumerate(chs):
                        _p(doc, f'{C[si]} {ch}', 'hard_choice')

            # ── Hard 융합 논술 ────────────────────────────────────
            elif bt == 'hard_wr':
                _p(doc, blk['header'], 'hard_blk_header')
                if blk.get('instruction'):
                    _p(doc, blk['instruction'], 'wr_instr')
                if blk.get('template'):
                    g_wt = _g('wr_template')
                    _bd_para(doc, blk['template'], g_wt['size'],
                             bold=g_wt['bold'], sb=g_wt['sb'], sa=g_wt['sa'],
                             ind=g_wt['ind'] or 0.3)
                if blk.get('score_line'):
                    _p(doc, blk['score_line'], 'wr_score')

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 4. 정답표 (새 페이지)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    _add_page_break(doc)
    _p(doc, '정  답  표', 'ans_title')

    for pas in section_data['passages']:
        if not pas.get('answer_table'): continue
        _p(doc, f'지문 {pas["num"]}', 'ans_passage')
        _content_table(doc, pas['answer_table'])

    # ── 저장 ──────────────────────────────────────────────────────
    out_path = output_dir / f'{section_key}.docx'
    doc.save(out_path)
    KO_FONT = _orig_font
    return out_path


def run_conversion(input_path, output_dir, cb, fmt=None):
    try:
        cb('log', '문서 파싱 중...')
        sections = parse_docx(input_path)
        cb('log', '파싱 완료')
        output_dir.mkdir(parents=True, exist_ok=True)
        order = [('easy','easy.docx'), ('standard_a','standard_a.docx'),
                 ('standard_b','standard_b.docx'), ('hard','hard.docx')]
        for key, _ in order:
            if not sections[key]['passages']:
                cb('skip', key); continue
            cb('start', key)
            try:
                out = build_doc(key, sections[key], output_dir, fmt=fmt)
                cb('done', (key, str(out)))
            except Exception as e:
                cb('error', (key, str(e)))
        cb('finished', None)
    except Exception as e:
        cb('fatal', str(e))


# ══════════════════════════════════════════════════════════════════
# GUI 헬퍼
# ══════════════════════════════════════════════════════════════════
# 학생 데이터 관리
# ══════════════════════════════════════════════════════════════════

def _new_id():
    return datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')


def load_students():
    if STUDENTS_FILE.exists():
        try:
            return json.loads(STUDENTS_FILE.read_text(encoding='utf-8'))
        except Exception:
            pass
    return []


def save_students(data):
    STUDENTS_FILE.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')


def load_scores():
    if SCORES_FILE.exists():
        try:
            return json.loads(SCORES_FILE.read_text(encoding='utf-8'))
        except Exception:
            pass
    return []


def save_scores(data):
    SCORES_FILE.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')


# ── Google 인증 / Drive / Sheets ──────────────────────────────────

TOKEN_FILE = DATA_DIR / 'google_token.json'   # _set_account()가 덮어씀
GOOGLE_SCOPES = [
    'https://www.googleapis.com/auth/spreadsheets',
    'https://www.googleapis.com/auth/drive.readonly',
]


def _google_creds(client_secrets: str):
    """OAuth2 자격증명 반환. 토큰이 없으면 브라우저 인증 실행."""
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from google_auth_oauthlib.flow import InstalledAppFlow

    creds = None
    if TOKEN_FILE.exists():
        try:
            creds = Credentials.from_authorized_user_file(
                str(TOKEN_FILE), GOOGLE_SCOPES)
        except Exception:
            creds = None

    if creds and creds.expired and creds.refresh_token:
        try:
            creds.refresh(Request())
            TOKEN_FILE.write_text(creds.to_json(), encoding='utf-8')
        except Exception:
            creds = None

    if not creds or not creds.valid:
        flow = InstalledAppFlow.from_client_secrets_file(
            client_secrets, GOOGLE_SCOPES)
        creds = flow.run_local_server(port=0)
        TOKEN_FILE.write_text(creds.to_json(), encoding='utf-8')

    return creds


def _google_email():
    """저장된 토큰에서 이메일 추출 (없으면 None)."""
    if not TOKEN_FILE.exists():
        return None
    try:
        import base64 as _b64
        d = json.loads(TOKEN_FILE.read_text(encoding='utf-8'))
        id_token = d.get('id_token', '')
        if id_token:
            payload = id_token.split('.')[1]
            payload += '=' * (-len(payload) % 4)
            info = json.loads(_b64.urlsafe_b64decode(payload))
            return info.get('email')
    except Exception:
        pass
    return '(로그인됨)'


def list_drive_sheets(creds):
    """Google Drive에서 스프레드시트 목록 반환 [(name, id, date), ...].
    googleapiclient 대신 gspread 내장 API 사용 (네임스페이스 충돌 방지).
    """
    import gspread
    gc    = gspread.authorize(creds)
    files = gc.list_spreadsheet_files()   # gspread 6.x+
    return [(f['name'], f['id'],
             f.get('modifiedTime', '')[:10])
            for f in files]


def read_sheet_problem_types(creds, sheet_id):
    """스프레드시트에서 문제 유형 목록 추출.

    시트 구조: A열=구분(번호), B열=유형(이름), 1행=헤더
    → B열(col 2)의 2행~끝까지 읽고 빈 셀 제외.
    """
    import gspread
    gc = gspread.authorize(creds)
    sh = gc.open_by_key(sheet_id)

    # '문제유형' 탭 우선, 없으면 첫 번째 탭
    ws = None
    for candidate in ['문제유형', 'ProblemTypes', '유형']:
        try:
            ws = sh.worksheet(candidate)
            break
        except Exception:
            pass
    if ws is None:
        ws = sh.get_worksheet(0)

    # B열(index 2) 읽기, 헤더(1행) 제외
    col_b = ws.col_values(2)          # ['유형', '인지편향', '', ...]
    return [r.strip() for r in col_b[1:] if r.strip()]


# ══════════════════════════════════════════════════════════════════

class _Btn:
    """호버·비활성 상태를 지원하는 커스텀 버튼."""
    def __init__(self, parent, text, cmd, fg_n=None, fg_h=None,
                 bg_n=None, bg_h=None, font=None, pad=(12, 6)):
        self._cmd = cmd; self._on = True
        fn = fg_n or SUB; fh = fg_h or TEXT
        bn = bg_n or CARD; bh = bg_h or HOV
        self._fn = fn; self._fh = fh; self._bn = bn; self._bh = bh

        outer = tk.Frame(parent, bg=BDR, bd=0)
        self._outer = outer
        self._lbl = tk.Label(outer, text=text, font=font or (KO, 9),
                             bg=bn, fg=fn, padx=pad[0], pady=pad[1],
                             cursor='hand2')
        self._lbl.pack(padx=1, pady=1)
        self._lbl.bind('<Enter>', lambda _: self._lbl.configure(bg=bh, fg=fh) if self._on else None)
        self._lbl.bind('<Leave>', lambda _: self._lbl.configure(bg=bn, fg=fn) if self._on else None)
        self._lbl.bind('<Button-1>', lambda _: self._fire())

    def _fire(self):
        if self._on: self._cmd()

    def pack(self, **kw):  self._outer.pack(**kw)
    def grid(self, **kw):  self._outer.grid(**kw)
    def place(self, **kw): self._outer.place(**kw)

    def set_enabled(self, enabled):
        self._on = enabled
        self._lbl.configure(
            fg=self._fn if enabled else DIM,
            cursor='hand2' if enabled else '',
            bg=self._bn)

    def set_text(self, t): self._lbl.configure(text=t)


# ══════════════════════════════════════════════════════════════════
# 메인 앱
# ══════════════════════════════════════════════════════════════════

FILE_META = {
    'easy':       ('📗', 'Easy 시험지',      '#2d6a4f'),
    'standard_a': ('📘', 'Standard A 시험지', '#1a4a8a'),
    'standard_b': ('📙', 'Standard B 시험지', '#7b4a00'),
    'hard':       ('📕', 'Hard 시험지',       '#6b2020'),
}

# 메뉴 구조: (key, icon, 이름, 활성화, sub_items)
# sub_items: [(key, icon, 이름, 활성화), ...]
MENU_ITEMS = [
    ('home',      '🏠', '대시보드',    True, []),
    ('converter', '📝', '문제 변환',   True, [
        ('conv_files',  '📂', '문제 파일 변환', True),
        ('conv_format', '🎨', '출력 서식 설정', True),
    ]),
    ('students',  '👤', '학생 데이터', True, [
        ('stu_manage', '👤', '학생 관리',       True),
        ('stu_mock',   '📊', '모의고사 · 내신', True),
        ('stu_sheet',  '📋', '워크시트 성적',   True),
    ]),
    ('analysis',  '📌', '약점 분석',   False, []),
    ('report',    '📊', '리포트',      False, []),
    ('settings',  '⚙️',  '설정',        True, [
        ('set_general',  '⚙️',  '일반 설정',  True),
        ('set_connect',  '🔗', '연동 설정',   True),
        ('set_security', '🔒', '보안 설정',   True),
    ]),
]


class ConverterApp(_ROOT):
    def __init__(self):
        super().__init__()
        self.title('English Worksheet System')
        self.geometry('1160x740')
        self.minsize(980, 640)
        self.configure(bg=BG)

        self._sel_file      = None
        self._converting    = False
        self._out_dir       = None
        self._file_items    = {}
        self._status        = {}
        self._prog_val      = tk.DoubleVar(value=0)
        self._cur_module    = None
        self._menu_btns     = {}
        self._module_frames = {}
        self._tpl_path      = None   # 현재 서식 파일 경로
        self._tpl_fmt       = {}     # 추출된 서식 설정

        # config.json 경로 확정 및 로드
        global CONFIG_FILE
        CONFIG_FILE = BASE_DIR / 'config.json'
        cfg = load_config()
        tp  = cfg.get('template_path')
        if tp and Path(tp).exists():
            self._tpl_path = Path(tp)
            self._tpl_fmt  = cfg.get('fmt', {})

        self._setup_ttk()
        self._build_ui()

        self._refresh_files()
        self._refresh_score_combos()
        self.bind_all('<MouseWheel>', self._on_mousewheel)
        # 로그인 상태에 따른 문제유형/워크시트 잠금 초기화
        self._apply_custom_types()
        self._refresh_types_display()
        self.after(50, self._update_sheet_google_lock)

        # ── 비밀번호 확인 / 최초 설정 (mainloop 시작 후 실행) ──
        self.withdraw()
        self.after(80, self._do_pw_check)

    # ── TTK 스타일 ──────────────────────────────────────────────

    def _setup_ttk(self):
        s = ttk.Style(self)
        s.theme_use('clam')
        s.configure('Prog.Horizontal.TProgressbar',
                    troughcolor=PRBG, background=PRFG,
                    thickness=5, relief='flat', borderwidth=0)
        s.configure('Vertical.TScrollbar', background=BDR,
                    troughcolor=ITEM, arrowcolor=SUB, relief='flat')

    # ── UI 골격 ──────────────────────────────────────────────────

    def _build_ui(self):
        self._build_header()
        self._build_statusbar()   # side='bottom' 이므로 먼저 배치

        # 메인 영역 (사이드바 + 콘텐츠)
        main = tk.Frame(self, bg=BG)
        main.pack(fill='both', expand=True)

        # 사이드바
        sb = tk.Frame(main, bg=ITEM, width=185)
        sb.pack(side='left', fill='y')
        sb.pack_propagate(False)
        self._build_sidebar(sb)

        tk.Frame(main, bg=BDR, width=1).pack(side='left', fill='y')

        # 콘텐츠 영역
        self._content = tk.Frame(main, bg=BG)
        self._content.pack(side='left', fill='both', expand=True)

        # 모듈 빌드
        self._build_home_page()
        self._build_module_converter()
        self._build_module_students()
        self._build_placeholder('analysis', '📌', '약점 분석',
                                '개인별 취약 유형을 분석하고 맞춤 개선 방향을 제시합니다.')
        self._build_placeholder('report',   '📊', '리포트',
                                '성적 추이, 유형별 정답률 등을 차트로 시각화합니다.')
        self._build_module_settings()

        # 기본 메뉴: 대시보드로 시작
        self._show_module('home')

    # ── 사이드바 ─────────────────────────────────────────────────

    def _build_sidebar(self, parent):
        self._sub_containers = {}   # parent_key → sub Frame
        self._arrow_labels   = {}   # parent_key → arrow Label
        self._expanded       = {}   # parent_key → bool

        tk.Frame(parent, bg=ITEM, height=14).pack(fill='x')
        tk.Label(parent, text='M  E  N  U', font=(KO, 8, 'bold'),
                 bg=ITEM, fg=DIM).pack(anchor='w', padx=18, pady=(0, 6))
        tk.Frame(parent, bg=BDR, height=1).pack(fill='x', padx=14, pady=(0, 6))

        for key, icon, name, active, subs in MENU_ITEMS:
            self._make_menu_item(parent, key, icon, name, active, has_subs=bool(subs))
            if subs:
                # 서브 아이템 컨테이너 (초기: 숨김)
                sub_fr = tk.Frame(parent, bg=ITEM)
                self._sub_containers[key] = sub_fr
                self._expanded[key] = False
                for skey, sicon, sname, sactive in subs:
                    self._make_sub_item(sub_fr, skey, sicon, sname, sactive)

        tk.Frame(parent, bg=ITEM).pack(fill='both', expand=True)

    def _make_menu_item(self, parent, key, icon, name, active, has_subs=False):
        wrap = tk.Frame(parent, bg=ITEM)
        wrap.pack(fill='x', padx=8, pady=(3, 0))
        bar  = tk.Frame(wrap, bg=ITEM, width=3)
        bar.pack(side='left', fill='y')
        # 서브 있는 활성 메뉴: 항상 hand2 (클릭 = 토글)
        cur = 'hand2' if active else ''
        body = tk.Frame(wrap, bg=ITEM, cursor=cur)
        body.pack(side='left', fill='both', expand=True, padx=(6, 8), pady=8)

        ico_box = tk.Frame(body, bg=ITEM, width=30)
        ico_box.pack(side='left', fill='y'); ico_box.pack_propagate(False)
        ico_lbl = tk.Label(ico_box, text=icon, font=(KO, 13), bg=ITEM,
                           fg=TEXT if active else DIM)
        ico_lbl.pack(fill='both', expand=True)

        txt = tk.Frame(body, bg=ITEM)
        txt.pack(side='left', fill='x', expand=True)
        name_lbl = tk.Label(txt, text=name,
                            font=(KO, 9, 'bold' if active else 'normal'),
                            bg=ITEM, fg=TEXT if active else DIM, anchor='w')
        name_lbl.pack(anchor='w', fill='x')
        if not active:
            tk.Label(txt, text='준비 중', font=(KO, 7), bg=ITEM,
                     fg=DIM, anchor='w').pack(anchor='w', fill='x')

        # 서브 있는 항목: 우측에 아코디언 화살표
        arr = None
        if has_subs and active:
            arr = tk.Label(body, text='▶', font=(KO, 8), bg=ITEM, fg=DIM)
            arr.pack(side='right', padx=(0, 2))
            self._arrow_labels[key] = arr

        ws = [wrap, bar, body, ico_box, ico_lbl, txt, name_lbl]
        if arr: ws.append(arr)
        self._menu_btns[key] = {'wrap': wrap, 'bar': bar, 'ws': ws, 'active': active}

        if active:
            cmd = (lambda _, k=key: self._toggle_parent(k)) if has_subs \
                  else (lambda _, k=key: self._show_module(k))
            for w in ws:
                w.bind('<Enter>', lambda _, k=key: self._menu_hover(k, True))
                w.bind('<Leave>', lambda _, k=key: self._menu_hover(k, False))
                w.bind('<Button-1>', cmd)

    def _make_sub_item(self, parent, key, icon, name, active):
        wrap = tk.Frame(parent, bg=ITEM)
        wrap.pack(fill='x', padx=8, pady=0)
        tk.Frame(wrap, bg=ITEM, width=18).pack(side='left', fill='y')
        bar = tk.Frame(wrap, bg=ITEM, width=2)
        bar.pack(side='left', fill='y')
        body = tk.Frame(wrap, bg=ITEM, cursor='hand2' if active else '')
        body.pack(side='left', fill='both', expand=True, padx=(6, 8), pady=5)
        ico_box = tk.Frame(body, bg=ITEM, width=22)
        ico_box.pack(side='left', fill='y'); ico_box.pack_propagate(False)
        ico_lbl = tk.Label(ico_box, text=icon, font=(KO, 10), bg=ITEM,
                           fg=SUB if active else DIM)
        ico_lbl.pack(fill='both', expand=True)
        name_lbl = tk.Label(body, text=name, font=(KO, 8),
                            bg=ITEM, fg=SUB if active else DIM, anchor='w')
        name_lbl.pack(side='left', fill='x', expand=True)
        ws = [wrap, bar, body, ico_box, ico_lbl, name_lbl]
        self._menu_btns[key] = {'wrap': wrap, 'bar': bar, 'ws': ws, 'active': active}
        if active:
            for w in ws:
                w.bind('<Enter>', lambda _, k=key: self._menu_hover(k, True))
                w.bind('<Leave>', lambda _, k=key: self._menu_hover(k, False))
                w.bind('<Button-1>', lambda _, k=key: self._show_module(k))

    # ── 아코디언 토글 ────────────────────────────────────────────

    def _toggle_parent(self, key):
        """상위 메뉴 클릭 시 서브 메뉴 펼치기/접기."""
        sub = self._sub_containers.get(key)
        if sub is None: return

        expanded = self._expanded.get(key, False)

        # 이미 열린 다른 부모 닫기 (one-open accordion)
        for k, sc in self._sub_containers.items():
            if k != key and self._expanded.get(k, False):
                sc.pack_forget()
                self._expanded[k] = False
                if k in self._arrow_labels:
                    self._arrow_labels[k].configure(text='▶')
                # 부모 강조 초기화 (자식이 선택 안 된 경우)
                if self._cur_module not in [sk for sk,_,_,_ in
                        next((s for pk,_,_,_,s in MENU_ITEMS if pk == k), [])]:
                    pb = self._menu_btns.get(k)
                    if pb:
                        pb['bar'].configure(bg=ITEM)
                        for w in pb['ws']:
                            try: w.configure(bg=ITEM)
                            except: pass

        if expanded:
            # 접기
            sub.pack_forget()
            self._expanded[key] = False
            if key in self._arrow_labels:
                self._arrow_labels[key].configure(text='▶')
        else:
            # 펼치기: 부모 wrap 바로 뒤에 삽입
            parent_wrap = self._menu_btns[key]['wrap']
            sub.pack(fill='x', after=parent_wrap)
            self._expanded[key] = True
            if key in self._arrow_labels:
                self._arrow_labels[key].configure(text='▼')

    def _parent_of(self, key):
        """서브 메뉴 키의 상위 메뉴 키 반환."""
        for pk, _, _, _, subs in MENU_ITEMS:
            if any(sk == key for sk, _, _, _ in subs):
                return pk
        return None

    def _menu_hover(self, key, entering):
        if key == self._cur_module: return
        btn = self._menu_btns.get(key)
        if not btn: return
        bg = HOV if entering else ITEM
        for w in btn['ws']:
            try: w.configure(bg=bg)
            except: pass

    def _show_module(self, key):
        # 이전 선택 초기화
        if self._cur_module and self._cur_module in self._menu_btns:
            prev_parent = self._parent_of(self._cur_module)
            for target_key in ([self._cur_module] +
                               ([prev_parent] if prev_parent else [])):
                btn = self._menu_btns.get(target_key)
                if btn:
                    btn['bar'].configure(bg=ITEM)
                    for w in btn['ws']:
                        try: w.configure(bg=ITEM)
                        except: pass

        self._cur_module = key
        btn = self._menu_btns[key]
        for w in btn['ws']:
            try: w.configure(bg=SEL)
            except: pass
        btn['bar'].configure(bg=GOLD)

        # 부모 메뉴도 미묘하게 강조
        parent_key = self._parent_of(key)
        if parent_key and parent_key in self._menu_btns:
            pb = self._menu_btns[parent_key]
            pb['bar'].configure(bg=GOLD)

        # 콘텐츠 프레임 전환
        for k, fr in self._module_frames.items():
            if k == key: fr.pack(fill='both', expand=True)
            else:        fr.pack_forget()

        # 홈 화면 전환 시 자동 새로고침
        if key == 'home' and hasattr(self, '_home_inner'):
            self.after(50, self._refresh_home)

        # 워크시트 성적: Google 로그인 잠금 상태 업데이트
        if key == 'stu_sheet':
            self.after(10, self._update_sheet_google_lock)

    # ── 전역 마우스 휠 ──────────────────────────────────────────

    def _on_mousewheel(self, event):
        """루트 레벨 마우스 휠.
        - 스크롤 가능한 위젯 발견 시: 내용이 넘치면 스크롤, 아니면 무시
        - 어느 경우든 그 위젯에서 완전히 중단 (부모 버블링 없음)."""
        delta = -1 * (event.delta // 120)
        widget = self.winfo_containing(event.x_root, event.y_root)
        while widget:
            if isinstance(widget, (tk.Canvas, tk.Text, ttk.Treeview)):
                try:
                    yv = widget.yview()
                    # yv[1] < 1.0 이면 아래로 스크롤할 내용 있음
                    # yv[0] > 0.0 이면 위로 스크롤할 내용 있음
                    if yv[0] > 0.0 or yv[1] < 1.0:
                        widget.yview_scroll(delta, 'units')
                except Exception:
                    pass
                return   # 스크롤 여부와 무관하게 버블링 차단
            widget = getattr(widget, 'master', None)

    # ── 모듈: 문제 파일 변환 ────────────────────────────────────

    def _build_module_converter(self):
        # ① 문제 파일 변환 화면
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['conv_files'] = fr

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=16, pady=(0, 12))

        lw = tk.Frame(body, bg=BG, width=295)
        lw.pack(side='left', fill='y', padx=(0, 12))
        lw.pack_propagate(False)
        self._build_left(lw)

        tk.Frame(body, bg=BDR, width=1).pack(side='left', fill='y', padx=(0, 12))

        rw = tk.Frame(body, bg=BG)
        rw.pack(side='left', fill='both', expand=True)
        self._build_right(rw)

        # ② 출력 서식 설정 화면
        self._build_format_page()

    # ── 모듈: 준비 중 플레이스홀더 ──────────────────────────────

    def _build_placeholder(self, key, icon, name, desc):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames[key] = fr

        center = tk.Frame(fr, bg=BG)
        center.place(relx=0.5, rely=0.42, anchor='center')

        tk.Label(center, text=icon, font=(KO, 48), bg=BG, fg=DIM).pack()
        tk.Label(center, text=name, font=(KO, 17, 'bold'),
                 bg=BG, fg=SUB).pack(pady=(10, 5))
        tk.Label(center, text=desc, font=(KO, 9),
                 bg=BG, fg=DIM, justify='center').pack()

        # 준비 중 배지
        badge = tk.Frame(center, bg=GOLD)
        badge.pack(pady=16)
        tk.Label(badge, text='  준비 중  ', font=(KO, 9, 'bold'),
                 bg=GOLD, fg='#ffffff', padx=4, pady=5).pack()

    # ══════════════════════════════════════════════════════════════
    # 학생 데이터 모듈
    # ══════════════════════════════════════════════════════════════

    # ══════════════════════════════════════════════════════════════
    # 대시보드 (홈 화면)
    # ══════════════════════════════════════════════════════════════
    # 비밀번호 잠금
    # ══════════════════════════════════════════════════════════════

    @staticmethod
    def _hash_pw(pw: str) -> str:
        return hashlib.sha256(pw.encode('utf-8')).hexdigest()

    @staticmethod
    def _check_pw_strength(pw: str) -> str:
        """비밀번호 강도 검사. 통과 시 '' 반환, 실패 시 오류 메시지."""
        if len(pw) < 8:
            return '⚠️  8자 이상 입력하세요'
        if not re.search(r'[A-Za-z]', pw):
            return '⚠️  영문자(A-Z)를 포함해야 합니다'
        if not re.search(r'[0-9]', pw):
            return '⚠️  숫자(0-9)를 포함해야 합니다'
        if not re.search(r'[!@#$%^&*()\-_=+\[\]{}|;:,.<>?/~`]', pw):
            return '⚠️  특수문자(!@#$ 등)를 포함해야 합니다'
        return ''

    def _do_pw_check(self):
        """mainloop 실행 후 라이센스 → 비밀번호 확인 또는 최초 설정."""
        # ─ 0. 라이센스 키 확인 ──────────────────────────────────
        cfg = load_config()
        if not cfg.get('license_hash'):
            self._lic_ok = False
            lic_win = self._create_license_window()
            self.wait_window(lic_win)
            if not self._lic_ok:
                self.destroy(); return

        self._pw_ok = False
        cfg         = load_config()
        pw_hash     = cfg.get('app_password_hash', '')

        # ① 로그인 or 최초 계정 생성
        win = self._create_pw_login_window(pw_hash) if pw_hash \
              else self._create_pw_setup_window()
        self.wait_window(win)
        if not self._pw_ok:
            self.destroy(); return

        # ② 30일 이상 경과 시 비밀번호 변경 강제
        cfg = load_config()
        last_changed = cfg.get('password_last_changed', '')
        force = False
        if last_changed:
            try:
                days = (datetime.date.today() -
                        datetime.date.fromisoformat(last_changed)).days
                force = days > 30
            except Exception:
                force = False
        else:
            force = bool(pw_hash)   # 날짜 기록 없으면 강제

        if force:
            self._pw_force_ok = False
            cwin = self._create_pw_force_change_window()
            self.wait_window(cwin)
            if not self._pw_force_ok:
                self.destroy(); return

        # ③ 로그인 성공 → 계정별 데이터 디렉토리 활성화 ─────────
        cfg = load_config()
        uid = cfg.get('app_user_id', '')
        if uid:
            _set_account(uid)
            acct = _load_acct_settings()
            tp   = acct.get('template_path', '')
            self._tpl_path = Path(tp) if tp and Path(tp).exists() else None
            self._tpl_fmt  = acct.get('fmt', {})

        # ④ 데이터 저장 위치 미설정 시 → 최초 설정 창 표시 ──────
        global_cfg = {}
        gcf = BASE_DIR / 'config.json'
        if gcf.exists():
            try:
                global_cfg = json.loads(gcf.read_text(encoding='utf-8'))
            except Exception:
                pass
        if 'data_dir' not in global_cfg:
            dwin = self._create_data_dir_setup_window()
            self.wait_window(dwin)

        self.deiconify()

    def _create_data_dir_setup_window(self):
        """최초 로그인 시 데이터 저장 위치 선택 창."""
        win = tk.Toplevel()
        win.title('English Worksheet System — 데이터 저장 위치 설정')
        win.geometry('480x420')
        win.resizable(False, False)
        win.configure(bg=CARD)
        win.attributes('-topmost', True)

        win.update_idletasks()
        sx = (win.winfo_screenwidth()  - 480) // 2
        sy = (win.winfo_screenheight() - 420) // 2
        win.geometry(f'480x420+{sx}+{sy}')

        tk.Frame(win, bg=GOLD, height=4).pack(fill='x')
        inner = tk.Frame(win, bg=CARD)
        inner.pack(expand=True, fill='both', padx=44, pady=22)

        tk.Label(inner, text='📁', font=(KO, 26), bg=CARD, fg=GOLD,
                 anchor='center').pack(fill='x')
        tk.Label(inner, text='데이터 저장 위치를 선택해 주세요',
                 font=(KO, 13, 'bold'), bg=CARD, fg=TEXT, anchor='center'
                 ).pack(fill='x', pady=(8, 4))
        tk.Label(inner, text='학생 성적, 설정 등 모든 데이터가 저장될 폴더입니다.\n'
                              '구글 드라이브나 OneDrive 폴더를 선택하면 자동 백업됩니다.',
                 font=(KO, 8), bg=CARD, fg=SUB, justify='center', anchor='center'
                 ).pack(fill='x', pady=(0, 16))

        # 현재(기본) 경로
        path_var = tk.StringVar(value=str(DATA_DIR))
        tk.Label(inner, text='저장 위치', font=(KO, 8), bg=CARD, fg=SUB,
                 anchor='w').pack(fill='x', pady=(0, 3))

        path_row = tk.Frame(inner, bg=CARD); path_row.pack(fill='x', pady=(0, 14))
        outer_p = tk.Frame(path_row, bg=BDR, bd=1, relief='flat')
        outer_p.pack(side='left', fill='x', expand=True, padx=(0, 8))
        path_lbl = tk.Label(outer_p, textvariable=path_var,
                            font=('Consolas', 8), bg=ITEM, fg=TEXT,
                            anchor='w', padx=6, pady=5)
        path_lbl.pack(fill='x')

        def _browse():
            chosen = filedialog.askdirectory(
                title='데이터 저장 위치 선택',
                initialdir=path_var.get())
            if chosen:
                path_var.set(chosen)

        _Btn(path_row, '📂 변경', _browse,
             font=(KO, 8), pad=(10, 6)).pack(side='left')

        # 안내
        note = tk.Frame(inner, bg=HOV)
        note.pack(fill='x', pady=(0, 18))
        nc = tk.Frame(note, bg=HOV); nc.pack(fill='x', padx=10, pady=8)
        for t in ['💡  구글 드라이브·OneDrive·학교 서버 폴더를 선택하면',
                  '     자동 백업 및 다른 컴퓨터와 공유가 가능합니다.']:
            tk.Label(nc, text=t, font=(KO, 8), bg=HOV, fg=TEXT, anchor='w'
                     ).pack(anchor='w')

        err = tk.Label(inner, text='', font=(KO, 8), bg=CARD, fg=ERR, anchor='center')
        err.pack(fill='x', pady=(0, 4))

        def _apply():
            chosen = Path(path_var.get().strip())
            try:
                chosen.mkdir(parents=True, exist_ok=True)
                # 전역 config에 저장
                gcf = BASE_DIR / 'config.json'
                try:
                    cfg = json.loads(gcf.read_text(encoding='utf-8')) \
                          if gcf.exists() else {}
                except Exception:
                    cfg = {}
                cfg['data_dir'] = str(chosen)
                gcf.write_text(
                    json.dumps(cfg, ensure_ascii=False, indent=2), encoding='utf-8')
                # 전역 변수 업데이트
                _init_custom_data_dir()
                win.destroy()
            except Exception as e:
                err.configure(text=f'❌  {e}')

        _Btn(inner, '  이 위치로 시작하기  ', _apply,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 11, 'bold'), pad=(24, 11)).pack()

        tk.Button(inner, text='나중에 설정하기 (기본 위치 사용)',
                  font=(KO, 8), bg=CARD, fg=DIM, relief='flat',
                  cursor='hand2',
                  command=win.destroy).pack(pady=(10, 0))

        win.protocol('WM_DELETE_WINDOW', win.destroy)
        return win

    def _create_license_window(self):
        """최초 실행 시 라이센스 키 입력 창."""
        win = tk.Toplevel()
        win.title('English Worksheet System — 라이센스 등록')
        win.geometry('440x340')
        win.resizable(False, False)
        win.configure(bg=CARD)
        win.attributes('-topmost', True)

        win.update_idletasks()
        sx = (win.winfo_screenwidth()  - 440) // 2
        sy = (win.winfo_screenheight() - 340) // 2
        win.geometry(f'440x340+{sx}+{sy}')

        tk.Frame(win, bg=GOLD, height=4).pack(fill='x')
        inner = tk.Frame(win, bg=CARD)
        inner.pack(expand=True, fill='both', padx=50, pady=20)

        tk.Label(inner, text='🔑', font=(KO, 26), bg=CARD, fg=GOLD, anchor='center'
                 ).pack(fill='x')
        tk.Label(inner, text='라이센스 등록', font=(KO, 13, 'bold'), bg=CARD, fg=TEXT,
                 anchor='center').pack(fill='x', pady=(6, 2))
        tk.Label(inner, text='발급받은 라이센스 키를 입력하세요\n(형식: EWS-XXXX-XXXX-XXXX-XXXX)',
                 font=(KO, 8), bg=CARD, fg=SUB, anchor='center',
                 justify='center').pack(fill='x', pady=(0, 14))

        outer = tk.Frame(inner, bg=BDR, bd=1, relief='flat')
        outer.pack(fill='x', pady=(0, 4))
        e_key = tk.Entry(outer, font=('Consolas', 11), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT, justify='center')
        e_key.pack(fill='x', padx=6, pady=6)
        e_key.focus_set()

        err = tk.Label(inner, text='', font=(KO, 8), bg=CARD, fg=ERR, anchor='center')
        err.pack(fill='x', pady=(2, 0))

        def _register():
            key = e_key.get().strip().upper()
            if not key.startswith('EWS-') or len(key.split('-')) != 5:
                err.configure(text='❌  올바른 형식이 아닙니다  (EWS-XXXX-XXXX-XXXX-XXXX)')
                return
            if not _validate_license(key):
                err.configure(text='❌  유효하지 않은 라이센스 키입니다'); return

            new_hash = hashlib.sha256(key.encode()).hexdigest()
            cfg = load_config()
            old_hash = cfg.get('license_hash', '')

            # ── 기존과 다른 키 입력 → 전체 설정 + 메모리 초기화 ──
            if old_hash and old_hash != new_hash:
                _reset_user_settings()
                cfg = load_config()
                self._apply_reset_state()   # 메모리·UI 동기화

            cfg['license_hash'] = new_hash
            save_config(cfg)
            self._lic_ok = True
            win.destroy()

        e_key.bind('<Return>', lambda _: _register())
        _Btn(inner, '  라이센스 등록  ', _register,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 11, 'bold'), pad=(24, 10)).pack(pady=(14, 0))

        win.protocol('WM_DELETE_WINDOW', lambda: (
            setattr(self, '_lic_ok', False), win.destroy()))
        return win

    def _create_pw_login_window(self, pw_hash: str):
        """시작 시 아이디 + 비밀번호 로그인 창."""
        win = tk.Toplevel()
        win.title('English Worksheet System')
        win.geometry('400x440')
        win.resizable(False, False)
        win.configure(bg=CARD)
        win.attributes('-topmost', True)

        win.update_idletasks()
        sx = (win.winfo_screenwidth()  - 400) // 2
        sy = (win.winfo_screenheight() - 440) // 2
        win.geometry(f'400x440+{sx}+{sy}')

        tk.Frame(win, bg=GOLD, height=4).pack(fill='x')
        inner = tk.Frame(win, bg=CARD)
        inner.pack(expand=True, fill='both', padx=50, pady=20)

        tk.Label(inner, text='🔒', font=(KO, 26), bg=CARD, fg=GOLD, anchor='center'
                 ).pack(fill='x')
        tk.Label(inner, text='English Worksheet System',
                 font=(KO, 12, 'bold'), bg=CARD, fg=TEXT, anchor='center'
                 ).pack(fill='x', pady=(6, 2))
        tk.Label(inner, text='아이디와 비밀번호를 입력하세요',
                 font=(KO, 8), bg=CARD, fg=SUB, anchor='center'
                 ).pack(fill='x', pady=(0, 14))

        saved_uid = load_config().get('app_user_id', '')

        def _field(label, show=''):
            blk = tk.Frame(inner, bg=CARD); blk.pack(fill='x', pady=4)
            tk.Label(blk, text=label, font=(KO, 8), bg=CARD, fg=SUB, anchor='w'
                     ).pack(fill='x', pady=(0, 2))
            o = tk.Frame(blk, bg=BDR, bd=1, relief='flat'); o.pack(fill='x')
            e = tk.Entry(o, font=(KO, 11), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT,
                         justify='center', show=show)
            e.pack(fill='x', padx=6, pady=4)
            return e

        e_id = _field('아이디')
        e_pw = _field('비밀번호', show='●')
        e_id.insert(0, saved_uid)
        e_id.focus_set()

        err = tk.Label(inner, text='', font=(KO, 8), bg=CARD, fg=ERR, anchor='center')
        err.pack(fill='x', pady=(4, 0))

        def _check():
            uid = e_id.get().strip()
            pw  = e_pw.get()
            if uid != saved_uid:
                err.configure(text='❌  아이디가 올바르지 않습니다'); return
            if self._hash_pw(pw) != pw_hash:
                err.configure(text='❌  비밀번호가 올바르지 않습니다')
                e_pw.delete(0, 'end'); e_pw.focus(); return
            self._pw_ok = True
            win.destroy()

        e_id.bind('<Return>', lambda _: e_pw.focus())
        e_pw.bind('<Return>', lambda _: _check())

        _Btn(inner, '  로그인  ', _check,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 11, 'bold'), pad=(28, 10)).pack(pady=(14, 0))

        # 구분선
        sep = tk.Frame(inner, bg=BG); sep.pack(fill='x', pady=(16, 0))
        tk.Frame(sep, bg=BDR, height=1).pack(fill='x')

        # 회원가입 링크
        signup_row = tk.Frame(inner, bg=CARD); signup_row.pack(pady=(10, 0))
        tk.Label(signup_row, text='계정이 없으신가요?',
                 font=(KO, 8), bg=CARD, fg=SUB).pack(side='left', padx=(0, 6))
        signup_lbl = tk.Label(signup_row, text='회원가입',
                              font=(KO, 8, 'bold'), bg=CARD, fg=GOLD,
                              cursor='hand2')
        signup_lbl.pack(side='left')

        def _go_signup():
            win.destroy()
            setup_win = self._create_pw_setup_window()
            self.wait_window(setup_win)

        signup_lbl.bind('<Enter>', lambda _: signup_lbl.configure(fg=GOLDB))
        signup_lbl.bind('<Leave>', lambda _: signup_lbl.configure(fg=GOLD))
        signup_lbl.bind('<Button-1>', lambda _: _go_signup())

        win.protocol('WM_DELETE_WINDOW', lambda: (
            setattr(self, '_pw_ok', False), win.destroy()))

        return win

    def _create_pw_force_change_window(self):
        """30일 경과 강제 비밀번호 변경 창."""
        cfg     = load_config()
        pw_hash = cfg.get('app_password_hash', '')
        last    = cfg.get('password_last_changed', '알 수 없음')

        win = tk.Toplevel()
        win.title('비밀번호 변경 필요')
        win.geometry('400x420')
        win.resizable(False, False)
        win.configure(bg=CARD)
        win.attributes('-topmost', True)

        win.update_idletasks()
        sx = (win.winfo_screenwidth()  - 400) // 2
        sy = (win.winfo_screenheight() - 420) // 2
        win.geometry(f'400x420+{sx}+{sy}')

        tk.Frame(win, bg=WARN, height=4).pack(fill='x')
        inner = tk.Frame(win, bg=CARD)
        inner.pack(expand=True, fill='both', padx=50, pady=20)

        tk.Label(inner, text='🔔', font=(KO, 26), bg=CARD, fg=WARN, anchor='center'
                 ).pack(fill='x')
        tk.Label(inner, text='비밀번호 변경 기간이 지났습니다',
                 font=(KO, 12, 'bold'), bg=CARD, fg=TEXT, anchor='center'
                 ).pack(fill='x', pady=(6, 2))
        tk.Label(inner, text=f'마지막 변경: {last}  |  30일마다 변경 필수',
                 font=(KO, 8), bg=CARD, fg=SUB, anchor='center'
                 ).pack(fill='x', pady=(0, 14))

        def _field(label, show=''):
            blk = tk.Frame(inner, bg=CARD); blk.pack(fill='x', pady=4)
            tk.Label(blk, text=label, font=(KO, 8), bg=CARD, fg=SUB, anchor='w'
                     ).pack(fill='x', pady=(0, 2))
            o = tk.Frame(blk, bg=BDR, bd=1, relief='flat'); o.pack(fill='x')
            e = tk.Entry(o, show=show, font=(KO, 11), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT, justify='center')
            e.pack(fill='x', padx=6, pady=4)
            return e

        e_cur = _field('현재 비밀번호', show='●')
        e_new = _field('새 비밀번호 (영문+숫자+특수문자, 8자↑)', show='●')
        e_con = _field('새 비밀번호 확인', show='●')
        err   = tk.Label(inner, text='', font=(KO, 8), bg=CARD, fg=ERR, anchor='center')
        err.pack(fill='x', pady=(4, 0))
        e_cur.focus_set()

        def _save():
            cur, pw, pw2 = e_cur.get(), e_new.get(), e_con.get()
            if self._hash_pw(cur) != pw_hash:
                err.configure(text='❌  현재 비밀번호가 올바르지 않습니다'); return
            _e = self._check_pw_strength(pw)
            if _e: err.configure(text=_e); return
            if pw != pw2:
                err.configure(text='❌  새 비밀번호가 일치하지 않습니다'); return
            c = load_config()
            c['app_password_hash']     = self._hash_pw(pw)
            c['password_last_changed'] = datetime.date.today().isoformat()
            save_config(c)
            self._pw_force_ok = True
            win.destroy()

        e_cur.bind('<Return>', lambda _: e_new.focus())
        e_new.bind('<Return>', lambda _: e_con.focus())
        e_con.bind('<Return>', lambda _: _save())

        _Btn(inner, '  비밀번호 변경 후 시작하기  ', _save,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=WARN, bg_h='#d97706',
             font=(KO, 10, 'bold'), pad=(20, 10)).pack(pady=(14, 0))

        win.protocol('WM_DELETE_WINDOW', lambda: (
            setattr(self, '_pw_force_ok', False), win.destroy()))
        return win

    def _create_pw_setup_window(self):
        """최초 실행 시 비밀번호 설정 창."""
        win = tk.Toplevel()
        win.title('English Worksheet System')
        win.geometry('400x520')
        win.resizable(False, False)
        win.configure(bg=CARD)
        win.attributes('-topmost', True)

        win.update_idletasks()
        sx = (win.winfo_screenwidth()  - 400) // 2
        sy = (win.winfo_screenheight() - 520) // 2
        win.geometry(f'400x520+{sx}+{sy}')

        tk.Frame(win, bg=GOLD, height=4).pack(fill='x')

        inner = tk.Frame(win, bg=CARD)
        inner.pack(expand=True, fill='both', padx=50, pady=20)

        # 아이콘 + 제목 (가운데 정렬)
        tk.Label(inner, text='🔐', font=(KO, 26), bg=CARD, fg=GOLD,
                 anchor='center').pack(fill='x')
        tk.Label(inner, text='처음 오셨군요!', font=(KO, 13, 'bold'), bg=CARD, fg=TEXT,
                 anchor='center').pack(fill='x', pady=(6, 2))
        tk.Label(inner, text='사용할 아이디와 비밀번호를 설정하세요',
                 font=(KO, 8), bg=CARD, fg=SUB, anchor='center').pack(fill='x', pady=(0, 12))

        def _field(label, show=''):
            blk = tk.Frame(inner, bg=CARD); blk.pack(fill='x', pady=4)
            tk.Label(blk, text=label, font=(KO, 8), bg=CARD, fg=SUB, anchor='w'
                     ).pack(fill='x', pady=(0, 2))
            outer = tk.Frame(blk, bg=BDR, bd=1, relief='flat'); outer.pack(fill='x')
            e = tk.Entry(outer, show=show, font=(KO, 11), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT, justify='center')
            e.pack(fill='x', padx=6, pady=4)
            return e

        e_id  = _field('아이디')
        e_new = _field('비밀번호 (영문+숫자+특수문자, 8자↑)', show='●')
        e_con = _field('비밀번호 확인', show='●')

        err = tk.Label(inner, text='', font=(KO, 8), bg=CARD, fg=ERR, anchor='center')
        err.pack(fill='x', pady=(4, 0))

        e_id.focus_set()

        def _save():
            uid, pw, pw2 = e_id.get().strip(), e_new.get(), e_con.get()
            if not uid:
                err.configure(text='⚠️  아이디를 입력하세요'); return
            _e = self._check_pw_strength(pw)
            if _e: err.configure(text=_e); return
            if pw != pw2:
                err.configure(text='❌  비밀번호가 일치하지 않습니다'); return

            # ── 기존 계정과 다른 아이디 → 설정 + 메모리 초기화 ─────
            cfg      = load_config()
            old_uid  = cfg.get('app_user_id', '')
            old_lic  = cfg.get('license_hash', '')
            if old_uid and old_uid != uid:
                _reset_user_settings()
                cfg = load_config()
                cfg['license_hash'] = old_lic
                self._apply_reset_state()   # 메모리·UI 동기화

            cfg['app_user_id']           = uid
            cfg['app_password_hash']     = self._hash_pw(pw)
            cfg['password_last_changed'] = datetime.date.today().isoformat()
            save_config(cfg)
            # 처음 가입이어도 메모리·UI 상태 동기화 (잔여값 제거)
            self._apply_reset_state()
            self._pw_ok = True
            win.destroy()

        e_id.bind('<Return>',  lambda _: e_new.focus())
        e_new.bind('<Return>', lambda _: e_con.focus())
        e_con.bind('<Return>', lambda _: _save())

        _Btn(inner, '  가입 완료  ', _save,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 11, 'bold'), pad=(28, 10)).pack(pady=(14, 0))

        win.protocol('WM_DELETE_WINDOW', lambda: (
            setattr(self, '_pw_ok', False), win.destroy()))

        return win

    # ── 보안 설정 페이지 ────────────────────────────────────────

    def _build_set_security_page(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['set_security'] = fr
        self._stu_page_header(fr, '🔒', '보안 설정',
                              '앱 시작 시 비밀번호 잠금을 설정합니다')

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=40, pady=24)

        # ── 현재 상태 카드 ──
        st_outer = tk.Frame(body, bg=BDR); st_outer.pack(fill='x', pady=(0, 16))
        st_inner = tk.Frame(st_outer, bg=CARD); st_inner.pack(fill='both', padx=1, pady=1)
        sc = tk.Frame(st_inner, bg=CARD); sc.pack(fill='x', padx=20, pady=14)

        tk.Label(sc, text='잠금 상태', font=(KO, 8, 'bold'), bg=CARD, fg=SUB
                 ).pack(anchor='w', pady=(0, 6))
        self._pw_status_lbl = tk.Label(sc, text='', font=(KO, 10, 'bold'), bg=CARD, fg=TEXT)
        self._pw_status_lbl.pack(anchor='w')
        self._update_pw_status()

        # ── 비밀번호 설정/변경 카드 ──
        pw_outer = tk.Frame(body, bg=BDR); pw_outer.pack(fill='x', pady=(0, 14))
        pw_inner = tk.Frame(pw_outer, bg=CARD); pw_inner.pack(fill='both', padx=1, pady=1)
        pc = tk.Frame(pw_inner, bg=CARD); pc.pack(fill='x', padx=20, pady=18)

        tk.Label(pc, text='비밀번호 설정 / 변경', font=(KO, 9, 'bold'),
                 bg=CARD, fg=TEXT).pack(anchor='w', pady=(0, 12))

        def pw_row(label, show='●'):
            r = tk.Frame(pc, bg=CARD); r.pack(fill='x', pady=4)
            tk.Label(r, text=label, font=(KO, 8), bg=CARD, fg=SUB,
                     width=12, anchor='w').pack(side='left', padx=(0, 8))
            outer_e = tk.Frame(r, bg=BDR, bd=1, relief='flat')
            outer_e.pack(side='left', fill='x', expand=True)
            e = tk.Entry(outer_e, show=show, font=(KO, 10), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT)
            e.pack(fill='x', padx=4, pady=3)
            return e

        cfg = load_config()
        has_pw = bool(cfg.get('app_password_hash', ''))

        if has_pw:
            self._e_cur  = pw_row('현재 비밀번호')
        else:
            self._e_cur  = None
        self._e_new  = pw_row('새 비밀번호')
        self._e_con  = pw_row('비밀번호 확인')

        self._pw_msg = tk.Label(pc, text='', font=(KO, 8), bg=CARD, fg=ERR)
        self._pw_msg.pack(anchor='w', pady=(6, 0))

        btn_row2 = tk.Frame(pc, bg=CARD); btn_row2.pack(fill='x', pady=(12, 0))
        _Btn(btn_row2, '💾  저장',
             self._save_password,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(18, 8)).pack(side='left')

        if has_pw:
            _Btn(btn_row2, '🔓  잠금 해제 (비밀번호 삭제)',
                 self._remove_password,
                 font=(KO, 8), pad=(14, 8)).pack(side='left', padx=(10, 0))

        # ── 아이디 변경 카드 ──────────────────────────────────────
        id_outer = tk.Frame(body, bg=BDR); id_outer.pack(fill='x', pady=(0, 14))
        id_inner = tk.Frame(id_outer, bg=CARD); id_inner.pack(fill='both', padx=1, pady=1)
        ic = tk.Frame(id_inner, bg=CARD); ic.pack(fill='x', padx=20, pady=18)

        cfg2 = load_config()
        cur_uid = cfg2.get('app_user_id', '(없음)')
        tk.Label(ic, text='아이디 변경', font=(KO, 9, 'bold'), bg=CARD, fg=TEXT
                 ).pack(anchor='w', pady=(0, 4))
        tk.Label(ic, text=f'현재 아이디: {cur_uid}', font=(KO, 8), bg=CARD, fg=SUB
                 ).pack(anchor='w', pady=(0, 10))

        def id_row(label, show=''):
            r = tk.Frame(ic, bg=CARD); r.pack(fill='x', pady=3)
            tk.Label(r, text=label, font=(KO, 8), bg=CARD, fg=SUB,
                     width=12, anchor='w').pack(side='left', padx=(0, 8))
            o = tk.Frame(r, bg=BDR, bd=1, relief='flat')
            o.pack(side='left', fill='x', expand=True)
            e = tk.Entry(o, show=show, font=(KO, 10), bg=CARD, fg=TEXT,
                         relief='flat', bd=0, insertbackground=TEXT)
            e.pack(fill='x', padx=4, pady=3)
            return e

        self._e_id_pw  = id_row('현재 비밀번호', show='●')
        self._e_new_id = id_row('새 아이디')
        self._id_msg   = tk.Label(ic, text='', font=(KO, 8), bg=CARD, fg=ERR)
        self._id_msg.pack(anchor='w', pady=(4, 0))
        _Btn(ic, '💾  아이디 변경', self._save_user_id,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(16, 8)).pack(anchor='w', pady=(10, 0))

        # ── 전체 데이터 초기화 카드 ───────────────────────────────
        rs_outer = tk.Frame(body, bg=ERR); rs_outer.pack(fill='x')
        rs_inner = tk.Frame(rs_outer, bg=CARD); rs_inner.pack(fill='both', padx=1, pady=1)
        rc = tk.Frame(rs_inner, bg=CARD); rc.pack(fill='x', padx=20, pady=16)

        tk.Label(rc, text='⚠️  전체 데이터 초기화', font=(KO, 9, 'bold'), bg=CARD, fg=ERR
                 ).pack(anchor='w', pady=(0, 4))
        tk.Label(rc, text='학생 정보, 성적 기록, 계정 정보가 모두 삭제됩니다.\n이 작업은 되돌릴 수 없습니다.',
                 font=(KO, 8), bg=CARD, fg=SUB).pack(anchor='w', pady=(0, 10))

        rst_row = tk.Frame(rc, bg=CARD); rst_row.pack(fill='x')
        tk.Label(rst_row, text='현재 비밀번호', font=(KO, 8), bg=CARD, fg=SUB,
                 width=12, anchor='w').pack(side='left', padx=(0, 8))
        o2 = tk.Frame(rst_row, bg=BDR, bd=1, relief='flat')
        o2.pack(side='left', fill='x', expand=True)
        self._e_reset_pw = tk.Entry(o2, show='●', font=(KO, 10), bg=CARD, fg=TEXT,
                                    relief='flat', bd=0, insertbackground=TEXT)
        self._e_reset_pw.pack(fill='x', padx=4, pady=3)

        self._reset_msg = tk.Label(rc, text='', font=(KO, 8), bg=CARD, fg=ERR)
        self._reset_msg.pack(anchor='w', pady=(4, 0))
        _Btn(rc, '🗑  전체 초기화 (되돌릴 수 없음)', self._reset_all_data,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=ERR, bg_h='#b91c1c',
             font=(KO, 8, 'bold'), pad=(14, 7)).pack(anchor='w', pady=(10, 0))

    def _update_pw_status(self):
        if not hasattr(self, '_pw_status_lbl'): return
        cfg = load_config()
        if cfg.get('app_password_hash'):
            self._pw_status_lbl.configure(text='🔒  잠금 설정됨  —  앱 시작 시 비밀번호가 필요합니다', fg=OK)
        else:
            self._pw_status_lbl.configure(text='🔓  잠금 없음  —  누구나 앱을 바로 실행할 수 있습니다', fg=DIM)

    def _save_password(self):
        cfg     = load_config()
        has_pw  = bool(cfg.get('app_password_hash', ''))
        new_pw  = self._e_new.get()
        con_pw  = self._e_con.get()

        # 현재 비밀번호 검증 (이미 설정된 경우)
        if has_pw and self._e_cur:
            cur_pw = self._e_cur.get()
            if self._hash_pw(cur_pw) != cfg['app_password_hash']:
                self._pw_msg.configure(text='❌  현재 비밀번호가 올바르지 않습니다', fg=ERR)
                return

        _e = self._check_pw_strength(new_pw)
        if _e:
            self._pw_msg.configure(text=_e, fg=WARN)
            return
        if new_pw != con_pw:
            self._pw_msg.configure(text='❌  새 비밀번호와 확인이 일치하지 않습니다', fg=ERR)
            return

        cfg['app_password_hash']     = self._hash_pw(new_pw)
        cfg['password_last_changed'] = datetime.date.today().isoformat()
        save_config(cfg)
        self._update_pw_status()
        self._pw_msg.configure(text='✅  비밀번호가 저장되었습니다', fg=OK)
        self._lbl_status.configure(text='✅  비밀번호 설정 완료', fg=OK)

        # 입력창 초기화
        for e in (self._e_cur, self._e_new, self._e_con):
            if e: e.delete(0, 'end')

    def _remove_password(self):
        cfg = load_config()
        if self._e_cur and self._hash_pw(self._e_cur.get()) != cfg.get('app_password_hash',''):
            self._pw_msg.configure(text='❌  현재 비밀번호가 올바르지 않습니다', fg=ERR)
            return
        cfg.pop('app_password_hash', None)
        cfg.pop('app_user_id', None)
        cfg.pop('password_last_changed', None)
        save_config(cfg)
        self._update_pw_status()
        self._pw_msg.configure(text='✅  잠금이 해제되었습니다', fg=OK)
        self._lbl_status.configure(text='🔓  비밀번호 잠금 해제됨', fg=SUB)

    def _save_user_id(self):
        """아이디 변경 저장."""
        cfg     = load_config()
        pw_hash = cfg.get('app_password_hash', '')
        cur_pw  = self._e_id_pw.get()
        new_id  = self._e_new_id.get().strip()

        if pw_hash and self._hash_pw(cur_pw) != pw_hash:
            self._id_msg.configure(text='❌  현재 비밀번호가 올바르지 않습니다', fg=ERR); return
        if not new_id:
            self._id_msg.configure(text='⚠️  새 아이디를 입력하세요', fg=WARN); return
        cfg['app_user_id'] = new_id
        save_config(cfg)
        self._id_msg.configure(text=f'✅  아이디가 "{new_id}"로 변경되었습니다', fg=OK)
        self._lbl_status.configure(text=f'✅  아이디 변경 완료: {new_id}', fg=OK)
        self._e_id_pw.delete(0, 'end')
        self._e_new_id.delete(0, 'end')

    def _reset_all_data(self):
        """모든 데이터 및 계정 초기화 (비밀번호 확인 후)."""
        cfg     = load_config()
        pw_hash = cfg.get('app_password_hash', '')
        if pw_hash and self._hash_pw(self._e_reset_pw.get()) != pw_hash:
            self._reset_msg.configure(text='❌  비밀번호가 올바르지 않습니다', fg=ERR); return

        # 현재 계정 데이터 파일 삭제
        for f in (STUDENTS_FILE, SCORES_FILE, TOKEN_FILE):
            try:
                if f.exists(): f.unlink()
            except Exception: pass
        # 현재 계정 settings 삭제
        sf = _acct_settings_file()
        if sf and sf.exists():
            try: sf.unlink()
            except Exception: pass
        # 글로벌 계정 인증 초기화
        _pop_config('app_user_id', 'app_password_hash', 'password_last_changed')
        self._reset_msg.configure(text='✅  초기화 완료 — 앱을 재시작하세요', fg=OK)
        self._lbl_status.configure(text='✅  전체 초기화 완료 — 재시작 필요', fg=WARN)

    # ══════════════════════════════════════════════════════════════

    def _build_home_page(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['home'] = fr

        # 헤더
        hdr = tk.Frame(fr, bg=BG2)
        hdr.pack(fill='x')
        tk.Frame(hdr, bg=GOLD, height=2).pack(side='bottom', fill='x')
        ph = tk.Frame(hdr, bg=BG2)
        ph.pack(fill='x', padx=24, pady=12)
        tk.Label(ph, text='🏠', font=(KO, 18), bg=BG2, fg=GOLD).pack(side='left', padx=(0, 10))
        tl = tk.Frame(ph, bg=BG2)
        tl.pack(side='left')
        tk.Label(tl, text='대시보드', font=(KO, 13, 'bold'), bg=BG2, fg=TEXT).pack(anchor='w')
        tk.Label(tl, text='학생 현황과 성적을 한눈에 확인합니다',
                 font=(KO, 8), bg=BG2, fg=SUB).pack(anchor='w')
        _Btn(ph, '🔄  새로고침', self._refresh_home,
             font=(KO, 8), pad=(10, 5)).pack(side='right')

        # 스크롤 캔버스
        cv = tk.Canvas(fr, bg=BG, highlightthickness=0)
        sb = ttk.Scrollbar(fr, orient='vertical', command=cv.yview,
                           style='Vertical.TScrollbar')
        cv.configure(yscrollcommand=sb.set)
        sb.pack(side='right', fill='y')
        cv.pack(fill='both', expand=True)
        inner = tk.Frame(cv, bg=BG)
        cwin  = cv.create_window((0, 0), window=inner, anchor='nw')
        def _on_cfg(e): cv.configure(scrollregion=(0, 0, e.width, e.height))
        inner.bind('<Configure>', _on_cfg)
        cv.bind('<Configure>', lambda e: cv.itemconfig(cwin, width=e.width))

        self._home_inner  = inner
        self._home_canvas = cv
        self._home_fill()

    def _refresh_home(self):
        for w in self._home_inner.winfo_children():
            w.destroy()
        self._home_fill()
        self.after(10, lambda: (
            self._home_canvas.configure(
                scrollregion=(0, 0,
                              self._home_inner.winfo_reqwidth(),
                              self._home_inner.winfo_reqheight())),
            self._home_canvas.yview_moveto(0)
        ))

    # ── 대시보드 컨텐츠 ──────────────────────────────────────────

    def _home_fill(self):
        """대시보드 전체 내용 빌드 (데이터 기반 동적 생성)."""
        body = self._home_inner
        students = load_students()
        scores   = load_scores()
        mock_s   = [r for r in scores if r.get('category') == 'mock']
        sheet_s  = [r for r in scores if r.get('category') == 'sheet']
        stu_map  = {s['id']: s for s in students}

        # ── 1. 요약 스탯 카드 (3개) ────────────────────────────────
        stats_row = tk.Frame(body, bg=BG)
        stats_row.pack(fill='x', padx=20, pady=(16, 8))
        stats_row.columnconfigure(0, weight=1)
        stats_row.columnconfigure(1, weight=1)
        stats_row.columnconfigure(2, weight=1)

        def stat_card(parent, col, icon, val, label, accent):
            o = tk.Frame(parent, bg=BDR); o.grid(row=0, column=col, padx=5, sticky='nsew')
            i = tk.Frame(o, bg=CARD);     i.pack(fill='both', padx=1, pady=1)
            tk.Frame(i, bg=accent, height=4).pack(fill='x')
            c = tk.Frame(i, bg=CARD);     c.pack(fill='x', padx=16, pady=14)
            tk.Label(c, text=icon, font=(KO, 20), bg=CARD, fg=accent).pack(anchor='w')
            tk.Label(c, text=str(val), font=(KO, 22, 'bold'), bg=CARD, fg=TEXT
                     ).pack(anchor='w', pady=(4, 0))
            tk.Label(c, text=label, font=(KO, 8), bg=CARD, fg=SUB).pack(anchor='w')

        stat_card(stats_row, 0, '👥', len(students),  '등록 학생',      '#1a4a8a')
        stat_card(stats_row, 1, '📊', len(mock_s),    '모의고사 기록',  '#2d6a4f')
        stat_card(stats_row, 2, '📋', len(sheet_s),   '워크시트 기록',  '#7b4a00')

        # ── 2. 학생별 현황 테이블 ──────────────────────────────────
        self._stu_section(body, '학생 현황')
        tk.Frame(body, bg=BG, height=2).pack()

        if not students:
            tk.Label(body, text='등록된 학생이 없습니다. 학생 데이터 → 학생 관리에서 추가하세요.',
                     font=(KO, 9), bg=BG, fg=DIM).pack(anchor='w', padx=20, pady=8)
        else:
            # 테이블 헤더
            def tbl_row(parent, cells, bold=False, bg_c=ITEM):
                row = tk.Frame(parent, bg=bg_c)
                row.pack(fill='x', padx=20, pady=1)
                widths = [100, 120, 90, 80, 80, 80]
                for txt, w in zip(cells, widths):
                    tk.Label(row, text=txt,
                             font=(KO, 8, 'bold' if bold else 'normal'),
                             bg=bg_c, fg=TEXT if bold else TEXT,
                             width=int(w/7), anchor='w'
                             ).pack(side='left', padx=4, pady=5)

            tbl_row(body, ['이름', '학교', '학년/반', '최근 모의고사', '등급', '워크시트'], bold=True, bg_c=BDR)

            for s in students:
                sid = s['id']
                # 최근 모의고사
                stu_mock = sorted(
                    [r for r in mock_s if r.get('student_id') == sid],
                    key=lambda r: r.get('date', ''), reverse=True)
                mock_txt  = stu_mock[0].get('score', '—') if stu_mock else '—'
                grade_txt = stu_mock[0].get('grade', '—') if stu_mock else '—'
                # 최근 워크시트
                stu_ws = sorted(
                    [r for r in sheet_s if r.get('student_id') == sid],
                    key=lambda r: r.get('date', ''), reverse=True)
                ws_txt = stu_ws[0].get('total_score', '—') if stu_ws else '—'

                tbl_row(body, [
                    s['name'],
                    s.get('school', '—')[:8],
                    s.get('grade_class', '—'),
                    str(mock_txt),
                    str(grade_txt),
                    str(ws_txt),
                ], bg_c=CARD)

        # ── 3. 문제 유형별 평균 점수 (데이터 있을 때만) ────────────
        if sheet_s:
            self._stu_section(body, '문제 유형별 평균 점수')
            type_scores = defaultdict(list)
            for r in sheet_s:
                qt = r.get('qtype', '')
                try:
                    sc = float(r.get('total_score', ''))
                    if qt: type_scores[qt].append(sc)
                except (ValueError, TypeError):
                    pass

            if type_scores:
                chart_data = [(qt, sum(vs)/len(vs)) for qt, vs in type_scores.items()]
                chart_data.sort(key=lambda x: -x[1])
                self._draw_bar_chart(body, chart_data)
            else:
                tk.Label(body, text='점수 데이터가 부족합니다.',
                         font=(KO, 9), bg=BG, fg=DIM).pack(anchor='w', padx=20, pady=6)

        # ── 4. 최근 모의고사 기록 (있을 때만) ─────────────────────
        if mock_s:
            self._stu_section(body, '최근 모의고사 기록 (최근 5건)')
            recent = sorted(mock_s, key=lambda r: r.get('date',''), reverse=True)[:5]
            cols = ['날짜', '학생', '구분', '시험명', '원점수', '등급', '백분위']
            widths = [80, 70, 80, 140, 55, 40, 55]
            # 헤더
            hr = tk.Frame(body, bg=BDR); hr.pack(fill='x', padx=20, pady=(4,1))
            for h, w in zip(cols, widths):
                tk.Label(hr, text=h, font=(KO, 8, 'bold'), bg=BDR, fg=TEXT,
                         width=int(w/7), anchor='w').pack(side='left', padx=4, pady=4)
            for r in recent:
                dr = tk.Frame(body, bg=CARD); dr.pack(fill='x', padx=20, pady=1)
                vals = [r.get('date',''), stu_map.get(r.get('student_id',''), {}).get('name','?'),
                        r.get('exam_type',''), r.get('exam_name','')[:10],
                        r.get('score',''), r.get('grade',''), r.get('percentile','')]
                for v, w in zip(vals, widths):
                    tk.Label(dr, text=str(v), font=(KO, 8), bg=CARD, fg=TEXT,
                             width=int(w/7), anchor='w').pack(side='left', padx=4, pady=4)

        # ── 5. 개발 예정 섹션 ────────────────────────────────────
        self._stu_section(body, '추가 기능 (개발 예정)')
        coming = tk.Frame(body, bg=BG)
        coming.pack(fill='x', padx=20, pady=(4, 20))
        coming.columnconfigure(0, weight=1); coming.columnconfigure(1, weight=1)

        def coming_card(parent, row, col, icon, title, desc):
            o = tk.Frame(parent, bg=BDR); o.grid(row=row, column=col, padx=5, pady=5, sticky='nsew')
            i = tk.Frame(o, bg=CARD);     i.pack(fill='both', padx=1, pady=1)
            c = tk.Frame(i, bg=CARD);     c.pack(fill='x', padx=14, pady=12)
            top = tk.Frame(c, bg=CARD);   top.pack(fill='x')
            tk.Label(top, text=icon, font=(KO, 16), bg=CARD, fg=DIM).pack(side='left', padx=(0,8))
            tk.Label(top, text=title, font=(KO, 10, 'bold'), bg=CARD, fg=SUB).pack(side='left')
            badge = tk.Frame(top, bg=BDR); badge.pack(side='right')
            tk.Label(badge, text=' 개발 예정 ', font=(KO, 7), bg=BDR,
                     fg=DIM, padx=4, pady=2).pack()
            tk.Label(c, text=desc, font=(KO, 8), bg=CARD, fg=DIM,
                     wraplength=300, justify='left').pack(anchor='w', pady=(6,0))

        coming_card(coming, 0, 0, '🔍', '약점 분석',
                    '문제 유형별 오답 패턴을 분석하여\n개인별 취약점과 개선 방향을 제시합니다.')
        coming_card(coming, 0, 1, '📈', '성적 추이',
                    '시간별 성적 변화를 차트로 시각화하여\n학습 효과와 성장 곡선을 확인합니다.')
        coming_card(coming, 1, 0, '🎯', '개인별 리포트',
                    '학생 한 명의 전체 성적 데이터를\nPDF 리포트로 출력합니다.')
        coming_card(coming, 1, 1, '📊', '비교 분석',
                    '같은 시험의 학생들 간 성적을\n비교하고 순위를 확인합니다.')

    def _draw_bar_chart(self, parent, data):
        """간단한 Canvas 막대 차트."""
        if not data: return
        wrap = tk.Frame(parent, bg=BG)
        wrap.pack(fill='x', padx=20, pady=(4, 8))
        cv_h = 140
        cv = tk.Canvas(wrap, bg=CARD, height=cv_h, highlightthickness=1,
                       highlightbackground=BDR)
        cv.pack(fill='x')

        def _draw(_event=None):
            cv.delete('all')
            w = cv.winfo_width() or 600
            n = len(data)
            if n == 0: return
            pad_l, pad_r, pad_t, pad_b = 10, 10, 20, 30
            bar_area_w = w - pad_l - pad_r
            bar_w      = max(20, bar_area_w // n - 8)
            max_val    = max(v for _, v in data) or 100
            colors     = ['#1a4a8a','#2d6a4f','#7b4a00','#6b2020','#4a1a6b']

            for idx, (label, val) in enumerate(data):
                x0 = pad_l + idx * (bar_w + 8)
                x1 = x0 + bar_w
                bar_h = int((val / max_val) * (cv_h - pad_t - pad_b))
                y0 = cv_h - pad_b - bar_h
                y1 = cv_h - pad_b
                color = colors[idx % len(colors)]
                cv.create_rectangle(x0, y0, x1, y1, fill=color, outline='', width=0)
                # 점수 텍스트
                cv.create_text((x0+x1)//2, y0 - 4, text=f'{val:.0f}',
                               font=(KO, 7), fill=TEXT, anchor='s')
                # 라벨 (4자 이하로 줄임)
                short = label[:5] if len(label) > 5 else label
                cv.create_text((x0+x1)//2, cv_h - pad_b + 4, text=short,
                               font=(KO, 7), fill=SUB, anchor='n')

        cv.bind('<Configure>', _draw)
        cv.after(50, _draw)

    # ══════════════════════════════════════════════════════════════

    def _build_module_students(self):
        self._build_stu_manage()
        self._build_stu_mock()
        self._build_stu_sheet()

    # ══════════════════════════════════════════════════════════════
    # 설정 모듈
    # ══════════════════════════════════════════════════════════════

    def _build_module_settings(self):
        self._build_set_general_page()
        self._build_set_connect_page()
        self._build_set_security_page()

    # ── 일반 설정 페이지 ────────────────────────────────────────

    def _build_set_general_page(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['set_general'] = fr
        self._stu_page_header(fr, '⚙️', '일반 설정',
                              '데이터 저장 위치 등 프로그램 기본 설정을 관리합니다')

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=40, pady=24)

        # ── 데이터 저장 위치 카드 ──
        c_out = tk.Frame(body, bg=BDR); c_out.pack(fill='x', pady=(0, 16))
        c_in  = tk.Frame(c_out, bg=CARD); c_in.pack(fill='both', padx=1, pady=1)
        cc = tk.Frame(c_in, bg=CARD); cc.pack(fill='x', padx=20, pady=18)

        tk.Label(cc, text='데이터 저장 위치', font=(KO, 9, 'bold'),
                 bg=CARD, fg=TEXT).pack(anchor='w', pady=(0, 8))

        # 현재 경로 표시
        path_row = tk.Frame(cc, bg=CARD); path_row.pack(fill='x', pady=(0, 8))
        tk.Label(path_row, text='현재 위치', font=(KO, 8), bg=CARD, fg=SUB,
                 width=8, anchor='w').pack(side='left', padx=(0, 8))
        self._data_dir_lbl = tk.Label(path_row, text=str(DATA_DIR),
                                      font=('Consolas', 8), bg=ITEM, fg=TEXT,
                                      anchor='w', padx=6, pady=4, relief='flat')
        self._data_dir_lbl.pack(side='left', fill='x', expand=True)

        # 크기 정보
        self._data_size_lbl = tk.Label(cc, text='', font=(KO, 8), bg=CARD, fg=SUB)
        self._data_size_lbl.pack(anchor='w', pady=(0, 10))
        self._refresh_data_size()

        # 이동 여부 체크박스
        self._move_data_var = tk.BooleanVar(value=True)
        tk.Checkbutton(cc, text='기존 데이터를 새 위치로 복사',
                       variable=self._move_data_var,
                       font=(KO, 8), bg=CARD, fg=TEXT,
                       activebackground=CARD, selectcolor=CARD
                       ).pack(anchor='w', pady=(0, 10))

        # 버튼 행
        btn_r = tk.Frame(cc, bg=CARD); btn_r.pack(fill='x')
        _Btn(btn_r, '📂  저장 위치 변경', self._pick_data_dir,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(16, 8)).pack(side='left', padx=(0, 8))
        _Btn(btn_r, '📁  현재 폴더 열기',
             lambda: subprocess.Popen(['explorer', str(DATA_DIR)]),
             font=(KO, 8), pad=(12, 7)).pack(side='left')

        self._gen_msg = tk.Label(cc, text='', font=(KO, 8), bg=CARD, fg=OK)
        self._gen_msg.pack(anchor='w', pady=(8, 0))

        # ── 안내 카드 ──
        g_out = tk.Frame(body, bg=BDR); g_out.pack(fill='x')
        g_in  = tk.Frame(g_out, bg=CARD); g_in.pack(fill='both', padx=1, pady=1)
        gc = tk.Frame(g_in, bg=CARD); gc.pack(fill='x', padx=20, pady=16)
        tk.Label(gc, text='💡  활용 팁', font=(KO, 9, 'bold'),
                 bg=CARD, fg=TEXT).pack(anchor='w', pady=(0, 8))
        for tip in [
            '구글 드라이브, OneDrive, Dropbox 폴더를 지정하면 자동 클라우드 백업됩니다.',
            '학교 네트워크 드라이브를 지정하면 어느 컴퓨터에서든 같은 데이터를 사용할 수 있습니다.',
            '저장 위치 변경 후 재시작하면 완전히 적용됩니다.',
        ]:
            tk.Label(gc, text=f'• {tip}', font=(KO, 8), bg=CARD, fg=SUB,
                     anchor='w', wraplength=540, justify='left'
                     ).pack(anchor='w', pady=2)

    def _refresh_data_size(self):
        """데이터 폴더 크기 계산."""
        try:
            total = sum(f.stat().st_size for f in DATA_DIR.rglob('*') if f.is_file())
            txt = f'{total/1024:.0f} KB' if total < 1_048_576 else f'{total/1_048_576:.1f} MB'
            self._data_size_lbl.configure(text=f'데이터 크기: {txt}')
        except Exception:
            self._data_size_lbl.configure(text='')

    def _pick_data_dir(self):
        """새 데이터 저장 위치 선택 및 적용."""
        new_path = filedialog.askdirectory(
            title='데이터 저장 위치 선택',
            initialdir=str(DATA_DIR))
        if not new_path:
            return

        new_dir = Path(new_path)
        if new_dir == DATA_DIR:
            self._gen_msg.configure(text='⚠️  현재 위치와 동일합니다', fg=WARN)
            return

        self._gen_msg.configure(text='📦  이동 중...', fg=WARN)
        self.update_idletasks()

        def _do():
            try:
                self._apply_new_data_dir(new_dir, self._move_data_var.get())
                self.after(0, lambda: (
                    self._data_dir_lbl.configure(text=str(new_dir)),
                    self._refresh_data_size(),
                    self._gen_msg.configure(
                        text='✅  저장 위치 변경 완료. 재시작하면 완전히 적용됩니다.', fg=OK),
                    self._lbl_status.configure(text='✅  데이터 위치 변경 완료', fg=OK),
                ))
            except Exception as e:
                msg = str(e)
                self.after(0, lambda m=msg: self._gen_msg.configure(
                    text=f'❌  {m}', fg=ERR))

        threading.Thread(target=_do, daemon=True).start()

    def _apply_new_data_dir(self, new_dir: Path, copy_data: bool):
        """데이터 저장 위치 실제 변경."""
        global DATA_DIR, STUDENTS_FILE, SCORES_FILE

        new_dir.mkdir(parents=True, exist_ok=True)

        if copy_data:
            for item in DATA_DIR.iterdir():
                dst = new_dir / item.name
                if item.is_dir():
                    shutil.copytree(str(item), str(dst), dirs_exist_ok=True)
                elif not dst.exists():
                    shutil.copy2(str(item), str(dst))

        # 경로 업데이트
        DATA_DIR      = new_dir
        STUDENTS_FILE = new_dir / 'students.json'
        SCORES_FILE   = new_dir / 'scores.json'

        # config.json에 저장 (전역)
        cfg_file = BASE_DIR / 'config.json'
        try:
            cfg = json.loads(cfg_file.read_text(encoding='utf-8')) \
                  if cfg_file.exists() else {}
        except Exception:
            cfg = {}
        cfg['data_dir'] = str(new_dir)
        cfg_file.write_text(
            json.dumps(cfg, ensure_ascii=False, indent=2), encoding='utf-8')


    # ─── 공통 헬퍼 ─────────────────────────────────────────────

    def _stu_page_header(self, parent, icon, title, desc):
        hdr = tk.Frame(parent, bg=BG2)
        hdr.pack(fill='x')
        tk.Frame(hdr, bg=GOLD, height=2).pack(side='bottom', fill='x')
        ph = tk.Frame(hdr, bg=BG2)
        ph.pack(fill='x', padx=24, pady=12)
        tk.Label(ph, text=icon, font=(KO, 18), bg=BG2, fg=GOLD).pack(side='left', padx=(0, 10))
        tl = tk.Frame(ph, bg=BG2)
        tl.pack(side='left')
        tk.Label(tl, text=title, font=(KO, 12, 'bold'), bg=BG2, fg=TEXT).pack(anchor='w')
        tk.Label(tl, text=desc, font=(KO, 8), bg=BG2, fg=SUB).pack(anchor='w', pady=(1, 0))

    def _stu_section(self, parent, title):
        """섹션 구분선 + 소제목."""
        f = tk.Frame(parent, bg=BG)
        f.pack(fill='x', pady=(12, 4))
        tk.Label(f, text=title, font=(KO, 8, 'bold'), bg=BG, fg=SUB).pack(side='left')
        tk.Frame(f, bg=BDR, height=1).pack(side='left', fill='x', expand=True, padx=(8, 0), pady=4)

    # ── 공통 폼 행 ─ 두 성적 화면이 동일 레이아웃 사용 ──────────
    _FORM_LBL_W = 9   # 레이블 고정 너비 (문자 단위)

    def _form_combo(self, parent, label, values, var):
        """레이블 + 콤보박스 행 (통일 간격)."""
        row = tk.Frame(parent, bg=BG)
        row.pack(fill='x', pady=4)
        tk.Label(row, text=label, font=(KO, 8), bg=BG, fg=SUB,
                 width=self._FORM_LBL_W, anchor='w').pack(side='left', padx=(0, 8))
        c = ttk.Combobox(row, textvariable=var, values=values,
                         state='readonly', font=(KO, 9))
        c.pack(side='left', fill='x', expand=True)
        return c

    def _form_entry(self, parent, label, default='', placeholder=''):
        """레이블 + 입력창 행. placeholder: 클릭 시 사라지는 힌트 텍스트."""
        row = tk.Frame(parent, bg=BG)
        row.pack(fill='x', pady=4)
        tk.Label(row, text=label, font=(KO, 8), bg=BG, fg=SUB,
                 width=self._FORM_LBL_W, anchor='w').pack(side='left', padx=(0, 8))
        outer = tk.Frame(row, bg=BDR, bd=1, relief='flat')
        outer.pack(side='left', fill='x', expand=True)
        e = tk.Entry(outer, font=(KO, 9), bg=CARD, fg=TEXT,
                     relief='flat', bd=0, insertbackground=TEXT)
        e.pack(fill='x', padx=4, pady=2)
        if placeholder:
            e.insert(0, placeholder); e.configure(fg=DIM)
            e.bind('<FocusIn>',  lambda _, en=e, ph=placeholder:
                   (en.delete(0, 'end'), en.configure(fg=TEXT)) if en.get() == ph else None)
            e.bind('<FocusOut>', lambda _, en=e, ph=placeholder:
                   (en.insert(0, ph), en.configure(fg=DIM)) if not en.get() else None)
        elif default:
            e.insert(0, default)
        return e

    def _stu_treeview(self, parent, columns, col_widths, height=8):
        """공통 Treeview (성적 내역 표)."""
        style = ttk.Style()
        style.configure('Stu.Treeview', font=(KO, 8), rowheight=22,
                        background=CARD, fieldbackground=CARD,
                        foreground=TEXT, borderwidth=0)
        style.configure('Stu.Treeview.Heading', font=(KO, 8, 'bold'),
                        background=ITEM, foreground=TEXT)
        style.map('Stu.Treeview', background=[('selected', SEL)])

        wrap = tk.Frame(parent, bg=BDR)
        wrap.pack(fill='both', expand=True)
        tv = ttk.Treeview(wrap, columns=columns, show='headings',
                          style='Stu.Treeview', height=height)
        sb = ttk.Scrollbar(wrap, orient='vertical', command=tv.yview,
                           style='Vertical.TScrollbar')
        tv.configure(yscrollcommand=sb.set)
        for col, w in zip(columns, col_widths):
            tv.heading(col, text=col)
            tv.column(col, width=w, anchor='center', minwidth=w)
        sb.pack(side='right', fill='y')
        tv.pack(fill='both', expand=True, padx=1, pady=1)
        return tv

    def _student_options(self):
        """콤보박스용 학생 목록 (이름+반)."""
        return [f'{s["name"]}  ({s.get("grade_class","")})' for s in load_students()]

    def _student_id_from_sel(self, sel_text):
        name = sel_text.split('  (')[0].strip()
        for s in load_students():
            if s['name'] == name:
                return s['id']
        return None

    # ─── 학생 관리 ──────────────────────────────────────────────

    def _build_stu_manage(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['stu_manage'] = fr
        self._stu_page_header(fr, '👤', '학생 관리', '학생을 등록하고 기본 정보를 관리합니다')

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=24, pady=14)

        # 상단 버튼
        top = tk.Frame(body, bg=BG)
        top.pack(fill='x', pady=(0, 10))
        self._stu_count_lbl = tk.Label(top, text='', font=(KO, 9), bg=BG, fg=SUB)
        self._stu_count_lbl.pack(side='left')
        _Btn(top, '➕  학생 추가', self._toggle_add_form,
             fg_n='#ffffff', fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(14, 6)).pack(side='right')

        # 학생 추가 폼 (처음엔 숨김)
        self._add_form_frame = tk.Frame(body, bg=CARD)
        self._add_form_visible = False
        self._build_add_form(self._add_form_frame)

        # 학생 목록 (스크롤)
        list_wrap = tk.Frame(body, bg=BDR, bd=1, relief='flat')
        list_wrap.pack(fill='both', expand=True)
        canvas = tk.Canvas(list_wrap, bg=BG, highlightthickness=0)
        sb2 = ttk.Scrollbar(list_wrap, orient='vertical', command=canvas.yview,
                            style='Vertical.TScrollbar')
        canvas.configure(yscrollcommand=sb2.set)
        sb2.pack(side='right', fill='y')
        canvas.pack(fill='both', expand=True)
        self._stu_list_canvas = canvas
        self._stu_list_inner = tk.Frame(canvas, bg=BG)
        cwin = canvas.create_window((0, 0), window=self._stu_list_inner, anchor='nw')
        def _on_stu_canvas_cfg(e):
            canvas.itemconfig(cwin, width=e.width)
            self.after(1, self._reset_stu_scroll)
        canvas.bind('<Configure>', _on_stu_canvas_cfg)
        canvas.bind('<MouseWheel>', lambda e: canvas.yview_scroll(-1*(e.delta//120), 'units'))

        self._refresh_stu_list()

    def _build_add_form(self, parent):
        f = tk.Frame(parent, bg=CARD)
        f.pack(fill='x', padx=1, pady=1)
        inner = tk.Frame(f, bg=CARD)
        inner.pack(fill='x', padx=16, pady=12)

        tk.Label(inner, text='학생 추가', font=(KO, 9, 'bold'), bg=CARD, fg=TEXT).grid(
            row=0, columnspan=4, sticky='w', pady=(0, 8))

        fields = [('이름', '이름 입력'), ('학교', '학교명 입력'), ('학년/반', '예: 3학년 2반')]
        self._add_entries = {}
        for col, (lbl, ph) in enumerate(fields):
            tk.Label(inner, text=lbl, font=(KO, 8), bg=CARD, fg=SUB).grid(
                row=1, column=col*2, sticky='w', padx=(0 if col==0 else 12, 4))
            e = tk.Entry(inner, font=(KO, 9), bg=ITEM, fg=TEXT,
                         relief='flat', bd=1, insertbackground=TEXT, width=14)
            e.insert(0, ph)
            e.configure(fg=DIM)
            e.bind('<FocusIn>',  lambda _, p=ph, en=e: (en.delete(0,'end'), en.configure(fg=TEXT)) if en.get()==p else None)
            e.bind('<FocusOut>', lambda _, p=ph, en=e: (en.insert(0,p), en.configure(fg=DIM)) if not en.get() else None)
            e.grid(row=1, column=col*2+1, sticky='ew', padx=(0, 4))
            inner.columnconfigure(col*2+1, weight=1)
            self._add_entries[lbl] = (e, ph)

        btn_row = tk.Frame(inner, bg=CARD)
        btn_row.grid(row=2, columnspan=6, sticky='e', pady=(10, 0))
        _Btn(btn_row, '취소', self._toggle_add_form, font=(KO, 8), pad=(10, 5)).pack(side='left', padx=(0,6))
        _Btn(btn_row, '저장', self._save_new_student,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 8, 'bold'), pad=(14, 5)).pack(side='left')

    def _toggle_add_form(self):
        if self._add_form_visible:
            self._add_form_frame.pack_forget()
            self._add_form_visible = False
        else:
            self._add_form_frame.pack(fill='x', pady=(0, 10))
            self._add_form_visible = True

    def _save_new_student(self):
        vals = {}
        for lbl, (e, ph) in self._add_entries.items():
            v = e.get().strip()
            vals[lbl] = '' if v == ph else v

        if not vals['이름']:
            self._lbl_status.configure(text='⚠️  이름을 입력해주세요', fg=WARN); return

        students = load_students()
        students.append({
            'id':          _new_id(),
            'name':        vals['이름'],
            'school':      vals['학교'],
            'grade_class': vals['학년/반'],
            'created_at':  datetime.datetime.now().isoformat(timespec='seconds'),
        })
        save_students(students)
        self._toggle_add_form()
        # 폼 초기화
        for lbl, (e, ph) in self._add_entries.items():
            e.delete(0, 'end'); e.insert(0, ph); e.configure(fg=DIM)
        self._refresh_stu_list()
        self._refresh_score_combos()
        self._lbl_status.configure(text=f'✅  {vals["이름"]} 학생 추가 완료', fg=OK)

    def _delete_student(self, sid):
        students = [s for s in load_students() if s['id'] != sid]
        save_students(students)
        scores = [r for r in load_scores() if r['student_id'] != sid]
        save_scores(scores)
        self._refresh_stu_list()
        self._refresh_score_combos()
        self._lbl_status.configure(text='학생 삭제됨', fg=SUB)

    def _refresh_stu_list(self):
        for w in self._stu_list_inner.winfo_children():
            w.destroy()
        students = load_students()
        self._stu_count_lbl.configure(text=f'총 {len(students)}명 등록')
        if not students:
            tk.Label(self._stu_list_inner,
                     text='\n등록된 학생이 없습니다\n오른쪽 상단 [➕ 학생 추가] 버튼을 눌러 추가하세요',
                     font=(KO, 9), bg=BG, fg=DIM, justify='center').pack(pady=30)
        else:
            for s in students:
                self._make_stu_card(self._stu_list_inner, s)

        self.after(1, self._reset_stu_scroll)

    def _make_stu_card(self, parent, s):
        outer = tk.Frame(parent, bg=BDR, bd=1, relief='flat')
        outer.pack(fill='x', padx=2, pady=3)
        inner = tk.Frame(outer, bg=CARD)
        inner.pack(fill='x', padx=1, pady=1)
        row = tk.Frame(inner, bg=CARD)
        row.pack(fill='x', padx=14, pady=10)

        tk.Label(row, text='👤', font=(KO, 16), bg=CARD, fg=TEXT).pack(side='left', padx=(0,10))
        info = tk.Frame(row, bg=CARD)
        info.pack(side='left', fill='both', expand=True)
        tk.Label(info, text=s['name'], font=(KO, 10, 'bold'), bg=CARD, fg=TEXT, anchor='w').pack(anchor='w')
        sub_txt = f'{s.get("school","─")}   |   {s.get("grade_class","─")}   |   등록: {s["created_at"][:10]}'
        tk.Label(info, text=sub_txt, font=(KO, 8), bg=CARD, fg=SUB, anchor='w').pack(anchor='w', pady=(2,0))

        _Btn(row, '✕', lambda sid=s['id']: self._delete_student(sid),
             fg_n=DIM, fg_h=ERR, pad=(6, 4)).pack(side='right')

    # ─── 모의고사 · 내신 성적 ───────────────────────────────────

    def _build_stu_mock(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['stu_mock'] = fr
        self._stu_page_header(fr, '📊', '모의고사 · 내신 성적',
                              '수능 모의고사 및 내신 영어 성적을 날짜별로 누적 저장합니다')

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=24, pady=14)

        left = tk.Frame(body, bg=BG, width=300)
        left.pack(side='left', fill='y', padx=(0, 16))
        left.pack_propagate(False)
        tk.Frame(body, bg=BDR, width=1).pack(side='left', fill='y', padx=(0, 16))
        right = tk.Frame(body, bg=BG)
        right.pack(side='left', fill='both', expand=True)

        # ── 입력 폼 ──
        self._stu_section(left, '성적 입력')

        self._mock_stu_var = tk.StringVar()
        self._mock_stu_combo = self._form_combo(left, '학생', [], self._mock_stu_var)
        self._mock_stu_combo.bind('<<ComboboxSelected>>', lambda _: self._refresh_mock_history())

        self._mock_type_var = tk.StringVar(value='모의고사')
        self._form_combo(left, '구분',
                         ['모의고사', '내신 (중간)', '내신 (기말)'],
                         self._mock_type_var)

        self._mock_name_e  = self._form_entry(left, '시험명', placeholder='예: 2026년 3월')
        self._mock_date_e  = self._form_entry(left, '날짜', default=datetime.date.today().isoformat())
        self._mock_score_e = self._form_entry(left, '원점수')
        self._mock_grade_e = self._form_entry(left, '등급')
        self._mock_pct_e   = self._form_entry(left, '백분위')
        self._mock_note_e  = self._form_entry(left, '메모')

        _Btn(left, '  💾  저장  ', self._save_mock_score,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(20, 8)).pack(pady=(12, 0))

        # ── 내역 테이블 ──
        self._stu_section(right, '성적 내역')
        cols = ['날짜', '구분', '시험명', '원점수', '등급', '백분위', '메모']
        widths = [80, 80, 120, 60, 50, 60, 100]
        self._mock_editing_id = None
        self._mock_tv = self._stu_treeview(right, cols, widths, height=12)
        self._mock_tv.bind('<Delete>',           lambda _: self._mock_delete_selected())
        self._mock_tv.bind('<<TreeviewSelect>>', lambda _: self._mock_load_selected())

        act_row = tk.Frame(right, bg=BG); act_row.pack(fill='x', pady=(4, 0))
        self._btn_mock_del = _Btn(act_row, '🗑  선택 행 삭제', self._mock_delete_selected,
                                  fg_n=ERR, fg_h='#ffffff', bg_n=CARD, bg_h=ERR,
                                  font=(KO, 8), pad=(10, 5))
        self._btn_mock_del.pack(side='right')
        self._btn_mock_del.set_enabled(False)
        tk.Label(act_row, text='행을 클릭하면 왼쪽 폼에서 수정할 수 있습니다',
                 font=(KO, 7), bg=BG, fg=DIM).pack(side='left')

    # ── 모의고사 행 선택/삭제 ────────────────────────────────────

    def _mock_load_selected(self):
        """모의고사 행 클릭 → 폼에 데이터 로드."""
        sel = self._mock_tv.selection()
        self._btn_mock_del.set_enabled(bool(sel))
        if not sel: return
        iid = sel[0]
        self._mock_editing_id = iid
        r = next((x for x in load_scores() if x['id'] == iid), None)
        if not r: return
        # 학생 설정
        s = next((x for x in load_students() if x['id'] == r.get('student_id')), None)
        if s:
            opt = f'{s["name"]}  ({s.get("grade_class","")})'
            self._mock_stu_var.set(opt)
        self._mock_type_var.set(r.get('exam_type', ''))
        for e, key in [(self._mock_name_e,'exam_name'), (self._mock_date_e,'date'),
                       (self._mock_score_e,'score'), (self._mock_grade_e,'grade'),
                       (self._mock_pct_e,'percentile'), (self._mock_note_e,'notes')]:
            e.delete(0, 'end'); e.insert(0, r.get(key, '')); e.configure(fg=TEXT)

    def _mock_delete_selected(self):
        sel = self._mock_tv.selection()
        if not sel: return
        scores = [r for r in load_scores() if r['id'] != sel[0]]
        save_scores(scores)
        self._mock_tv.delete(sel[0])
        self._mock_editing_id = None
        self._btn_mock_del.set_enabled(False)
        self._lbl_status.configure(text='삭제됨', fg=SUB)

    def _save_mock_score(self):
        stu_sel = self._mock_stu_var.get()
        if not stu_sel:
            self._lbl_status.configure(text='⚠️  학생을 선택하세요', fg=WARN); return
        sid = self._student_id_from_sel(stu_sel)
        score_val = self._mock_score_e.get().strip()
        if not score_val:
            self._lbl_status.configure(text='⚠️  원점수를 입력하세요', fg=WARN); return

        now    = datetime.datetime.now()
        record = {
            'id':          self._mock_editing_id or _new_id(),
            'student_id':  sid,
            'category':    'mock',
            'exam_type':   self._mock_type_var.get(),
            'exam_name':   self._mock_name_e.get().strip(),
            'date':        self._mock_date_e.get().strip() or now.strftime('%Y-%m-%d'),
            'time':        now.strftime('%H:%M:%S'),
            'score':       score_val,
            'grade':       self._mock_grade_e.get().strip(),
            'percentile':  self._mock_pct_e.get().strip(),
            'notes':       self._mock_note_e.get().strip(),
            'created_at':  now.isoformat(timespec='seconds'),
        }
        scores = load_scores()
        if self._mock_editing_id:
            scores = [r if r['id'] != self._mock_editing_id else record for r in scores]
            label = '✅  성적 수정 완료'
        else:
            scores.append(record)
            label = '✅  성적 저장 완료'
        save_scores(scores)
        self._mock_editing_id = None
        self._btn_mock_del.set_enabled(False)
        self._refresh_mock_history()
        self._lbl_status.configure(text=label, fg=OK)
        self._auto_sync_sheets()

    def _refresh_mock_history(self):
        self._mock_tv.delete(*self._mock_tv.get_children())
        stu_sel = self._mock_stu_var.get()
        if not stu_sel: return
        sid = self._student_id_from_sel(stu_sel)
        rows = sorted(
            [r for r in load_scores() if r.get('student_id') == sid and r.get('category') == 'mock'],
            key=lambda r: (r.get('date',''), r.get('time','')), reverse=True)
        for r in rows:
            self._mock_tv.insert('', 'end', iid=r['id'],
                values=(r['date'], r.get('exam_type',''), r.get('exam_name',''),
                        r.get('score',''), r.get('grade',''), r.get('percentile',''),
                        r.get('notes','')))

    # ─── 워크시트 성적 ──────────────────────────────────────────

    def _update_sheet_google_lock(self):
        """워크시트 성적 화면의 Google 로그인 잠금 상태 업데이트."""
        if not hasattr(self, '_sheet_google_overlay'): return
        if TOKEN_FILE.exists():
            self._sheet_google_overlay.place_forget()
        else:
            self._sheet_google_overlay.place(relx=0, rely=0, relwidth=1, relheight=1)
            self._sheet_google_overlay.lift()

    def _build_stu_sheet(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['stu_sheet'] = fr

        # ── Google 로그인 잠금 오버레이 (미로그인 시 표시) ──────
        overlay = tk.Frame(fr, bg=BG)
        self._sheet_google_overlay = overlay
        center = tk.Frame(overlay, bg=BG)
        center.place(relx=0.5, rely=0.42, anchor='center')
        tk.Label(center, text='🔒', font=(KO, 44), bg=BG, fg=GOLD).pack()
        tk.Label(center, text='워크시트 성적은 Google 로그인이 필요합니다',
                 font=(KO, 13, 'bold'), bg=BG, fg=TEXT).pack(pady=(10, 5))
        tk.Label(center, text='설정 → 연동 설정에서 Google 계정에 로그인하세요.',
                 font=(KO, 9), bg=BG, fg=SUB).pack()
        _Btn(center, '⚙️  연동 설정으로 이동',
             lambda: self._show_module('set_connect'),
             fg_n=GOLD, fg_h=GOLDB, bg_n=CARD, bg_h=HOV,
             font=(KO, 10, 'bold'), pad=(16, 10)).pack(pady=(16, 0))

        self._stu_page_header(fr, '📋', '워크시트 성적',
                              'Easy / Standard A·B / Hard 성적을 지문별로 누적 저장합니다')

        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=24, pady=14)

        left = tk.Frame(body, bg=BG, width=300)
        left.pack(side='left', fill='y', padx=(0, 16))
        left.pack_propagate(False)
        tk.Frame(body, bg=BDR, width=1).pack(side='left', fill='y', padx=(0, 16))
        right = tk.Frame(body, bg=BG)
        right.pack(side='left', fill='both', expand=True)

        # ── 입력 폼 ──
        self._stu_section(left, '성적 입력')

        self._sheet_stu_var = tk.StringVar()
        self._sheet_stu_combo = self._form_combo(left, '학생', [], self._sheet_stu_var)
        self._sheet_stu_combo.bind('<<ComboboxSelected>>', lambda _: self._refresh_sheet_history())

        self._sheet_level_var = tk.StringVar(value='Easy')
        self._sheet_type_var  = tk.StringVar(value='')
        self._form_combo(left, '레벨',
                         ['Easy', 'Standard A', 'Standard B', 'Hard'],
                         self._sheet_level_var)
        # Google 로그인 상태일 때만 문제 유형 불러오기
        cfg = load_config()
        default_qtypes = cfg.get('custom_problem_types', []) if TOKEN_FILE.exists() else []
        self._sheet_type_combo = self._form_combo(
            left, '문제 유형', default_qtypes, self._sheet_type_var)
        self._sheet_type_var.set(default_qtypes[0] if default_qtypes else '')

        self._sheet_date_e  = self._form_entry(left, '날짜', default=datetime.date.today().isoformat())
        self._sheet_total_e = self._form_entry(left, '총점')
        self._sheet_note_e  = self._form_entry(left, '메모')

        _Btn(left, '  💾  저장  ', self._save_sheet_score,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(20, 8)).pack(pady=(12, 0))

        # ── 내역 테이블 ──
        self._stu_section(right, '성적 내역')
        cols = ['날짜', '레벨', '문제 유형', '총점', '메모']
        widths = [80, 90, 120, 60, 150]
        self._sheet_editing_id = None
        self._sheet_tv = self._stu_treeview(right, cols, widths, height=12)
        self._sheet_tv.bind('<Delete>',           lambda _: self._sheet_delete_selected())
        self._sheet_tv.bind('<<TreeviewSelect>>', lambda _: self._sheet_load_selected())

        act_row2 = tk.Frame(right, bg=BG); act_row2.pack(fill='x', pady=(4, 0))
        self._btn_sheet_del = _Btn(act_row2, '🗑  선택 행 삭제', self._sheet_delete_selected,
                                   fg_n=ERR, fg_h='#ffffff', bg_n=CARD, bg_h=ERR,
                                   font=(KO, 8), pad=(10, 5))
        self._btn_sheet_del.pack(side='right')
        self._btn_sheet_del.set_enabled(False)
        tk.Label(act_row2, text='행을 클릭하면 왼쪽 폼에서 수정할 수 있습니다',
                 font=(KO, 7), bg=BG, fg=DIM).pack(side='left')

    # ── 워크시트 행 선택/삭제 ────────────────────────────────────

    def _sheet_load_selected(self):
        sel = self._sheet_tv.selection()
        self._btn_sheet_del.set_enabled(bool(sel))
        if not sel: return
        iid = sel[0]
        self._sheet_editing_id = iid
        r = next((x for x in load_scores() if x['id'] == iid), None)
        if not r: return
        s = next((x for x in load_students() if x['id'] == r.get('student_id')), None)
        if s:
            self._sheet_stu_var.set(f'{s["name"]}  ({s.get("grade_class","")})')
        self._sheet_level_var.set(r.get('level', 'Easy'))
        self._sheet_type_var.set(r.get('qtype', ''))
        for e, key in [(self._sheet_date_e, 'date'),
                       (self._sheet_total_e, 'total_score'),
                       (self._sheet_note_e, 'notes')]:
            e.delete(0, 'end'); e.insert(0, r.get(key, '')); e.configure(fg=TEXT)

    def _sheet_delete_selected(self):
        sel = self._sheet_tv.selection()
        if not sel: return
        scores = [r for r in load_scores() if r['id'] != sel[0]]
        save_scores(scores)
        self._sheet_tv.delete(sel[0])
        self._sheet_editing_id = None
        self._btn_sheet_del.set_enabled(False)
        self._lbl_status.configure(text='삭제됨', fg=SUB)

    def _save_sheet_score(self):
        stu_sel = self._sheet_stu_var.get()
        if not stu_sel:
            self._lbl_status.configure(text='⚠️  학생을 선택하세요', fg=WARN); return

        sid    = self._student_id_from_sel(stu_sel)
        now    = datetime.datetime.now()
        record = {
            'id':          self._sheet_editing_id or _new_id(),
            'student_id':  sid,
            'category':    'sheet',
            'level':       self._sheet_level_var.get(),
            'qtype':       self._sheet_type_var.get(),
            'date':        self._sheet_date_e.get().strip() or now.strftime('%Y-%m-%d'),
            'time':        now.strftime('%H:%M:%S'),
            'total_score': self._sheet_total_e.get().strip(),
            'notes':       self._sheet_note_e.get().strip(),
            'created_at':  now.isoformat(timespec='seconds'),
        }
        scores = load_scores()
        if self._sheet_editing_id:
            scores = [r if r['id'] != self._sheet_editing_id else record for r in scores]
            label = '✅  성적 수정 완료'
        else:
            scores.append(record)
            label = '✅  성적 저장 완료'
        save_scores(scores)
        self._sheet_editing_id = None
        self._btn_sheet_del.set_enabled(False)
        self._refresh_sheet_history()
        self._lbl_status.configure(text=label, fg=OK)
        self._auto_sync_sheets()

    def _refresh_sheet_history(self):
        self._sheet_tv.delete(*self._sheet_tv.get_children())
        stu_sel = self._sheet_stu_var.get()
        if not stu_sel: return
        sid = self._student_id_from_sel(stu_sel)
        rows = sorted(
            [r for r in load_scores() if r.get('student_id') == sid and r.get('category') == 'sheet'],
            key=lambda r: (r.get('date',''), r.get('time','')), reverse=True)
        for r in rows:
            self._sheet_tv.insert('', 'end', iid=r['id'],
                values=(r['date'], r.get('level',''), r.get('qtype',''),
                        r.get('total_score',''), r.get('notes','')))

    def _delete_score_row(self, tv, _category):
        sel = tv.selection()
        if not sel: return
        rid = sel[0]
        scores = [r for r in load_scores() if r['id'] != rid]
        save_scores(scores)
        tv.delete(rid)
        self._lbl_status.configure(text='삭제됨', fg=SUB)

    # ─── 설정 → 연동 설정 ──────────────────────────────────────

    def _build_set_connect_page(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['set_connect'] = fr
        self._stu_page_header(fr, '🔗', '연동 설정',
                              'Google 계정으로 로그인하여 Drive 파일을 연결합니다')

        # 스크롤 가능한 body
        canvas = tk.Canvas(fr, bg=BG, highlightthickness=0)
        sb = ttk.Scrollbar(fr, orient='vertical', command=canvas.yview, style='Vertical.TScrollbar')
        canvas.configure(yscrollcommand=sb.set)
        sb.pack(side='right', fill='y')
        canvas.pack(fill='both', expand=True)
        body = tk.Frame(canvas, bg=BG)
        cwin = canvas.create_window((0, 0), window=body, anchor='nw')
        body.bind('<Configure>', lambda e: canvas.configure(scrollregion=(0, 0, e.width, e.height)))
        canvas.bind('<Configure>', lambda e: canvas.itemconfig(cwin, width=e.width))
        canvas.bind('<MouseWheel>', lambda e: canvas.yview_scroll(-1*(e.delta//120), 'units'))

        pad = {'padx': 32, 'pady': (0, 14), 'fill': 'x'}

        def card(title):
            wrap = tk.Frame(body, bg=BG)
            wrap.pack(**pad)
            o = tk.Frame(wrap, bg=BDR); o.pack(fill='x')
            i = tk.Frame(o, bg=CARD); i.pack(fill='both', padx=1, pady=1)
            c = tk.Frame(i, bg=CARD); c.pack(fill='x', padx=20, pady=16)
            tk.Label(c, text=title, font=(KO, 9, 'bold'), bg=CARD, fg=TEXT
                     ).pack(anchor='w', pady=(0, 10))
            return c

        # ── 1. OAuth 클라이언트 키 선택 ──────────────────────────
        cc = card('① OAuth 클라이언트 키')
        cfg = load_config()
        self._oauth_key_var = tk.StringVar(value=cfg.get('oauth_key', ''))

        kr = tk.Frame(cc, bg=CARD); kr.pack(fill='x', pady=4)
        tk.Label(kr, text='client_secrets.json', font=(KO, 8), bg=CARD, fg=SUB,
                 width=18, anchor='w').pack(side='left', padx=(0, 8))
        tk.Entry(kr, textvariable=self._oauth_key_var, font=(KO, 8),
                 bg=ITEM, fg=TEXT, relief='flat', bd=1, state='readonly'
                 ).pack(side='left', fill='x', expand=True, padx=(0, 6))
        _Btn(kr, '📂 파일 선택', self._pick_oauth_key,
             font=(KO, 8), pad=(8, 4)).pack(side='left')

        tk.Label(cc, text='Google Cloud Console → API 및 서비스 → 사용자 인증 정보 →\n'
                 'OAuth 2.0 클라이언트 ID (데스크톱 앱) 생성 후 JSON 다운로드',
                 font=(KO, 7), bg=CARD, fg=DIM, justify='left').pack(anchor='w', pady=(6, 0))

        # ── 2. Google 로그인 상태 ────────────────────────────────
        lc = card('② Google 계정 로그인')
        self._google_email_lbl = tk.Label(lc, text='', font=(KO, 9), bg=CARD, fg=SUB)
        self._google_email_lbl.pack(anchor='w', pady=(0, 8))
        self._update_google_status()

        btn_row = tk.Frame(lc, bg=CARD); btn_row.pack(anchor='w')
        _Btn(btn_row, '🔓  Google 로그인', self._do_google_login,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(16, 8)).pack(side='left', padx=(0, 8))
        _Btn(btn_row, '로그아웃', self._do_google_logout,
             font=(KO, 8), pad=(10, 6)).pack(side='left')

        # ── 3. Drive 파일 선택 ───────────────────────────────────
        dc = card('③ Google Drive 파일 선택')
        tk.Label(dc, text='로그인 후 Drive에서 연결할 Google Sheets 파일을 선택하세요.',
                 font=(KO, 8), bg=CARD, fg=SUB).pack(anchor='w', pady=(0, 8))
        _Btn(dc, '🔄  파일 목록 새로고침', self._refresh_drive_list,
             font=(KO, 8), pad=(12, 6)).pack(anchor='w', pady=(0, 8))

        # 파일 목록
        list_fr = tk.Frame(dc, bg=ITEM, bd=1, relief='flat')
        list_fr.pack(fill='x', pady=(0, 6))
        self._drive_list_inner = tk.Frame(list_fr, bg=ITEM)
        self._drive_list_inner.pack(fill='x', padx=1, pady=1)
        tk.Label(self._drive_list_inner, text='(로그인 후 새로고침)',
                 font=(KO, 8), bg=ITEM, fg=DIM).pack(pady=12)

        # 선택 상태
        sel_r = tk.Frame(dc, bg=CARD); sel_r.pack(fill='x', pady=(0, 4))
        tk.Label(sel_r, text='선택된 파일', font=(KO, 8), bg=CARD, fg=SUB,
                 width=10, anchor='w').pack(side='left')
        self._drive_sel_lbl = tk.Label(sel_r, text='없음', font=(KO, 8, 'bold'),
                                       bg=CARD, fg=TEXT)
        self._drive_sel_lbl.pack(side='left')

        # ── 4. 문제 유형 불러오기 ────────────────────────────────
        tc = card('④ 문제 유형 읽어오기')
        tk.Label(tc,
                 text='선택한 시트의 B열(유형 컬럼)에 문제 유형 이름을 나열하면\n'
                      '워크시트 성적 입력 화면의 드롭다운에 자동 반영됩니다.\n'
                      '(1행=헤더 자동 제외 · 탭명: 아무 탭이나 첫 번째 탭 사용)',
                 font=(KO, 8), bg=CARD, fg=SUB, justify='left').pack(anchor='w', pady=(0, 10))

        btn_t = tk.Frame(tc, bg=CARD); btn_t.pack(anchor='w')
        _Btn(btn_t, '📥  문제 유형 불러오기', self._load_problem_types_from_sheet,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(16, 8)).pack(side='left', padx=(0, 10))
        self._types_lbl = tk.Label(btn_t, text='', font=(KO, 8), bg=CARD, fg=OK)
        self._types_lbl.pack(side='left')

        # 현재 적용된 문제 유형 표시
        self._cur_types_lbl = tk.Label(tc, text='', font=(KO, 8), bg=CARD, fg=SUB)
        self._cur_types_lbl.pack(anchor='w', pady=(6, 0))
        self._refresh_types_display()

        # ── 5. 성적 동기화 ───────────────────────────────────────
        sc2 = card('⑤ 성적 데이터 동기화')
        tk.Label(sc2, text='성적 저장 시 자동으로 Drive 파일에 반영됩니다.\n'
                           '아래 버튼으로 전체 데이터를 즉시 동기화할 수 있습니다.',
                 font=(KO, 8), bg=CARD, fg=SUB, justify='left').pack(anchor='w', pady=(0, 10))
        _Btn(sc2, '☁️  지금 전체 동기화', self._sync_all_to_sheets,
             fg_n='#ffffff', fg_h='#ffffff', bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(16, 8)).pack(anchor='w')

        # 여백
        tk.Frame(body, bg=BG, height=20).pack()

    # ─── OAuth / Drive 헬퍼 ─────────────────────────────────────

    def _pick_oauth_key(self):
        path = filedialog.askopenfilename(
            title='OAuth client_secrets.json 선택',
            filetypes=[('JSON 파일', '*.json'), ('모든 파일', '*.*')],
            initialdir=str(Path.home()))
        if path:
            self._oauth_key_var.set(path)
            cfg = load_config(); cfg['oauth_key'] = path; save_config(cfg)

    def _update_google_status(self):
        email = _google_email()
        if email:
            self._google_email_lbl.configure(
                text=f'✅  {email} (로그인됨)', fg=OK)
        else:
            self._google_email_lbl.configure(
                text='로그인되지 않았습니다.', fg=DIM)

    def _do_google_login(self):
        key = self._oauth_key_var.get().strip()
        if not key or not Path(key).exists():
            self._lbl_status.configure(
                text='⚠️  먼저 client_secrets.json 파일을 선택하세요', fg=WARN); return
        self._lbl_status.configure(text='🌐 브라우저에서 Google 로그인 진행 중...', fg=WARN)
        self.update_idletasks()

        def _auth():
            try:
                _google_creds(key)
                self.after(0, self._update_google_status)
                self.after(0, self._update_sheet_google_lock)
                self.after(0, lambda: self._lbl_status.configure(
                    text='✅  Google 로그인 성공', fg=OK))
            except Exception as e:
                msg = str(e)
                self.after(0, lambda m=msg: self._lbl_status.configure(
                    text=f'❌  로그인 실패: {m}', fg=ERR))

        threading.Thread(target=_auth, daemon=True).start()

    def _do_google_logout(self):
        if TOKEN_FILE.exists():
            TOKEN_FILE.unlink()
        # 연동 관련 설정 초기화 (계정별 settings.json에서 제거)
        _pop_config('custom_problem_types', 'sheets_id')
        # UI 초기화
        self._update_google_status()
        self._drive_sel_lbl.configure(text='없음')
        for w in self._drive_list_inner.winfo_children(): w.destroy()
        tk.Label(self._drive_list_inner, text='(로그인 후 새로고침)',
                 font=(KO, 8), bg=ITEM, fg=DIM).pack(pady=10)
        self._refresh_types_display()
        self._apply_custom_types()
        self._update_sheet_google_lock()
        self._lbl_status.configure(text='로그아웃 완료 — 연동 정보 초기화됨', fg=SUB)

    def _refresh_drive_list(self):
        key = load_config().get('oauth_key', '')
        if not key:
            self._lbl_status.configure(text='⚠️  먼저 로그인하세요', fg=WARN); return
        self._lbl_status.configure(text='📂 Drive 파일 목록 불러오는 중...', fg=WARN)
        # 목록 영역에 로딩 표시
        for w in self._drive_list_inner.winfo_children(): w.destroy()
        tk.Label(self._drive_list_inner, text='불러오는 중...', font=(KO, 8),
                 bg=ITEM, fg=DIM).pack(pady=10)

        def _fetch():
            try:
                creds = _google_creds(key)
                files = list_drive_sheets(creds)
                self.after(0, lambda: self._populate_drive_list(files))
            except Exception as e:
                msg = str(e)
                self.after(0, lambda m=msg: self._lbl_status.configure(
                    text=f'❌  {m}', fg=ERR))

        threading.Thread(target=_fetch, daemon=True).start()

    def _populate_drive_list(self, files):
        for w in self._drive_list_inner.winfo_children(): w.destroy()
        if not files:
            tk.Label(self._drive_list_inner, text='파일 없음', font=(KO, 8),
                     bg=ITEM, fg=DIM).pack(pady=8)
            return
        for name, fid, mtime in files:
            row = tk.Frame(self._drive_list_inner, bg=ITEM, cursor='hand2')
            row.pack(fill='x', padx=2, pady=1)
            tk.Label(row, text='📊', font=(KO, 10), bg=ITEM, fg=GOLD,
                     width=2).pack(side='left', padx=(6, 4), pady=5)
            tk.Label(row, text=name, font=(KO, 9), bg=ITEM, fg=TEXT,
                     anchor='w').pack(side='left', fill='x', expand=True)
            tk.Label(row, text=mtime, font=(KO, 7), bg=ITEM, fg=DIM,
                     padx=8).pack(side='right')
            for w in (row,) + tuple(row.winfo_children()):
                w.bind('<Enter>', lambda _, r=row: r.configure(bg=HOV))
                w.bind('<Leave>', lambda _, r=row: r.configure(bg=ITEM))
                w.bind('<Button-1>', lambda _, n=name, fid=fid: self._select_drive_file(n, fid))
        self._lbl_status.configure(text=f'✅  {len(files)}개 파일 로드됨', fg=OK)

    def _select_drive_file(self, name, fid):
        self._drive_sel_lbl.configure(text=f'📊  {name}', fg=GOLD)
        cfg = load_config(); cfg['sheets_id'] = fid; save_config(cfg)
        self._lbl_status.configure(text=f'✅  선택됨: {name}', fg=OK)
        # 파일 선택 후 목록 숨기기 (새로고침 버튼으로 다시 표시 가능)
        for w in self._drive_list_inner.winfo_children():
            w.destroy()
        tk.Label(self._drive_list_inner,
                 text=f'✓  {name}\n(변경하려면 "파일 목록 새로고침" 클릭)',
                 font=(KO, 8), bg=ITEM, fg=SUB, justify='center'
                 ).pack(pady=10)

    def _load_problem_types_from_sheet(self):
        cfg = load_config()
        key = cfg.get('oauth_key', ''); sid = cfg.get('sheets_id', '')
        if not key:
            self._lbl_status.configure(text='⚠️  먼저 Google 로그인하세요', fg=WARN); return
        if not sid:
            self._lbl_status.configure(text='⚠️  ③에서 Drive 파일을 먼저 선택하세요', fg=WARN); return
        self._types_lbl.configure(text='불러오는 중...', fg=WARN)
        self._lbl_status.configure(text='시트에서 문제 유형 읽는 중...', fg=WARN)

        def _fetch():
            try:
                creds  = _google_creds(key)
                qtypes = read_sheet_problem_types(creds, sid)
                if not qtypes:
                    self.after(0, lambda: (
                        self._types_lbl.configure(
                            text='⚠️  B열에 유형 데이터가 없습니다 (2행부터 입력)', fg=WARN),
                        self._lbl_status.configure(
                            text='⚠️  B열 2행부터 유형 이름을 입력 후 다시 시도', fg=WARN),
                    )); return
                c = load_config()
                c['custom_problem_types'] = qtypes; save_config(c)
                self.after(0, lambda: (
                    self._types_lbl.configure(
                        text=f'✅  {len(qtypes)}개 유형 적용됨', fg=OK),
                    self._lbl_status.configure(
                        text=f'✅  문제 유형 {len(qtypes)}개 적용: ' + ', '.join(qtypes[:3]) + ('...' if len(qtypes)>3 else ''),
                        fg=OK),
                    self._refresh_types_display(),
                    self._apply_custom_types(),
                ))
            except Exception as e:
                msg = str(e) or repr(e) or traceback.format_exc().strip().split('\n')[-1]
                self.after(0, lambda m=msg: (
                    self._types_lbl.configure(text=f'❌  오류', fg=ERR),
                    self._lbl_status.configure(text=f'❌  {m}', fg=ERR),
                ))

        threading.Thread(target=_fetch, daemon=True).start()

    def _refresh_types_display(self):
        if not hasattr(self, '_cur_types_lbl'): return
        # Google 로그인 상태일 때만 저장된 유형 표시
        if not TOKEN_FILE.exists():
            self._cur_types_lbl.configure(text='(Google 미로그인 — Drive 시트 미연결)')
            return
        cfg   = load_config()
        types = cfg.get('custom_problem_types', [])
        if types:
            self._cur_types_lbl.configure(
                text='적용 중: ' + ', '.join(types[:6]) + ('…' if len(types) > 6 else ''))
        else:
            self._cur_types_lbl.configure(text='(Drive 시트 미연결 — 기본 유형 사용 중)')

    def _apply_custom_types(self):
        """워크시트 성적 화면의 문제 유형 드롭다운 업데이트.
        Google 미로그인 시 빈 목록."""
        if TOKEN_FILE.exists():
            cfg   = load_config()
            types = cfg.get('custom_problem_types', [])
        else:
            types = []
        c = getattr(self, '_sheet_type_combo', None)
        if c:
            c['values'] = types
            self._sheet_type_var.set(types[0] if types else '')

    # ─── 공통 Google Sheets 동기화 ──────────────────────────────

    def _get_or_create_ws(self, sh, name, headers):
        import gspread
        try:
            return sh.worksheet(name)
        except gspread.WorksheetNotFound:
            ws = sh.add_worksheet(name, rows=1000, cols=len(headers))
            ws.append_row(headers); return ws

    def _sync_all_to_sheets(self):
        cfg = load_config()
        key = cfg.get('oauth_key', ''); sid = cfg.get('sheets_id', '')
        if not key or not sid:
            self._lbl_status.configure(
                text='⚠️  설정 → 연동 설정에서 로그인 및 파일을 선택하세요', fg=WARN); return
        self._lbl_status.configure(text='☁️  동기화 중...', fg=WARN)

        def _do():
            try:
                import gspread
                creds = _google_creds(key)
                gc    = gspread.authorize(creds)
                sh    = gc.open_by_key(sid)
                stu_map = {s['id']: s['name'] for s in load_students()}

                ws1 = self._get_or_create_ws(sh, '학생 목록',
                    ['ID', '이름', '학교', '학년/반', '등록일'])
                ws1.clear(); ws1.append_row(['ID', '이름', '학교', '학년/반', '등록일'])
                for s in load_students():
                    ws1.append_row([s['id'], s['name'],
                                    s.get('school',''), s.get('grade_class',''),
                                    s['created_at'][:10]])

                ws2 = self._get_or_create_ws(sh, '모의고사_내신',
                    ['학생', '구분', '시험명', '날짜', '원점수', '등급', '백분위', '메모'])
                ws2.clear(); ws2.append_row(['학생', '구분', '시험명', '날짜', '원점수', '등급', '백분위', '메모'])
                for r in sorted(load_scores(), key=lambda x: x.get('date',''), reverse=True):
                    if r.get('category') != 'mock': continue
                    ws2.append_row([stu_map.get(r['student_id'],'?'),
                                    r.get('exam_type',''), r.get('exam_name',''),
                                    r.get('date',''), r.get('score',''),
                                    r.get('grade',''), r.get('percentile',''), r.get('notes','')])

                ws3 = self._get_or_create_ws(sh, '워크시트_성적',
                    ['학생', '레벨', '문제유형', '날짜', '총점', '메모'])
                ws3.clear(); ws3.append_row(['학생', '레벨', '문제유형', '날짜', '총점', '메모'])
                for r in sorted(load_scores(), key=lambda x: x.get('date',''), reverse=True):
                    if r.get('category') != 'sheet': continue
                    ws3.append_row([stu_map.get(r['student_id'],'?'),
                                    r.get('level',''), r.get('qtype',''),
                                    r.get('date',''), r.get('total_score',''), r.get('notes','')])

                self.after(0, lambda: self._lbl_status.configure(
                    text='✅  Google Sheets 동기화 완료', fg=OK))
            except Exception as e:
                msg = str(e)
                self.after(0, lambda m=msg: self._lbl_status.configure(
                    text=f'❌  동기화 실패: {m}', fg=ERR))

        threading.Thread(target=_do, daemon=True).start()

    def _auto_sync_sheets(self):
        cfg = load_config()
        if cfg.get('oauth_key') and cfg.get('sheets_id'):
            threading.Thread(target=self._sync_all_to_sheets, daemon=True).start()

    def _refresh_score_combos(self):
        """학생 콤보박스 목록 갱신."""
        opts = self._student_options()
        for combo in (getattr(self,'_mock_stu_combo',None),
                      getattr(self,'_sheet_stu_combo',None)):
            if combo:
                combo['values'] = opts

    # ── 출력 서식 설정 전용 화면 ────────────────────────────────

    def _build_format_page(self):
        fr = tk.Frame(self._content, bg=BG)
        self._module_frames['conv_format'] = fr

        # 페이지 제목
        hdr = tk.Frame(fr, bg=BG2)
        hdr.pack(fill='x')
        tk.Frame(hdr, bg=GOLD, height=2).pack(side='bottom', fill='x')
        ph = tk.Frame(hdr, bg=BG2)
        ph.pack(fill='x', padx=24, pady=14)
        tk.Label(ph, text='🎨', font=(KO, 20), bg=BG2, fg=GOLD).pack(side='left', padx=(0, 10))
        tl = tk.Frame(ph, bg=BG2)
        tl.pack(side='left')
        tk.Label(tl, text='출력 서식 설정', font=(KO, 13, 'bold'), bg=BG2, fg=TEXT).pack(anchor='w')
        tk.Label(tl, text='Word 파일(.docx)을 불러와 글꼴·여백 설정을 출력에 적용합니다',
                 font=(KO, 8), bg=BG2, fg=SUB).pack(anchor='w', pady=(2, 0))

        # 본문 (스크롤 없는 단순 frame)
        body = tk.Frame(fr, bg=BG)
        body.pack(fill='both', expand=True, padx=40, pady=28)

        # ── 현재 서식 카드 (outer 직접 pack) ──────────────────────
        c_outer = tk.Frame(body, bg=BDR)
        c_outer.pack(fill='x', pady=(0, 16))
        c_inner = tk.Frame(c_outer, bg=CARD)
        c_inner.pack(fill='both', expand=True, padx=1, pady=1)
        cc = tk.Frame(c_inner, bg=CARD)
        cc.pack(fill='x', padx=22, pady=20)

        tk.Label(cc, text='현재 적용 서식', font=(KO, 8, 'bold'),
                 bg=CARD, fg=SUB).pack(anchor='w', pady=(0, 10))
        self._lbl_fmt_name = tk.Label(cc, text='', font=(KO, 13, 'bold'),
                                      bg=CARD, fg=TEXT)
        self._lbl_fmt_name.pack(anchor='w', pady=(0, 4))
        self._lbl_fmt_info = tk.Label(cc, text='', font=(KO, 9),
                                      bg=CARD, fg=SUB)
        self._lbl_fmt_info.pack(anchor='w', pady=(0, 18))

        br = tk.Frame(cc, bg=CARD)
        br.pack(fill='x')
        _Btn(br, '📂   서식 파일 불러오기',
             self._pick_template,
             fg_n='#1a1a1a', fg_h='#000000',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 10, 'bold'), pad=(20, 10)
             ).pack(side='left', padx=(0, 10))
        self._btn_reset_fmt = _Btn(br, '↩  기본값으로',
                                   self._clear_template,
                                   font=(KO, 9), pad=(14, 10))
        self._btn_reset_fmt.pack(side='left')

        # ── 안내 카드 (outer 직접 pack) ────────────────────────────
        g_outer = tk.Frame(body, bg=BDR)
        g_outer.pack(fill='x')
        g_inner = tk.Frame(g_outer, bg=CARD)
        g_inner.pack(fill='both', expand=True, padx=1, pady=1)
        gc = tk.Frame(g_inner, bg=CARD)
        gc.pack(fill='x', padx=22, pady=18)

        tk.Label(gc, text='💡  지원 파일 형식 및 추출 정보', font=(KO, 9, 'bold'),
                 bg=CARD, fg=TEXT).pack(anchor='w', pady=(0, 10))
        for row in [
            ('📝 Word (.docx)',      '글꼴 이름·크기, 상하좌우 여백  →  출력에 그대로 적용'),
            ('📋 한글 (.hwp)',       '글꼴 이름  →  여백은 기본값(상하 1.5cm / 좌우 2.0cm) 적용'),
            ('📊 PowerPoint (.pptx)','슬라이드 크기, 테마 글꼴  →  여백은 기본값 적용'),
        ]:
            fr_row = tk.Frame(gc, bg=CARD)
            fr_row.pack(fill='x', pady=3)
            tk.Label(fr_row, text=row[0], font=(KO, 8, 'bold'),
                     bg=CARD, fg=TEXT, width=20, anchor='w').pack(side='left')
            tk.Label(fr_row, text=row[1], font=(KO, 8),
                     bg=CARD, fg=SUB, anchor='w').pack(side='left')

        self._update_format_display(); self._update_fmt_badge()

    # ── 헤더 ────────────────────────────────────────────────────

    def _build_header(self):
        hdr = tk.Frame(self, bg=BG2)
        hdr.pack(fill='x')
        tk.Frame(hdr, bg=GOLD, height=2).pack(side='bottom', fill='x')

        inner = tk.Frame(hdr, bg=BG2)
        inner.pack(fill='x', padx=22, pady=0)

        left = tk.Frame(inner, bg=BG2)
        left.pack(side='left', fill='y', pady=14)
        tk.Label(left, text='📚', font=(KO, 20), bg=BG2, fg=GOLD).pack(side='left', padx=(0, 12))
        tbox = tk.Frame(left, bg=BG2)
        tbox.pack(side='left')
        tk.Label(tbox, text='English Worksheet System',
                 font=(KO, 14, 'bold'), bg=BG2, fg=TEXT).pack(anchor='w')
        tk.Label(tbox, text='문제 변환  ·  성적 입력  ·  학생 데이터  ·  약점 분석  ·  리포트',
                 font=(KO, 8), bg=BG2, fg=SUB).pack(anchor='w', pady=(2, 0))

        right = tk.Frame(inner, bg=BG2)
        right.pack(side='right', fill='y', pady=14)
        _Btn(right, '📂  문제 폴더', lambda: self._explorer(INPUT_DIR)).pack(side='left', padx=3)
        _Btn(right, '📁  결과 폴더', lambda: self._explorer(OUTPUT_DIR)).pack(side='left', padx=3)

    # ── 왼쪽 파일 패널 ──────────────────────────────────────────

    def _build_left(self, parent):
        # 타이틀
        th = tk.Frame(parent, bg=BG)
        th.pack(fill='x', pady=(14, 8))
        tk.Label(th, text='문  제  파  일', font=(KO, 10, 'bold'),
                 bg=BG, fg=GOLD).pack(side='left')

        _Btn(th, '↻ 새로고침', self._refresh_files).pack(side='right')

        # 스크롤 파일 목록
        wrap = tk.Frame(parent, bg=BDR, bd=1, relief='flat')
        wrap.pack(fill='both', expand=True)

        self._fc = tk.Canvas(wrap, bg=BG2, highlightthickness=0, bd=0)
        sb = ttk.Scrollbar(wrap, orient='vertical', command=self._fc.yview,
                           style='Vertical.TScrollbar')
        self._fc.configure(yscrollcommand=sb.set)
        sb.pack(side='right', fill='y')
        self._fc.pack(fill='both', expand=True)

        self._fi = tk.Frame(self._fc, bg=BG2)
        self._cwin = self._fc.create_window((0, 0), window=self._fi, anchor='nw')

        # Configure 이벤트: 캔버스 너비 변경 시 내부 프레임 너비 맞추고 scrollregion 재계산
        def _on_fc_configure(e):
            self._fc.itemconfig(self._cwin, width=e.width)
            self.after(1, self._reset_file_scroll)
        self._fc.bind('<Configure>', _on_fc_configure)
        self._fc.bind('<MouseWheel>', lambda e: self._fc.yview_scroll(-1*(e.delta//120), 'units'))

        # 드래그 앤 드롭 등록
        if DND_FILES:
            for w in (self._fc, self._fi):
                w.drop_target_register(DND_FILES)
                w.dnd_bind('<<Drop>>', self._on_drop)

        # 드롭 안내 텍스트 (DnD 사용 가능 시)
        if DND_FILES:
            tk.Label(parent, text='파일을 여기로 드래그하거나 아래 버튼으로 추가하세요',
                     font=(KO, 7), bg=BG, fg=DIM).pack(pady=(6, 0))

        # 하단 버튼 2개
        btn_row = tk.Frame(parent, bg=BG)
        btn_row.pack(fill='x', pady=(6, 0))
        _Btn(btn_row, '➕  파일 추가',
             self._pick_files,
             fg_n=TEXT, fg_h='#ffffff',
             bg_n=GOLD, bg_h=GOLDB,
             font=(KO, 9, 'bold'), pad=(10, 7)
             ).pack(side='left', fill='x', expand=True, padx=(0, 4))
        _Btn(btn_row, '📂  폴더 열기',
             lambda: self._explorer(INPUT_DIR),
             fg_n=SUB, fg_h=TEXT,
             pad=(10, 7)
             ).pack(side='left')

    # ── 파일 가져오기 ────────────────────────────────────────────

    def _pick_files(self):
        """파일 선택 대화상자 — PC 어디서든 .docx 파일 선택 후 문제 폴더로 복사."""
        paths = filedialog.askopenfilenames(
            title='문제 파일 선택',
            filetypes=[('Word 문서', '*.docx'), ('모든 파일', '*.*')],
            initialdir=str(Path.home()),
        )
        if paths:
            self._import_files([Path(p) for p in paths])

    def _on_drop(self, event):
        """드래그 앤 드롭 핸들러."""
        try:
            paths = self.tk.splitlist(event.data)
            self._import_files([Path(p) for p in paths])
        except Exception:
            pass

    def _import_files(self, paths):
        """파일 목록을 문제 폴더로 복사 후 목록 갱신."""
        copied = 0
        skipped = []
        for src in paths:
            src = Path(src)
            if src.suffix.lower() != '.docx':
                skipped.append(src.name)
                continue
            dst = INPUT_DIR / src.name
            if src.resolve() == dst.resolve():
                continue          # 이미 문제 폴더 안에 있음
            shutil.copy2(src, dst)
            copied += 1

        if copied:
            self._refresh_files()
            self._lbl_status.configure(
                text=f'✅  {copied}개 파일 추가됨', fg=OK)
        if skipped:
            self._lbl_status.configure(
                text=f'⚠️  .docx 파일만 지원합니다 ({", ".join(skipped[:2])})',
                fg=WARN)

    def _refresh_files(self):
        for w in self._fi.winfo_children(): w.destroy()
        self._file_items.clear()

        files = sorted(INPUT_DIR.glob('*.docx'))
        if not files:
            tk.Label(self._fi,
                     text='\n\n문제 폴더가 비어 있습니다\n\n.docx 파일을 추가한 뒤\n새로고침 해주세요',
                     font=(KO, 9), bg=BG2, fg=SUB, justify='center'
                     ).pack(expand=True, pady=30)
        else:
            for f in files:
                self._make_file_banner(f)

        # 한 프레임 후 실제 렌더링된 크기로 scrollregion 계산
        self.after(1, self._reset_file_scroll)

    def _reset_stu_scroll(self):
        """학생 목록 scrollregion 재설정."""
        try:
            children = self._stu_list_inner.winfo_children()
            if not children:
                content_h = 0
            else:
                self._stu_list_inner.update_idletasks()
                content_h = max(
                    (w.winfo_y() + w.winfo_height()
                     for w in children if w.winfo_exists()),
                    default=0)
            c = self._stu_list_canvas
            c.configure(scrollregion=(0, 0, c.winfo_width() or 300, content_h))
            c.yview_moveto(0)
        except Exception:
            pass

    def _reset_file_scroll(self):
        """파일 목록 scrollregion을 실제 렌더링 크기 기준으로 재설정."""
        try:
            children = self._fi.winfo_children()
            if not children:
                content_h = 0
            else:
                # 모든 자식 위젯의 실제 위치+높이 계산
                self._fi.update_idletasks()
                content_h = max(
                    (w.winfo_y() + w.winfo_height()
                     for w in children
                     if w.winfo_exists()),
                    default=0)
            canvas_w = self._fc.winfo_width() or 300
            self._fc.configure(scrollregion=(0, 0, canvas_w, content_h))
            self._fc.yview_moveto(0)
        except Exception:
            pass

    def _make_file_banner(self, path: Path):
        st   = path.stat()
        size = f'{st.st_size/1024:.0f} KB' if st.st_size < 1_048_576 else f'{st.st_size/1_048_576:.1f} MB'
        date = datetime.datetime.fromtimestamp(st.st_mtime).strftime('%Y.%m.%d')

        outer = tk.Frame(self._fi, bg=BG2)
        outer.pack(fill='x', padx=0, pady=1)

        # 배너 프레임
        fr = tk.Frame(outer, bg=ITEM, cursor='hand2')
        fr.pack(fill='x')

        # 선택 강조 바 (왼쪽)
        bar = tk.Frame(fr, bg=ITEM, width=4)
        bar.pack(side='left', fill='y')

        # 아이콘
        ico = tk.Label(fr, text='📄', font=(KO, 18), bg=ITEM, fg=TEXT, width=3)
        ico.pack(side='left', pady=12)

        # 텍스트
        txt = tk.Frame(fr, bg=ITEM)
        txt.pack(side='left', fill='both', expand=True, pady=12)
        name = tk.Label(txt, text=path.name, font=(KO, 10, 'bold'),
                        bg=ITEM, fg=TEXT, anchor='w')
        name.pack(fill='x')
        sub  = tk.Label(txt, text=f'{date}  ·  {size}',
                        font=(KO, 8), bg=ITEM, fg=SUB, anchor='w')
        sub.pack(fill='x')

        # 화살표
        arr = tk.Label(fr, text='›', font=(KO, 16), bg=ITEM, fg=DIM, padx=10)
        arr.pack(side='right')

        item_widgets = [fr, bar, ico, txt, name, sub, arr]
        self._file_items[path] = {'outer': outer, 'fr': fr, 'bar': bar,
                                  'widgets': item_widgets}

        for w in item_widgets:
            w.bind('<Enter>', lambda _, p=path: self._hover_banner(p, True))
            w.bind('<Leave>', lambda _, p=path: self._hover_banner(p, False))
            w.bind('<Button-1>', lambda _, p=path: self._select_file(p))

    def _hover_banner(self, path, entering):
        if path not in self._file_items: return
        item = self._file_items[path]
        if path == self._sel_file: return
        bg = HOV if entering else ITEM
        for w in item['widgets']:
            try: w.configure(bg=bg)
            except: pass

    def _select_file(self, path):
        if self._converting: return
        prev = self._sel_file
        self._sel_file = path

        # 이전 선택 해제
        if prev and prev in self._file_items:
            item = self._file_items[prev]
            item['bar'].configure(bg=ITEM)
            for w in item['widgets']:
                try: w.configure(bg=ITEM)
                except: pass

        # 새 선택 강조
        if path in self._file_items:
            item = self._file_items[path]
            item['bar'].configure(bg=GOLD)
            for w in item['widgets']:
                try: w.configure(bg=SEL)
                except: pass
            item['bar'].configure(bg=GOLD)

        self._update_right()

    # ── 서식 카드 ───────────────────────────────────────────────

    def _build_format_card(self, parent):
        card = self._card(parent)
        card.pack(fill='x', pady=(0, 10))
        cb = tk.Frame(card, bg=CARD)
        cb.pack(fill='x', padx=16, pady=12)

        # 타이틀
        th = tk.Frame(cb, bg=CARD)
        th.pack(fill='x', pady=(0, 6))
        tk.Label(th, text='출력 서식', font=(KO, 8, 'bold'),
                 bg=CARD, fg=SUB).pack(side='left')

        # 현재 서식 이름
        self._lbl_fmt_name = tk.Label(cb, text='', font=(KO, 10, 'bold'),
                                      bg=CARD, fg=TEXT, anchor='w')
        self._lbl_fmt_name.pack(fill='x', pady=(0, 2))

        # 서식 상세 정보
        self._lbl_fmt_info = tk.Label(cb, text='', font=(KO, 8),
                                      bg=CARD, fg=SUB, anchor='w')
        self._lbl_fmt_info.pack(fill='x', pady=(0, 8))

        # 버튼 행
        br = tk.Frame(cb, bg=CARD)
        br.pack(fill='x')
        _Btn(br, '📂  서식 파일 불러오기', self._pick_template,
             font=(KO, 8), pad=(10, 5)).pack(side='left', padx=(0, 6))
        self._btn_reset_fmt = _Btn(br, '↩  기본값으로', self._clear_template,
                                   font=(KO, 8), pad=(10, 5))
        self._btn_reset_fmt.pack(side='left')

        self._update_format_display(); self._update_fmt_badge()

    _FMT_ICONS = {'docx': '📝', 'hwp': '📋', 'pptx': '📊', 'ppt': '📊'}

    def _apply_reset_state(self):
        """설정 초기화 후 메모리 상태와 UI를 모두 기본값으로 동기화."""
        # ── 메모리 상태 초기화 ──
        self._tpl_path = None
        self._tpl_fmt  = {}

        # ── UI 업데이트 (widget 존재 시에만) ──
        for fn in (self._update_fmt_badge,
                   self._refresh_types_display,
                   self._apply_custom_types,
                   self._update_sheet_google_lock,
                   self._update_google_status):
            try:
                fn()
            except Exception:
                pass

        # Google 연결 UI 초기화
        try:
            self._drive_sel_lbl.configure(text='없음')
            for w in self._drive_list_inner.winfo_children():
                w.destroy()
        except Exception:
            pass

    def _update_fmt_badge(self):
        """변환 패널의 서식 배지를 현재 상태로 갱신."""
        if not hasattr(self, '_lbl_fmt_badge'):
            return
        if self._tpl_path:
            f       = self._tpl_fmt
            orient  = '세로' if f.get('orientation', 'landscape') == 'portrait' else '가로'
            fn      = f.get('font_name', '기본')
            ico     = self._FMT_ICONS.get(f.get('source_type', 'docx'), '📄')
            self._lbl_fmt_badge.configure(
                text=f'{ico}  {self._tpl_path.name}  ({orient} · {fn})',
                fg=GOLD)
        else:
            self._lbl_fmt_badge.configure(
                text='📋  기본 서식  (가로 A4 · 맑은 고딕)',
                fg=DIM)

    def _update_format_display(self):
        if self._tpl_path:
            name = self._tpl_path.name
            f    = self._tpl_fmt
            src  = f.get('source_type', 'docx')
            ico  = self._FMT_ICONS.get(src, '📄')

            if src == 'docx':
                fn  = f.get('font_name', '맑은 고딕')
                fs  = f.get('font_size', 10)
                lm  = f.get('left_margin_cm', 2.0)
                tm  = f.get('top_margin_cm', 1.5)
                orient = '가로' if f.get('orientation', 'landscape') == 'landscape' else '세로'
                info = f'글꼴: {fn} · {fs}pt  |  여백: 상하 {tm}cm / 좌우 {lm}cm  |  {orient}'
            elif src == 'hwp':
                fn     = f.get('font_name', '맑은 고딕')
                orient = '가로' if f.get('orientation', 'landscape') == 'landscape' else '세로'
                w      = f.get('page_width_cm', '?')
                h      = f.get('page_height_cm', '?')
                lm     = f.get('left_margin_cm', '기본')
                tm     = f.get('top_margin_cm', '기본')
                info   = f'글꼴: {fn}  |  {orient} {w}×{h}cm  |  여백: 상하{tm} / 좌우{lm}cm'
            elif src in ('pptx', 'ppt'):
                w  = f.get('slide_width_cm',  '?')
                h  = f.get('slide_height_cm', '?')
                fn = f.get('font_name', '맑은 고딕')
                info = f'PPT 참조  |  슬라이드: {w} × {h} cm  |  글꼴: {fn}'
            else:
                info = '서식 참조 중'

            self._lbl_fmt_name.configure(text=f'{ico}  {name}', fg=TEXT)
            self._lbl_fmt_info.configure(text=info)
            self._btn_reset_fmt.set_enabled(True)
        else:
            self._lbl_fmt_name.configure(text='📋  기본 서식 (내장)', fg=SUB)
            self._lbl_fmt_info.configure(
                text='맑은 고딕 · 10pt  |  여백 상하 1.5cm / 좌우 2.0cm')
            self._btn_reset_fmt.set_enabled(False)

    def _pick_template(self):
        path = filedialog.askopenfilename(
            title='서식 파일 선택',
            filetypes=[
                ('지원 파일',   '*.docx;*.hwp;*.pptx;*.ppt'),
                ('Word 문서',   '*.docx'),
                ('한글 문서',   '*.hwp'),
                ('PowerPoint',  '*.pptx;*.ppt'),
                ('모든 파일',   '*.*'),
            ],
            initialdir=str(Path.home()),
        )
        if not path:
            return
        p   = Path(path)
        fmt = read_template_fmt(p)
        self._tpl_path = p
        self._tpl_fmt  = fmt
        save_config({'template_path': str(p), 'fmt': fmt})
        self._update_format_display(); self._update_fmt_badge()
        self._lbl_status.configure(
            text=f'✅  서식 파일 적용됨: {p.name}', fg=OK)

    def _clear_template(self):
        self._tpl_path = None
        self._tpl_fmt  = {}
        save_config({'template_path': None, 'fmt': {}})
        self._update_format_display(); self._update_fmt_badge()
        self._lbl_status.configure(text='기본 서식으로 초기화됨', fg=SUB)

    # ── 오른쪽 패널 ─────────────────────────────────────────────

    def _build_right(self, parent):
        self._rp = parent

        # ① 빈 상태
        self._view_empty = tk.Frame(parent, bg=BG)
        self._view_empty.place(relx=0, rely=0, relwidth=1, relheight=1)
        tk.Label(self._view_empty, text='←  변환할 파일을 선택하세요',
                 font=(KO, 13), bg=BG, fg=SUB).place(relx=0.5, rely=0.44, anchor='center')
        tk.Label(self._view_empty, text='왼쪽 목록에서 .docx 파일을 클릭하면 시작할 수 있습니다',
                 font=(KO, 9), bg=BG, fg=DIM).place(relx=0.5, rely=0.52, anchor='center')

        # ② 준비 + 변환 중 뷰
        self._view_main = tk.Frame(parent, bg=BG)

        # 파일 정보 카드
        info_card = self._card(self._view_main)
        info_card.pack(fill='x', pady=(14, 10))
        ib = tk.Frame(info_card, bg=CARD)
        ib.pack(fill='x', padx=16, pady=14)
        tk.Label(ib, text='선택된 파일', font=(KO, 8), bg=CARD, fg=SUB).pack(anchor='w')
        self._lbl_fname = tk.Label(ib, text='', font=(KO, 12, 'bold'), bg=CARD, fg=TEXT)
        self._lbl_fname.pack(anchor='w', pady=(3, 6))
        tk.Label(ib, text='저장 경로', font=(KO, 8), bg=CARD, fg=SUB).pack(anchor='w')
        self._lbl_fpath = tk.Label(ib, text='', font=(KO, 9), bg=CARD, fg=GOLD)
        self._lbl_fpath.pack(anchor='w', pady=(2, 0))

        # 적용 서식 배지
        fmt_row = tk.Frame(ib, bg=CARD)
        fmt_row.pack(fill='x', pady=(6, 0))
        tk.Frame(fmt_row, bg=BDR, height=1).pack(fill='x', pady=(0, 6))
        lf = tk.Frame(fmt_row, bg=CARD)
        lf.pack(fill='x')
        tk.Label(lf, text='적용 서식', font=(KO, 8), bg=CARD, fg=SUB).pack(side='left')
        self._lbl_fmt_badge = tk.Label(lf, text='', font=(KO, 8),
                                       bg=CARD, fg=DIM)
        self._lbl_fmt_badge.pack(side='left', padx=(8, 0))
        self._update_fmt_badge()

        # 변환 시작 버튼
        self._btn_convert = _Btn(
            self._view_main, '  ▶   변환 시작  ',
            self._start_conversion,
            fg_n='#1a1a1a', fg_h='#000000',
            bg_n=GOLD, bg_h=GOLDB,
            font=(KO, 12, 'bold'), pad=(24, 12))
        self._btn_convert.pack(fill='x', pady=(0, 10))

        # 진행 카드
        prog_card = self._card(self._view_main)
        prog_card.pack(fill='x', pady=(0, 10))
        pb = tk.Frame(prog_card, bg=CARD)
        pb.pack(fill='x', padx=16, pady=14)

        ph = tk.Frame(pb, bg=CARD)
        ph.pack(fill='x', pady=(0, 6))
        self._lbl_prog = tk.Label(ph, text='변환 대기', font=(KO, 10, 'bold'), bg=CARD, fg=TEXT)
        self._lbl_prog.pack(side='left')
        self._lbl_pct  = tk.Label(ph, text='', font=(KO, 9), bg=CARD, fg=GOLD)
        self._lbl_pct.pack(side='right')

        self._progbar = ttk.Progressbar(pb, variable=self._prog_val,
                                         maximum=4, style='Prog.Horizontal.TProgressbar')
        self._progbar.pack(fill='x', pady=(0, 10))

        # 파일별 상태 행
        for key, (ico_txt, label, _) in FILE_META.items():
            row = tk.Frame(pb, bg=CARD)
            row.pack(fill='x', pady=2)
            ico_l = tk.Label(row, text='○', font=(KO, 11), bg=CARD, fg=DIM, width=2)
            ico_l.pack(side='left')
            tk.Label(row, text=f' {ico_txt}  {label}', font=(KO, 9), bg=CARD, fg=SUB).pack(side='left', padx=4)
            st_l = tk.Label(row, text='대기 중', font=(KO, 8), bg=CARD, fg=DIM)
            st_l.pack(side='right')
            self._status[key] = {'ico': ico_l, 'st': st_l}

        # 로그 카드
        log_card = self._card(self._view_main)
        log_card.pack(fill='both', expand=True, pady=(0, 0))
        lf = tk.Frame(log_card, bg=CARD)
        lf.pack(fill='both', expand=True, padx=16, pady=12)
        tk.Label(lf, text='변환 로그', font=(KO, 8, 'bold'), bg=CARD, fg=SUB
                 ).pack(anchor='w', pady=(0, 4))

        logwrap = tk.Frame(lf, bg=ITEM, bd=1, relief='flat')
        logwrap.pack(fill='both', expand=True)
        self._log = tk.Text(logwrap, bg=ITEM, fg=TEXT, font=(MONO, 8),
                            relief='flat', bd=0, state='disabled',
                            selectbackground=SEL, wrap='word', height=6)
        lsb = ttk.Scrollbar(logwrap, command=self._log.yview, style='Vertical.TScrollbar')
        self._log.configure(yscrollcommand=lsb.set)
        lsb.pack(side='right', fill='y')
        self._log.pack(fill='both', expand=True, padx=6, pady=6)
        for tag, fg in (('ok',OK),('err',ERR),('warn',WARN),('gold',GOLD),('sub',SUB)):
            self._log.tag_configure(tag, foreground=fg)

        # ③ 결과 뷰 (변환 완료 후)
        self._view_result = tk.Frame(parent, bg=BG)
        self._build_result_view()

        # 초기: 빈 뷰
        self._show_view('empty')

    def _build_result_view(self):
        rv = self._view_result

        # 성공 배너
        banner = tk.Frame(rv, bg=OK, height=52)
        banner.pack(fill='x', pady=(14, 14))
        banner.pack_propagate(False)
        tk.Label(banner, text='✅   변환 완료 — 결과 파일을 클릭하여 확인하세요',
                 font=(KO, 11, 'bold'), bg=OK, fg='#ffffff').pack(expand=True)

        # 파일 가로 목록 (1×4)
        grid = tk.Frame(rv, bg=BG)
        grid.pack(fill='x', padx=4)
        for i in range(4):
            grid.columnconfigure(i, weight=1)

        self._result_cards = {}
        for col, key in enumerate(('easy', 'standard_a', 'standard_b', 'hard')):
            ico_txt, label, accent = FILE_META[key]
            card = self._make_result_card(grid, key, ico_txt, label, accent)
            card.grid(row=0, column=col, padx=5, pady=5, sticky='nsew')

        # 폴더 열기 버튼
        btn_row = tk.Frame(rv, bg=BG)
        btn_row.pack(fill='x', pady=(10, 0), padx=4)
        _Btn(btn_row, '📁   결과 폴더 열기',
             lambda: self._explorer(self._out_dir),
             fg_n=GOLD, fg_h=GOLDB, bg_n=CARD, bg_h=HOV,
             font=(KO, 10, 'bold'), pad=(16, 10)
             ).pack(side='left')
        _Btn(btn_row, '↩   다른 파일 변환',
             self._back_to_main,
             font=(KO, 9), pad=(14, 10)
             ).pack(side='right')

    def _make_result_card(self, parent, key, ico, label, accent):
        outer = tk.Frame(parent, bg=BDR, bd=1, relief='flat')
        inner = tk.Frame(outer, bg=CARD, cursor='hand2')
        inner.pack(fill='both', expand=True, padx=1, pady=1)

        top = tk.Frame(inner, bg=accent, height=4)
        top.pack(fill='x')

        body = tk.Frame(inner, bg=CARD)
        body.pack(fill='both', expand=True, padx=10, pady=10)

        ico_l  = tk.Label(body, text=ico, font=(KO, 20), bg=CARD, fg=TEXT)
        ico_l.pack(anchor='w')
        name_l = tk.Label(body, text=f'{key}.docx', font=(KO, 8, 'bold'), bg=CARD, fg=TEXT)
        name_l.pack(anchor='w', pady=(3, 1))
        lbl_l  = tk.Label(body, text=label, font=(KO, 7), bg=CARD, fg=SUB)
        lbl_l.pack(anchor='w')
        size_l = tk.Label(body, text='—', font=(KO, 7), bg=CARD, fg=DIM)
        size_l.pack(anchor='w', pady=(3, 0))
        hint_l = tk.Label(body, text='클릭하여 열기  ↗', font=(KO, 7), bg=CARD, fg=DIM)
        hint_l.pack(anchor='w', pady=(1, 0))

        all_w = [outer, inner, top, body, ico_l, name_l, lbl_l, size_l, hint_l]
        for w in [inner, body, ico_l, name_l, lbl_l, size_l, hint_l]:
            w.bind('<Enter>', lambda _, ws=all_w: [x.configure(bg=HOV) for x in ws if x not in (outer, top)])
            w.bind('<Leave>', lambda _, ws=all_w: [x.configure(bg=CARD) for x in ws if x not in (outer, top)])
            w.bind('<Button-1>', lambda _, k=key: self._open_result(k))

        self._result_cards[key] = {'size': size_l}
        return outer

    def _update_result_sizes(self):
        if not self._out_dir: return
        for key in FILE_META:
            path = self._out_dir / f'{key}.docx'
            lbl  = self._result_cards.get(key, {}).get('size')
            if not lbl: continue
            if path.exists():
                size = path.stat().st_size
                txt  = f'{size//1024} KB'
            else:
                txt = '—'
            lbl.configure(text=txt)

    def _open_result(self, key):
        if not self._out_dir: return
        path = self._out_dir / f'{key}.docx'
        if path.exists():
            os.startfile(str(path))

    def _back_to_main(self):
        self._show_view('main')
        self._reset_progress()

    def _show_view(self, which):
        for v in (self._view_empty, self._view_main, self._view_result):
            v.place_forget()
        if which == 'empty':
            self._view_empty.place(relx=0, rely=0, relwidth=1, relheight=1)
        elif which == 'main':
            self._view_main.place(relx=0, rely=0, relwidth=1, relheight=1)
        elif which == 'result':
            self._view_result.place(relx=0, rely=0, relwidth=1, relheight=1)

    def _update_right(self):
        if self._sel_file:
            self._lbl_fname.configure(text=f'📄  {self._sel_file.name}')
            rel = OUTPUT_DIR / self._sel_file.stem
            self._lbl_fpath.configure(text=f'결과  ▸  {rel.relative_to(BASE_DIR)}')
            self._out_dir = rel
            self._reset_progress()
            self._show_view('main')
        else:
            self._show_view('empty')

    def _reset_progress(self):
        self._prog_val.set(0)
        self._lbl_prog.configure(text='변환 대기', fg=TEXT)
        self._lbl_pct.configure(text='')
        for v in self._status.values():
            v['ico'].configure(text='○', fg=DIM)
            v['st'].configure(text='대기 중', fg=DIM)
        self._log_clear()

    # ── 상태바 ──────────────────────────────────────────────────

    def _build_statusbar(self):
        bar = tk.Frame(self, bg=BG2, height=38)
        bar.pack(fill='x', side='bottom')
        tk.Frame(bar, bg=BDR, height=1).pack(fill='x', side='top')
        inner = tk.Frame(bar, bg=BG2)
        inner.pack(fill='both', expand=True, padx=20)
        self._lbl_status = tk.Label(inner, text='파일을 선택하세요',
                                    font=(KO, 9), bg=BG2, fg=SUB)
        self._lbl_status.pack(side='left', pady=10)

    # ── 변환 실행 ────────────────────────────────────────────────

    def _start_conversion(self):
        if not self._sel_file or self._converting: return
        self._converting = True
        self._btn_convert.set_enabled(False)
        self._btn_convert.set_text('  ⏳  변환 중...  ')
        self._prog_val.set(0)
        self._lbl_prog.configure(text='변환 중...', fg=WARN)
        self._lbl_status.configure(text='변환 중...', fg=WARN)
        self._log_clear()
        fmt_label = self._tpl_path.name if self._tpl_path else '기본 서식'
        self._log_write(f'▶  {self._sel_file.name}  |  서식: {fmt_label}', 'gold')
        # 실제 적용되는 설정 로그 출력
        f = self._tpl_fmt
        if f:
            fn  = f.get('font_name',        '맑은 고딕 (기본)')
            fs  = f.get('font_size',        10)
            tm  = f.get('top_margin_cm',    1.5)
            lm  = f.get('left_margin_cm',   2.0)
            self._log_write(f'  글꼴: {fn} / {fs}pt  |  여백: 상하{tm}cm · 좌우{lm}cm')
        threading.Thread(target=run_conversion,
                         args=(self._sel_file, self._out_dir, self._cb),
                         kwargs={'fmt': self._tpl_fmt},
                         daemon=True).start()

    def _cb(self, event, data):
        self.after(0, self._handle_cb, event, data)

    def _handle_cb(self, event, data):
        if event == 'log':
            self._log_write(f'  {data}')
        elif event == 'start':
            v = self._status[data]
            v['ico'].configure(text='⏳', fg=WARN); v['st'].configure(text='변환 중...', fg=WARN)
            self._log_write(f'  ▶ {data} 생성 중...')
        elif event == 'done':
            key, _ = data
            v = self._status[key]
            v['ico'].configure(text='✓', fg=OK); v['st'].configure(text='완료', fg=OK)
            done = sum(1 for x in self._status.values() if x['st'].cget('text') == '완료')
            self._prog_val.set(done)
            self._lbl_pct.configure(text=f'{done} / 4')
            self._log_write(f'  ✅  {key}.docx 생성 완료', 'ok')
        elif event == 'skip':
            v = self._status[data]
            v['ico'].configure(text='–', fg=SUB); v['st'].configure(text='없음', fg=SUB)
        elif event == 'error':
            key, msg = data
            v = self._status[key]
            v['ico'].configure(text='✗', fg=ERR); v['st'].configure(text='오류', fg=ERR)
            self._log_write(f'  ❌  {key}: {msg}', 'err')
        elif event == 'finished':
            self._lbl_prog.configure(text='✅  완료!', fg=OK)
            self._lbl_status.configure(text='✅  변환 완료', fg=OK)
            self._log_write('✅  모든 파일 변환 완료', 'ok')
            self._converting = False
            self._btn_convert.set_enabled(True)
            self._btn_convert.set_text('  ▶   변환 시작  ')
            self._update_result_sizes()
            self.after(800, lambda: self._show_view('result'))
        elif event == 'fatal':
            self._lbl_prog.configure(text='❌  실패', fg=ERR)
            self._lbl_status.configure(text='❌  변환 실패', fg=ERR)
            self._log_write(f'❌  오류: {data}', 'err')
            self._converting = False
            self._btn_convert.set_enabled(True)
            self._btn_convert.set_text('  ▶   변환 시작  ')

    # ── 헬퍼 ─────────────────────────────────────────────────────

    def _card(self, parent):
        outer = tk.Frame(parent, bg=BDR, bd=1, relief='flat')
        inner = tk.Frame(outer, bg=CARD)
        inner.pack(fill='both', expand=True, padx=1, pady=1)
        return inner

    def _explorer(self, path: Path):
        if path and path.exists():
            subprocess.Popen(['explorer', str(path)])

    def _log_clear(self):
        self._log.configure(state='normal')
        self._log.delete('1.0', 'end')
        self._log.configure(state='disabled')

    def _log_write(self, msg, tag=None):
        ts = datetime.datetime.now().strftime('%H:%M:%S')
        self._log.configure(state='normal')
        self._log.insert('end', f'{ts}  {msg}\n', tag or '')
        self._log.see('end')
        self._log.configure(state='disabled')


# ══════════════════════════════════════════════════════════════════
# 엔트리포인트
# ══════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    app = ConverterApp()
    app.mainloop()
